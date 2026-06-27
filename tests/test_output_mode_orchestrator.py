"""Integration tests: output_mode threading through process_files().

AC-6: process_files() passes output_mode to translate_docx / translate_pptx.
AC-7: with len(targets) > 1, process_files() clamps output_mode to "append".
AC-8 (office-output-mode): process_files() threads output_mode to translate_xlsx_xls;
      bilingual degrades to append for non-DOCX files with a warnings_callback notice.

Anti-tautology:
- Patch at consumer-bound names (app.backend.processors.orchestrator.translate_docx,
  …translate_pptx, …translate_xlsx_xls) — module-level import pattern per CLAUDE.md.
- Assert call_args.kwargs["output_mode"] value (selection assertion, not count).
- Call process_files() directly, not translate_document() (avoids tautology via
  wrong entry-point pattern per CLAUDE.md).

Stubs accept terms_getter=None AND output_mode=None kwargs per implementation-plan.md
handoff constraints.
"""

from __future__ import annotations

from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest

_FIXTURE_DOCX = Path(__file__).parent / "fixtures" / "minimal_phase0.docx"


def _make_mock_client():
    mock = MagicMock()
    mock.system_prompt = ""
    mock.model_type = "general"
    mock._is_translation_dedicated.return_value = False
    mock._is_translategemma_model.return_value = False
    mock.health_check.return_value = (True, "ok")
    return mock


def _make_mock_pptx(tmp_path: Path) -> Path:
    """Create a minimal PPTX fixture for pptx tests."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    txBox.text_frame.text = "Hello world"
    p = tmp_path / "test.pptx"
    prs.save(str(p))
    return p


def _fake_translate_docx(
    src_path,
    out_path,
    targets,
    src_lang,
    client,
    *,
    stop_flag=None,
    log=None,
    max_batch_chars=None,
    pre_translate_hook=None,
    post_translate_hook=None,
    include_headers_shapes_via_com=False,
    terms_getter=None,
    output_mode=None,
):
    """Minimal translate_docx stub that satisfies the orchestrator call signature."""
    import shutil
    shutil.copy2(src_path, out_path)
    return False


def _fake_translate_pptx(
    src_path,
    out_path,
    targets,
    src_lang,
    client,
    *,
    stop_flag=None,
    log=None,
    max_batch_chars=None,
    pre_translate_hook=None,
    post_translate_hook=None,
    terms_getter=None,
    output_mode=None,
):
    """Minimal translate_pptx stub that satisfies the orchestrator call signature."""
    import shutil
    shutil.copy2(src_path, out_path)
    return False


def _fake_translate_xlsx(
    in_path,
    out_path,
    targets,
    src_lang,
    client,
    *,
    stop_flag=None,
    log=None,
    max_batch_chars=None,
    pre_translate_hook=None,
    post_translate_hook=None,
    terms_getter=None,
    block_overrides=None,
    status_callback=None,
    output_mode=None,
):
    """Minimal translate_xlsx_xls stub; creates a minimal output XLSX."""
    import openpyxl
    from pathlib import Path as _Path
    wb = openpyxl.Workbook()
    _Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(_Path(out_path).with_suffix(".xlsx")))
    return False


def _make_xlsx_fixture(tmp_path: Path) -> Path:
    """Create a minimal XLSX file for orchestrator tests."""
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "Hello world"
    p = tmp_path / "test.xlsx"
    wb.save(str(p))
    return p


# ---------------------------------------------------------------------------
# AC-6: output_mode threaded to translate_docx
# ---------------------------------------------------------------------------

def test_orchestrator_threads_output_mode_to_translate_docx(tmp_path):
    """process_files with output_mode='replace' passes output_mode='replace' to translate_docx."""
    from app.backend.processors.orchestrator import process_files

    output_dir = tmp_path / "out"
    output_dir.mkdir(parents=True, exist_ok=True)

    with (
        patch(
            "app.backend.processors.orchestrator.translate_docx",
            side_effect=_fake_translate_docx,
        ) as mock_docx,
        patch(
            "app.backend.processors.orchestrator.OllamaClient",
            return_value=_make_mock_client(),
        ),
    ):
        process_files(
            files=[_FIXTURE_DOCX],
            output_dir=output_dir,
            targets=["vi"],
            src_lang="en",
            include_headers_shapes_via_com=False,
            ollama_model="qwen2.5:32b",
            output_mode="replace",
            log=lambda s: None,
        )

    assert mock_docx.called, "translate_docx was not called"
    call_kwargs = mock_docx.call_args.kwargs
    assert call_kwargs.get("output_mode") == "replace", (
        f"Expected output_mode='replace' but got: {call_kwargs.get('output_mode')!r}"
    )


# ---------------------------------------------------------------------------
# AC-6: output_mode threaded to translate_pptx
# ---------------------------------------------------------------------------

def test_orchestrator_threads_output_mode_to_translate_pptx(tmp_path):
    """process_files with output_mode='replace' passes output_mode='replace' to translate_pptx."""
    from app.backend.processors.orchestrator import process_files

    pptx_path = _make_mock_pptx(tmp_path)
    output_dir = tmp_path / "out"
    output_dir.mkdir(parents=True, exist_ok=True)

    with (
        patch(
            "app.backend.processors.orchestrator.translate_pptx",
            side_effect=_fake_translate_pptx,
        ) as mock_pptx,
        patch(
            "app.backend.processors.orchestrator.OllamaClient",
            return_value=_make_mock_client(),
        ),
    ):
        process_files(
            files=[pptx_path],
            output_dir=output_dir,
            targets=["vi"],
            src_lang="en",
            include_headers_shapes_via_com=False,
            ollama_model="qwen2.5:32b",
            output_mode="replace",
            log=lambda s: None,
        )

    assert mock_pptx.called, "translate_pptx was not called"
    call_kwargs = mock_pptx.call_args.kwargs
    assert call_kwargs.get("output_mode") == "replace", (
        f"Expected output_mode='replace' but got: {call_kwargs.get('output_mode')!r}"
    )


# ---------------------------------------------------------------------------
# AC-7: multi-target clamps replace → append
# ---------------------------------------------------------------------------

def test_orchestrator_clamps_replace_to_append_for_multi_target(tmp_path):
    """BR-67: with >1 targets, process_files clamps output_mode='replace' to 'append'."""
    from app.backend.processors.orchestrator import process_files

    output_dir = tmp_path / "out"
    output_dir.mkdir(parents=True, exist_ok=True)

    with (
        patch(
            "app.backend.processors.orchestrator.translate_docx",
            side_effect=_fake_translate_docx,
        ) as mock_docx,
        patch(
            "app.backend.processors.orchestrator.OllamaClient",
            return_value=_make_mock_client(),
        ),
    ):
        process_files(
            files=[_FIXTURE_DOCX],
            output_dir=output_dir,
            targets=["vi", "de"],
            src_lang="en",
            include_headers_shapes_via_com=False,
            ollama_model="qwen2.5:32b",
            output_mode="replace",
            log=lambda s: None,
        )

    assert mock_docx.called, "translate_docx was not called"
    call_kwargs = mock_docx.call_args.kwargs
    # BR-67: multi-target forces append
    assert call_kwargs.get("output_mode") == "append", (
        f"Expected output_mode clamped to 'append' for multi-target, "
        f"got: {call_kwargs.get('output_mode')!r}"
    )


# ---------------------------------------------------------------------------
# office-output-mode: XLSX output_mode threading
# ---------------------------------------------------------------------------

def test_orchestrator_threads_output_mode_to_translate_xlsx(tmp_path):
    """process_files with output_mode='replace' and an XLSX file passes output_mode='replace'
    to translate_xlsx_xls (the missing kwarg was the pre-change bug)."""
    from app.backend.processors.orchestrator import process_files

    xlsx_path = _make_xlsx_fixture(tmp_path)
    output_dir = tmp_path / "out"
    output_dir.mkdir(parents=True, exist_ok=True)

    with (
        patch(
            "app.backend.processors.orchestrator.translate_xlsx_xls",
            side_effect=_fake_translate_xlsx,
        ) as mock_xlsx,
        patch(
            "app.backend.processors.orchestrator.OllamaClient",
            return_value=_make_mock_client(),
        ),
    ):
        process_files(
            files=[xlsx_path],
            output_dir=output_dir,
            targets=["vi"],
            src_lang="en",
            include_headers_shapes_via_com=False,
            ollama_model="qwen2.5:32b",
            output_mode="replace",
            log=lambda s: None,
        )

    assert mock_xlsx.called, "translate_xlsx_xls was not called"
    call_kwargs = mock_xlsx.call_args.kwargs
    assert call_kwargs.get("output_mode") == "replace", (
        f"Expected output_mode='replace' threaded to translate_xlsx_xls, "
        f"got: {call_kwargs.get('output_mode')!r}"
    )


# ---------------------------------------------------------------------------
# office-output-mode: bilingual degrades to append for non-DOCX with warning
# ---------------------------------------------------------------------------

def test_orchestrator_degrades_bilingual_to_append_for_non_docx_with_warning(tmp_path):
    """process_files with output_mode='bilingual' and an XLSX file must:
    - pass output_mode='append' to translate_xlsx_xls (degrade), and
    - call warnings_callback with a notice about the degradation.
    """
    from app.backend.processors.orchestrator import process_files

    xlsx_path = _make_xlsx_fixture(tmp_path)
    output_dir = tmp_path / "out"
    output_dir.mkdir(parents=True, exist_ok=True)

    warnings_received: list = []

    with (
        patch(
            "app.backend.processors.orchestrator.translate_xlsx_xls",
            side_effect=_fake_translate_xlsx,
        ) as mock_xlsx,
        patch(
            "app.backend.processors.orchestrator.OllamaClient",
            return_value=_make_mock_client(),
        ),
    ):
        process_files(
            files=[xlsx_path],
            output_dir=output_dir,
            targets=["vi"],
            src_lang="en",
            include_headers_shapes_via_com=False,
            ollama_model="qwen2.5:32b",
            output_mode="bilingual",
            log=lambda s: None,
            warnings_callback=warnings_received.append,
        )

    assert mock_xlsx.called, "translate_xlsx_xls was not called"
    call_kwargs = mock_xlsx.call_args.kwargs
    # output_mode must be degraded to "append" for the XLSX file
    assert call_kwargs.get("output_mode") == "append", (
        f"Expected bilingual degraded to 'append' for XLSX, "
        f"got: {call_kwargs.get('output_mode')!r}"
    )
    # A warning notice must have been emitted
    assert len(warnings_received) >= 1, (
        "Expected at least one warning about bilingual degradation, got none"
    )
    assert any("bilingual" in w.lower() for w in warnings_received), (
        f"Expected 'bilingual' in warning text, got: {warnings_received}"
    )
