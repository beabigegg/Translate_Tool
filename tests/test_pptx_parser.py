"""Tests for PptxParser."""

from __future__ import annotations

import tempfile
from pathlib import Path

import pptx
import pytest

from app.backend.models.translatable_document import ElementType
from app.backend.parsers.pptx_parser import PptxParser


class TestPptxParser:
    """Tests for PptxParser class."""

    @pytest.fixture
    def parser(self):
        """Create a parser instance."""
        return PptxParser()

    @pytest.fixture
    def simple_pptx(self):
        """Create a simple PPTX file for testing."""
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            temp_path = f.name

        prs = pptx.Presentation()
        layout = prs.slide_layouts[0]  # Title slide
        slide = prs.slides.add_slide(layout)

        title = slide.shapes.title
        title.text = "Presentation Title"

        subtitle = slide.placeholders[1]
        subtitle.text = "Subtitle text here"

        prs.save(temp_path)

        yield temp_path

        Path(temp_path).unlink(missing_ok=True)

    @pytest.fixture
    def multi_slide_pptx(self):
        """Create a multi-slide PPTX file."""
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            temp_path = f.name

        prs = pptx.Presentation()

        # Slide 1
        layout = prs.slide_layouts[0]
        slide1 = prs.slides.add_slide(layout)
        slide1.shapes.title.text = "First Slide"
        slide1.placeholders[1].text = "First slide content"

        # Slide 2
        layout = prs.slide_layouts[1]  # Title and content
        slide2 = prs.slides.add_slide(layout)
        slide2.shapes.title.text = "Second Slide"

        prs.save(temp_path)

        yield temp_path

        Path(temp_path).unlink(missing_ok=True)

    @pytest.fixture
    def table_pptx(self):
        """Create a PPTX file with a table."""
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            temp_path = f.name

        prs = pptx.Presentation()
        layout = prs.slide_layouts[5]  # Blank
        slide = prs.slides.add_slide(layout)

        # Add a 2x2 table
        from pptx.util import Inches

        x, y, cx, cy = Inches(1), Inches(1), Inches(6), Inches(2)
        table = slide.shapes.add_table(2, 2, x, y, cx, cy).table

        table.cell(0, 0).text = "Cell A1"
        table.cell(0, 1).text = "Cell B1"
        table.cell(1, 0).text = "Cell A2"
        table.cell(1, 1).text = "Cell B2"

        prs.save(temp_path)

        yield temp_path

        Path(temp_path).unlink(missing_ok=True)

    def test_supported_extensions(self, parser):
        """Test that parser declares PPTX support."""
        assert ".pptx" in parser.supported_extensions

    def test_file_not_found(self, parser):
        """Test handling of non-existent file."""
        with pytest.raises(FileNotFoundError):
            parser.parse("/nonexistent/file.pptx")

    def test_invalid_extension(self, parser):
        """Test handling of non-PPTX file."""
        with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as f:
            f.write(b"Not a PPTX")
            f.flush()
            temp_path = f.name

        try:
            with pytest.raises(ValueError, match="Not a PPTX"):
                parser.parse(temp_path)
        finally:
            Path(temp_path).unlink(missing_ok=True)

    def test_parse_simple_pptx(self, parser, simple_pptx):
        """Test parsing a simple PPTX file."""
        doc = parser.parse(simple_pptx)

        assert doc is not None
        assert doc.source_type == "pptx"
        assert len(doc.elements) >= 2  # Title and subtitle

        contents = [e.content for e in doc.elements]
        assert any("Presentation Title" in c for c in contents)
        assert any("Subtitle" in c for c in contents)

    def test_parse_multi_slide(self, parser, multi_slide_pptx):
        """Test parsing multi-slide PPTX."""
        doc = parser.parse(multi_slide_pptx)

        assert doc.metadata.page_count == 2
        assert len(doc.pages) == 2

        # Check elements from different slides
        slide1_elements = [e for e in doc.elements if e.page_num == 1]
        slide2_elements = [e for e in doc.elements if e.page_num == 2]

        assert len(slide1_elements) >= 1
        assert len(slide2_elements) >= 1

    def test_parse_table(self, parser, table_pptx):
        """Test parsing PPTX with table."""
        doc = parser.parse(table_pptx)

        # Find table cell elements
        table_cells = [e for e in doc.elements if e.element_type == ElementType.TABLE_CELL]
        assert len(table_cells) >= 4  # 2x2 table

        # Check that cells are marked correctly
        for cell in table_cells:
            assert cell.metadata.get("in_table") is True

        # Check cell content
        contents = [e.content for e in table_cells]
        assert "Cell A1" in contents
        assert "Cell B2" in contents

    def test_title_classification(self, parser, simple_pptx):
        """Test that title placeholders are classified as TITLE."""
        doc = parser.parse(simple_pptx)

        title_elements = [e for e in doc.elements if e.element_type == ElementType.TITLE]
        assert len(title_elements) >= 1

    def test_element_ids_unique(self, parser, multi_slide_pptx):
        """Test that element IDs are unique."""
        doc = parser.parse(multi_slide_pptx)

        ids = [e.element_id for e in doc.elements]
        assert len(ids) == len(set(ids)), "Element IDs should be unique"

    def test_empty_text_skipped(self, parser):
        """Test that empty text frames are skipped."""
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            temp_path = f.name

        prs = pptx.Presentation()
        layout = prs.slide_layouts[5]  # Blank
        slide = prs.slides.add_slide(layout)

        # Add an empty text box
        from pptx.util import Inches

        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        txBox.text_frame.text = ""  # Empty

        # Add a non-empty text box
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(3), Inches(1))
        txBox2.text_frame.text = "Actual content"

        prs.save(temp_path)

        try:
            doc = parser.parse(temp_path)
            contents = [e.content for e in doc.elements]

            # Only non-empty content should be present
            assert "Actual content" in contents
            assert "" not in contents
        finally:
            Path(temp_path).unlink(missing_ok=True)

    def test_bbox_extraction(self, parser, simple_pptx):
        """Test that bounding boxes are extracted."""
        doc = parser.parse(simple_pptx)

        # Check that at least some elements have bboxes
        elements_with_bbox = [e for e in doc.elements if e.bbox is not None]
        assert len(elements_with_bbox) > 0

        for elem in elements_with_bbox:
            assert elem.bbox.width > 0
            assert elem.bbox.height > 0

    def test_page_info(self, parser, simple_pptx):
        """Test page info extraction."""
        doc = parser.parse(simple_pptx)

        assert len(doc.pages) >= 1
        page = doc.pages[0]
        assert page.width > 0
        assert page.height > 0


class TestPptxParserMetadata:
    """Tests for PPTX metadata extraction."""

    @pytest.fixture
    def parser(self):
        """Create a parser instance."""
        return PptxParser()

    def test_metadata_extraction(self, parser):
        """Test document metadata extraction."""
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            temp_path = f.name

        prs = pptx.Presentation()
        prs.core_properties.title = "Test Presentation"
        prs.core_properties.author = "Test Author"

        layout = prs.slide_layouts[5]
        prs.slides.add_slide(layout)
        prs.save(temp_path)

        try:
            doc = parser.parse(temp_path)
            assert doc.metadata.title == "Test Presentation"
            assert doc.metadata.author == "Test Author"
        finally:
            Path(temp_path).unlink(missing_ok=True)
