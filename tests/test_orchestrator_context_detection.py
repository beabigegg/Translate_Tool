"""Tests for cloud-path document-context summary parity (BR-109).

Change: cloud-doc-context-summary

Covers AC-1..AC-7 (change-classification.md):
  AC-1: cloud active client generates the summary, no local Ollama required.
  AC-2: summary injected as "Document context: <summary>" into the system prompt
        via the existing downstream wiring.
  AC-3: both CONTEXT_DETECTION_ENABLED and QWEN_CONTEXT_FLOW_ENABLED still gate
        the behavior (AND-gated) on the cloud path.
  AC-4: client._is_translation_dedicated() still skips summary generation.
  AC-5: a failed/empty cloud summary call degrades gracefully (no preamble,
        job never aborts).
  AC-6: local-Ollama context detection is unchanged.
  AC-7: no scope creep into the injection wiring or JSON structured I/O.

Mock boundaries (per test-plan.md Notes):
  - requests.Session.post for integration/AC-1/AC-2/AC-6 (network boundary).
  - app.backend.processors.orchestrator._detect_document_context (same-module
    attribute; call site is unqualified) for AC-3/AC-4 guard tests, asserted
    by call-count.
  - Direct calls to _detect_document_context with a capability-restricted
    MagicMock(spec=["complete"]) for the AC-1/AC-5 unit tests, so an
    accidental fall-back to an Ollama-only private method raises
    AttributeError (caught internally) rather than silently succeeding --
    turning a wrong-seam regression into a loud assertion failure on the
    returned content, not just a call-count check.

Torch-free: no COMET/QE imports anywhere in this file.
"""

from __future__ import annotations

import inspect
import logging
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest

from app.backend.clients.ollama_client import OllamaClient
from app.backend.clients.openai_compatible_client import OpenAICompatibleClient
from app.backend.processors import libreoffice_helpers as lo
from app.backend.processors.orchestrator import (
    _detect_document_context,
    _sample_file_text,
    process_files,
)


# ---------------------------------------------------------------------------
# Shared fixtures / helpers
# ---------------------------------------------------------------------------

_PANJIT_CFG = {
    "providers": [
        {
            "id": "panjit",
            "type": "openai",
            "enabled": True,
            "base_url": "http://panjit-mock:8080",
            "api_key": "test-key-panjit",
            "models": {"translate": "gpt-oss:120b"},
        },
        {
            "id": "ollama-local",
            "type": "ollama",
            "enabled": True,
            "base_url": "http://localhost:11434",
            "api_key": "",
            "models": {"translate": "qwen3.5:9b"},
        },
    ],
    "routing": {"default": {"model": "gpt-oss:120b", "provider": "panjit", "profile": "general"}},
    "fallback_chain": ["panjit", "deepseek"],
}


def _make_real_docx(tmp_path: Path, text: str) -> Path:
    """Write a minimal but structurally valid .docx so _sample_file_text
    extracts real, non-empty text (required to satisfy the `sample` guard
    condition at orchestrator.py's context-detection call site)."""
    import docx as _docx

    src = tmp_path / "test.docx"
    doc = _docx.Document()
    doc.add_paragraph(text)
    doc.save(str(src))
    return src


def _mock_post_with_content(content: str):
    """Build a requests.Session.post side_effect returning a fixed chat-completion body."""

    def _mock_post(url, **kwargs):
        resp = MagicMock()
        resp.status_code = 200
        resp.json.return_value = {"choices": [{"message": {"content": content}}]}
        return resp

    return _mock_post


def _make_capturing_translate_docx(captured: dict, stopped: bool = False):
    """Build a translate_docx stub that snapshots client.system_prompt at
    call time.

    orchestrator.process_files' per-file `finally:` clause unconditionally
    resets `client.system_prompt = base_system_prompt` after every file
    (pre-existing cleanup, out of scope for this change) -- so the
    scenario-injected value is only observable at translate_docx call time,
    not by inspecting `client.system_prompt` after process_files() returns.
    """

    def _fn(src_path, out_path, targets, src_lang, client, **kwargs):
        captured["system_prompt"] = client.system_prompt
        return stopped

    return _fn


# ---------------------------------------------------------------------------
# AC-1: cloud active client used for the summary; no local Ollama required
# ---------------------------------------------------------------------------

def test_cloud_active_client_used_for_summary_not_local_ollama():
    """AC-1 (unit): _detect_document_context calls client.complete(), the shared
    seam BOTH clients implement -- not an Ollama-only private method.

    The fake cloud client is capability-restricted to only `complete` via
    MagicMock(spec=["complete"]); if the implementation regressed to an
    Ollama-only method, accessing it would raise AttributeError (caught by
    _detect_document_context's try/except) and silently return "" -- which
    the content assertion below would catch as a hard failure.
    """
    cloud_client = MagicMock(spec=["complete"])
    cloud_client.complete.return_value = (True, "This document is a purchase order.")

    result = _detect_document_context(
        cloud_client, "PO number 12345, widget parts, quantity 500.", log=lambda s: None, target_lang="en"
    )

    cloud_client.complete.assert_called_once()
    prompt_arg = cloud_client.complete.call_args[0][0]
    assert "PO number 12345" in prompt_arg
    assert result == "This document is a purchase order."


def test_cloud_summary_generated_without_local_ollama_present(tmp_path):
    """AC-1 (integration): the document-context summary is produced via the
    active cloud client's HTTP path; the local OllamaClient's raw-completion
    method is never invoked."""
    src_docx = _make_real_docx(tmp_path, "This is a purchase order for widget parts.")
    out_dir = tmp_path / "out"
    out_dir.mkdir()

    post_urls = []

    def _mock_post(url, **kwargs):
        post_urls.append(url)
        resp = MagicMock()
        resp.status_code = 200
        resp.json.return_value = {"choices": [{"message": {"content": "A purchase order document."}}]}
        return resp

    with patch("app.backend.config.load_providers_config", return_value=_PANJIT_CFG), \
         patch("requests.Session.post", side_effect=_mock_post), \
         patch("app.backend.processors.orchestrator.translate_docx", return_value=False), \
         patch.object(OllamaClient, "_call_ollama") as mock_ollama_call:
        result = process_files(
            files=[src_docx],
            output_dir=out_dir,
            targets=["French"],
            src_lang="English",
            include_headers_shapes_via_com=False,
            ollama_model="qwen3.5:9b",
            model_type="general",
            system_prompt="",
            profile_id="general",
            provider_id="panjit",
        )

    _, _, stopped, _client, _term, winning_provider = result
    assert stopped is False
    assert winning_provider == "panjit"
    mock_ollama_call.assert_not_called()
    assert any("/v1/chat/completions" in u for u in post_urls), (
        "cloud completions endpoint must have been called for the summary"
    )


# ---------------------------------------------------------------------------
# AC-2: summary injected as "Document context: <summary>" into the system prompt
# ---------------------------------------------------------------------------

def test_cloud_summary_injected_as_document_context_in_system_prompt(tmp_path):
    """AC-2 / BR-109 / ADR-0016: the cloud-generated summary reaches the
    OUTGOING request payload as system-channel content on the real
    translate_once() call -- not merely assigned to the `client.system_prompt`
    attribute (which was formerly an orchestrator-compatibility stub whose
    writes were silently discarded on the cloud path). The translatable user
    payload must stay clean (no bleed of the preamble into it).

    DYNAMIC_SCENARIO_STRATEGY_ENABLED is patched off to isolate this
    assertion from scenario-detection heuristics (out of scope; already
    covered by tests/test_translation_strategy.py), exercising the
    non-dynamic else-branch literal at orchestrator.py instead.

    The `translate_docx` stub simulates the one real behavior that matters
    here: dispatching a segment through the resolved client via
    `client.translate_once(...)`, so the scenario-injected
    `client.system_prompt` actually flows through to the outgoing HTTP
    request -- this is the seam a "silently ignored" regression would break,
    unlike inspecting the attribute after the call (asserted as a tautology
    in the prior version of this test).
    """
    src_docx = _make_real_docx(tmp_path, "This document is an engineering change request for a wafer fab line.")
    out_dir = tmp_path / "out"
    out_dir.mkdir()

    summary_text = "An engineering change request for a semiconductor fab line."
    post_bodies: list = []

    def _mock_post(url, **kwargs):
        post_bodies.append(kwargs.get("json", {}))
        resp = MagicMock()
        resp.status_code = 200
        # 1st POST = the context-detection summary call; every subsequent
        # POST = the real translate_once call driven by the stub below.
        content = summary_text if len(post_bodies) == 1 else "Bonjour, demande de modification technique."
        resp.json.return_value = {"choices": [{"message": {"content": content}}]}
        return resp

    def _translate_docx_dispatches_one_segment(src_path, out_path, targets, src_lang, client, **kwargs):
        client.translate_once("Engineering change request text.", targets[0], src_lang)
        return False

    with patch("app.backend.config.load_providers_config", return_value=_PANJIT_CFG), \
         patch("requests.Session.post", side_effect=_mock_post), \
         patch("app.backend.processors.orchestrator.translate_docx", side_effect=_translate_docx_dispatches_one_segment), \
         patch("app.backend.processors.orchestrator.DYNAMIC_SCENARIO_STRATEGY_ENABLED", False):
        result = process_files(
            files=[src_docx],
            output_dir=out_dir,
            targets=["French"],
            src_lang="English",
            include_headers_shapes_via_com=False,
            ollama_model="qwen3.5:9b",
            model_type="general",
            system_prompt="You are a professional translator.",
            profile_id="general",
            provider_id="panjit",
        )

    _, _, stopped, _client, _term, winning_provider = result
    assert stopped is False
    assert winning_provider == "panjit"

    assert len(post_bodies) >= 2, "expected a summary POST plus a translate_once POST"
    translate_body = post_bodies[-1]
    messages = translate_body["messages"]

    system_messages = [m for m in messages if m["role"] == "system"]
    assert len(system_messages) == 1, f"expected exactly one system message, got {system_messages!r}"
    assert f"Document context: {summary_text}" in system_messages[0]["content"]

    user_messages = [m for m in messages if m["role"] == "user"]
    assert not any(f"Document context: {summary_text}" in m["content"] for m in user_messages), (
        "the injected summary must never bleed into the translatable user payload"
    )


# ---------------------------------------------------------------------------
# AC-3: both flags AND-gate the cloud path (negative cases, call-count asserts)
# ---------------------------------------------------------------------------

def test_context_detection_disabled_skips_cloud_summary(tmp_path):
    """AC-3: CONTEXT_DETECTION_ENABLED=False must skip the summary even on the
    cloud path (call-count assertion, not happy-path)."""
    src_docx = _make_real_docx(tmp_path, "Some engineering document text.")
    out_dir = tmp_path / "out"
    out_dir.mkdir()

    with patch("app.backend.config.load_providers_config", return_value=_PANJIT_CFG), \
         patch("requests.Session.post", side_effect=_mock_post_with_content("n/a")), \
         patch("app.backend.processors.orchestrator.translate_docx", return_value=False), \
         patch("app.backend.processors.orchestrator.CONTEXT_DETECTION_ENABLED", False), \
         patch("app.backend.processors.orchestrator._detect_document_context") as mock_detect:
        process_files(
            files=[src_docx],
            output_dir=out_dir,
            targets=["French"],
            src_lang="English",
            include_headers_shapes_via_com=False,
            ollama_model="qwen3.5:9b",
            model_type="general",
            system_prompt="",
            profile_id="general",
            provider_id="panjit",
        )

    mock_detect.assert_not_called()


def test_qwen_context_flow_disabled_skips_cloud_summary(tmp_path):
    """AC-3: QWEN_CONTEXT_FLOW_ENABLED=False must skip the summary even on the
    cloud path (call-count assertion, not happy-path)."""
    src_docx = _make_real_docx(tmp_path, "Some engineering document text.")
    out_dir = tmp_path / "out"
    out_dir.mkdir()

    with patch("app.backend.config.load_providers_config", return_value=_PANJIT_CFG), \
         patch("requests.Session.post", side_effect=_mock_post_with_content("n/a")), \
         patch("app.backend.processors.orchestrator.translate_docx", return_value=False), \
         patch("app.backend.processors.orchestrator.QWEN_CONTEXT_FLOW_ENABLED", False), \
         patch("app.backend.processors.orchestrator._detect_document_context") as mock_detect:
        process_files(
            files=[src_docx],
            output_dir=out_dir,
            targets=["French"],
            src_lang="English",
            include_headers_shapes_via_com=False,
            ollama_model="qwen3.5:9b",
            model_type="general",
            system_prompt="",
            profile_id="general",
            provider_id="panjit",
        )

    mock_detect.assert_not_called()


# ---------------------------------------------------------------------------
# AC-4: translation-dedicated ACTIVE client still skips summary generation
# ---------------------------------------------------------------------------

def test_translation_dedicated_cloud_client_skips_summary(tmp_path):
    """AC-4: when the ACTIVE cloud client reports _is_translation_dedicated()
    == True, the summary is skipped exactly as before (call-count assertion)."""
    src_docx = _make_real_docx(tmp_path, "Some engineering document text.")
    out_dir = tmp_path / "out"
    out_dir.mkdir()

    with patch("app.backend.config.load_providers_config", return_value=_PANJIT_CFG), \
         patch("requests.Session.post", side_effect=_mock_post_with_content("n/a")), \
         patch("app.backend.processors.orchestrator.translate_docx", return_value=False), \
         patch.object(OpenAICompatibleClient, "_is_translation_dedicated", return_value=True), \
         patch("app.backend.processors.orchestrator._detect_document_context") as mock_detect:
        process_files(
            files=[src_docx],
            output_dir=out_dir,
            targets=["French"],
            src_lang="English",
            include_headers_shapes_via_com=False,
            ollama_model="qwen3.5:9b",
            model_type="general",
            system_prompt="",
            profile_id="general",
            provider_id="panjit",
        )

    mock_detect.assert_not_called()


# ---------------------------------------------------------------------------
# AC-5: graceful degradation -- exception or empty result never aborts the job
# ---------------------------------------------------------------------------

def test_detect_document_context_returns_empty_on_cloud_call_exception():
    """AC-5 (unit): client.complete() raising, or returning (False, ""), must
    not propagate out of _detect_document_context -- it degrades to ""."""
    cloud_client_raises = MagicMock(spec=["complete"])
    cloud_client_raises.complete.side_effect = RuntimeError("cloud provider unreachable")
    assert _detect_document_context(cloud_client_raises, "sample text", log=lambda s: None) == ""

    cloud_client_empty = MagicMock(spec=["complete"])
    cloud_client_empty.complete.return_value = (False, "")
    assert _detect_document_context(cloud_client_empty, "sample text", log=lambda s: None) == ""


def test_job_continues_with_no_preamble_when_cloud_summary_empty(tmp_path):
    """AC-5 (resilience): an empty cloud summary response degrades to no
    preamble; the job completes normally (stopped=False), never aborts."""
    src_docx = _make_real_docx(tmp_path, "Some engineering document text.")
    out_dir = tmp_path / "out"
    out_dir.mkdir()

    captured: dict = {}

    with patch("app.backend.config.load_providers_config", return_value=_PANJIT_CFG), \
         patch("requests.Session.post", side_effect=_mock_post_with_content("")), \
         patch("app.backend.processors.orchestrator.translate_docx", side_effect=_make_capturing_translate_docx(captured)), \
         patch("app.backend.processors.orchestrator.DYNAMIC_SCENARIO_STRATEGY_ENABLED", False):
        result = process_files(
            files=[src_docx],
            output_dir=out_dir,
            targets=["French"],
            src_lang="English",
            include_headers_shapes_via_com=False,
            ollama_model="qwen3.5:9b",
            model_type="general",
            system_prompt="Base prompt.",
            profile_id="general",
            provider_id="panjit",
        )

    _, _, stopped, _client, _term, winning_provider = result
    assert stopped is False
    assert winning_provider == "panjit"
    assert "Document context:" not in captured["system_prompt"]


# ---------------------------------------------------------------------------
# AC-6: local-Ollama context detection is unchanged
# ---------------------------------------------------------------------------

def test_local_ollama_context_detection_unchanged(tmp_path):
    """AC-6: the local-Ollama path still generates and injects the summary
    identically to before this change (byte-identical HTTP call via
    complete() -> _call_ollama(_build_no_system_payload(prompt)))."""
    src_docx = _make_real_docx(tmp_path, "Some engineering document text.")
    out_dir = tmp_path / "out"
    out_dir.mkdir()

    summary_text = "An engineering specification document."
    captured: dict = {}

    with patch.object(OllamaClient, "_call_ollama", return_value=(True, summary_text)) as mock_call, \
         patch("app.backend.processors.orchestrator.translate_docx", side_effect=_make_capturing_translate_docx(captured)), \
         patch("app.backend.processors.orchestrator.DYNAMIC_SCENARIO_STRATEGY_ENABLED", False):
        result = process_files(
            files=[src_docx],
            output_dir=out_dir,
            targets=["French"],
            src_lang="English",
            include_headers_shapes_via_com=False,
            ollama_model="qwen3.5:9b",
            model_type="general",
            system_prompt="Base prompt.",
            profile_id="general",
            provider_id=None,
        )

    _, _, stopped, _client, _term, winning_provider = result
    assert stopped is False
    assert winning_provider == "ollama-local"
    assert mock_call.called
    assert captured["system_prompt"] == f"Base prompt.\n\nDocument context: {summary_text}"


# ---------------------------------------------------------------------------
# AC-7: no scope creep -- injection wiring and JSON structured I/O untouched
# ---------------------------------------------------------------------------

def test_no_scope_creep_into_injection_wiring_or_json_io():
    """AC-7: only the generation seam changed. The LLMClient Protocol stays
    at exactly 5 methods (complete() intentionally excluded), the downstream
    "Document context:" injection wiring is untouched, and no JSON
    structured translation I/O (Step 3, a separate tracked change) was
    introduced into _detect_document_context."""
    from app.backend.clients.base_llm_client import LLMClient
    from app.backend.services import translation_strategy
    import app.backend.processors.orchestrator as orch_module

    methods = [
        name for name, member in inspect.getmembers(LLMClient)
        if not name.startswith("_") and callable(getattr(LLMClient, name, None))
    ]
    assert len(methods) == 5, f"Expected 5 Protocol methods, got {len(methods)}: {methods}"
    assert "complete" not in methods, "complete() must stay off the LLMClient Protocol"

    strategy_src = inspect.getsource(translation_strategy.build_strategy)
    assert 'f"Document context: {detected_context' in strategy_src, (
        "build_strategy's Document-context injection literal must be unchanged"
    )

    detect_src = inspect.getsource(orch_module._detect_document_context)
    assert "json.loads" not in detect_src
    assert "{\"translation\"" not in detect_src


# ---------------------------------------------------------------------------
# doc-context-sampling-fix (BR-109 valid-sample coverage + INFO observability)
#
# Covers AC-1..AC-8 (test-plan.md):
#   AC-1: legacy .xls sampling reads real text via the LibreOffice conversion.
#   AC-2: table-only .docx sampling reads doc.tables cell text.
#   AC-3: table/graphic-frame .pptx sampling reads table cell text.
#   AC-4: a skipped/failed sample is visible at INFO level with a reason.
#   AC-5: successful detection is visible at INFO level ([CONTEXT] Detected:).
#   AC-6: a sampler exception/degradation never aborts the job and never
#         injects a "Document context:" preamble.
#   AC-7: the sampler-side .xls conversion does not double-invoke LibreOffice.
#   AC-8: both a legacy .xls and a table-only .docx emit [CONTEXT] Detected:
#         in the same job.
#
# All new tests must fail RED against the unfixed orchestrator.py (empty-string
# sampler branches for .docx tables / .pptx tables / .xls, and a
# logger.debug-only swallow) before IP-1..IP-6 land.
# ---------------------------------------------------------------------------


@pytest.fixture(autouse=True)
def _reset_libreoffice_cache():
    """Reset the module-level LibreOffice detection cache around every test
    in this file (harmless for tests that never touch it)."""
    lo._LIBREOFFICE_BINARY = None
    lo._DETECTION_DONE = False
    yield
    lo._LIBREOFFICE_BINARY = None
    lo._DETECTION_DONE = False


def _make_xlsx_writing_popen(token: str):
    """Build a Popen-fake class that writes a REAL openpyxl-authored .xlsx
    (carrying *token*) at the --outdir/<stem>.xlsx path _libreoffice_convert
    expects to find.

    Unlike tests/test_libreoffice_helpers.py::_FakePopen (which writes literal
    b"converted-bytes" -- openpyxl cannot open that), this fake produces a
    genuine workbook so the orchestrator's openpyxl-based .xls sampler branch
    can actually read the distinctive token back out of it.
    """

    class _FakeXlsConversionPopen:
        def __init__(self, cmd, stdout=None, stderr=None, text=None, start_new_session=None):
            self.cmd = cmd
            self.pid = 87654
            self.returncode = 0
            outdir = cmd[cmd.index("--outdir") + 1]
            input_path = cmd[-1]
            target_format = cmd[cmd.index("--convert-to") + 1]
            stem = Path(input_path).stem

            from openpyxl import Workbook

            wb = Workbook()
            ws = wb.active
            ws["A1"] = token
            wb.save(str(Path(outdir) / f"{stem}.{target_format}"))

        def communicate(self, timeout=None):
            return ("", "")

    return _FakeXlsConversionPopen


def _make_soffice_available_patch():
    """shutil.which side_effect that reports 'soffice' present on PATH."""
    return lambda name: "/usr/bin/soffice" if name == "soffice" else None


# ---------------------------------------------------------------------------
# AC-1: legacy .xls sampling reads real text via LibreOffice conversion
# ---------------------------------------------------------------------------

def test_sample_file_text_reads_legacy_xls_via_conversion(tmp_path):
    """AC-1 (unit): _sample_file_text's .xls branch converts via the
    LibreOffice-headless boundary (mocked at subprocess.Popen, never at
    xls_to_xlsx's own internals) and reads the resulting .xlsx with openpyxl,
    returning a sample that contains the distinctive token -- not merely a
    non-empty string."""
    token = "PANJIT-XLS-TOKEN-771"
    input_path = tmp_path / "legacy.xls"
    input_path.write_bytes(b"dummy-xls-bytes-never-read-popen-is-mocked")

    with patch("shutil.which", side_effect=_make_soffice_available_patch()), \
         patch("subprocess.Popen", side_effect=_make_xlsx_writing_popen(token)):
        sample = _sample_file_text(input_path)

    assert token in sample, f"expected distinctive xls token in sample, got: {sample!r}"


def test_process_files_context_detected_for_legacy_xls(tmp_path, caplog):
    """AC-1 (integration): a full process_files() run on a real .xls file
    (LibreOffice mocked at subprocess.Popen) emits [CONTEXT] Detected: through
    the SAME channel JobLogger.log() uses in production (JobLogger.log -> the
    "TranslateTool" logger -> the RotatingFileHandler that writes
    translator.log; app/backend/utils/logging_utils.py). `log` is wired to
    that real logger here (mirroring JobLogger.log's own behavior) and the
    assertion is against the "TranslateTool" logger name -- asserting against
    the module logger alone would pass even if the line never reached
    translator.log, which is the precise silent-skip failure mode BR-109
    exists to eliminate."""
    from app.backend.utils import logging_utils

    src = tmp_path / "legacy.xls"
    src.write_bytes(b"dummy-xls-bytes-never-read-popen-is-mocked")
    out_dir = tmp_path / "out"
    out_dir.mkdir()

    def _log_via_translate_tool_logger(message: str) -> None:
        logging_utils.logger.info(message)

    with patch("shutil.which", side_effect=_make_soffice_available_patch()), \
         patch("subprocess.Popen", side_effect=_make_xlsx_writing_popen("PANJIT-XLS-TOKEN-771")), \
         patch("app.backend.processors.orchestrator.translate_xlsx_xls", return_value=False), \
         patch.object(OllamaClient, "_call_ollama", return_value=(True, "A spreadsheet document.")), \
         patch("app.backend.processors.orchestrator.DYNAMIC_SCENARIO_STRATEGY_ENABLED", False), \
         caplog.at_level(logging.INFO, logger="TranslateTool"):
        result = process_files(
            files=[src],
            output_dir=out_dir,
            targets=["French"],
            src_lang="English",
            include_headers_shapes_via_com=False,
            ollama_model="qwen3.5:9b",
            model_type="general",
            system_prompt="Base prompt.",
            profile_id="general",
            provider_id=None,
            log=_log_via_translate_tool_logger,
        )

    _, _, stopped, _client, _term, winning_provider = result
    assert stopped is False
    assert winning_provider == "ollama-local"
    # Filter on r.name: caplog's handler lives on the ROOT logger, so
    # caplog.at_level(..., logger="TranslateTool") does NOT restrict which
    # loggers' records land in caplog.records. Without this filter the
    # orchestrator module logger's additive `logger.info(...)` satisfies the
    # assertion on its own, and the test stays green even when the line never
    # reaches translator.log — the exact silent failure BR-109 forbids.
    assert any(
        r.name == "TranslateTool" and "[CONTEXT] Detected:" in r.message
        for r in caplog.records
    ), "expected a [CONTEXT] Detected: INFO record to reach the TranslateTool logger channel"


# ---------------------------------------------------------------------------
# AC-2: table-only .docx sampling reads doc.tables cell text
# ---------------------------------------------------------------------------

def test_sample_file_text_docx_table_only_includes_cell_text(tmp_path):
    """AC-2 (unit): a .docx whose only content lives in a table (no
    paragraph text) still yields a sample containing the table-cell token --
    proving the sampler reads doc.tables, not just doc.paragraphs."""
    import docx as _docx

    token = "PANJIT-TABLE-TOKEN-771"
    doc = _docx.Document()
    table = doc.add_table(rows=1, cols=1)
    table.cell(0, 0).text = token
    src = tmp_path / "table_only.docx"
    doc.save(str(src))

    sample = _sample_file_text(src)

    assert token in sample, f"expected table-cell token in sample, got: {sample!r}"


# ---------------------------------------------------------------------------
# AC-3: .pptx table/graphic-frame sampling reads table cell text
# ---------------------------------------------------------------------------

def test_sample_file_text_pptx_table_includes_cell_text(tmp_path):
    """AC-3 (unit): a .pptx whose only shape is a table (a GraphicFrame, no
    has_text_frame shape) still yields a sample containing the table-cell
    token -- proving the sampler reads table cells, not just text frames."""
    from pptx import Presentation
    from pptx.util import Inches

    token = "PANJIT-PPTX-TABLE-TOKEN-556"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    graphic_frame = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
    graphic_frame.table.cell(0, 0).text = token
    src = tmp_path / "table_only.pptx"
    prs.save(str(src))

    sample = _sample_file_text(src)

    assert token in sample, f"expected table-cell token in sample, got: {sample!r}"


# ---------------------------------------------------------------------------
# AC-4: a skipped/failed sample is visible at INFO level with a reason
# ---------------------------------------------------------------------------

def test_detect_document_context_logs_info_reason_on_exception():
    """AC-4 (unit): _detect_document_context's exception handler routes the
    swallow reason through the `log` callback -- the channel that actually
    reaches translator.log in production (JobLogger.log -> the
    "TranslateTool" logger -> its RotatingFileHandler; see
    app/backend/utils/logging_utils.py). The module-level `logger` alone is
    NOT wired to any production handler, so asserting only against it (via
    caplog on "app.backend.processors.orchestrator") would pass even if the
    message never reached a real user-visible log -- the same
    wrong-boundary defect this change repairs. Capturing what is actually
    passed to `log` proves delivery."""
    client = MagicMock(spec=["complete"])
    client.complete.side_effect = RuntimeError("cloud provider unreachable")
    logged: list = []

    result = _detect_document_context(client, "some sample text", log=logged.append)

    assert result == ""
    assert any("cloud provider unreachable" in m for m in logged), (
        f"expected the log(...) callback to receive the failure reason, got: {logged!r}"
    )


def test_detect_document_context_logs_info_reason_when_provider_call_fails():
    """AC-4 (unit): client.complete() returning ok=False (a failed call with
    NO exception) must not silently fall through to an indistinguishable
    empty return -- it is a distinct, nameable failure reason routed through
    `log`."""
    client = MagicMock(spec=["complete"])
    client.complete.return_value = (False, "")
    logged: list = []

    result = _detect_document_context(client, "some sample text", log=logged.append)

    assert result == ""
    assert any("provider call failed" in m for m in logged), (
        f"expected log(...) to receive a distinct 'provider call failed' reason, got: {logged!r}"
    )


def test_detect_document_context_logs_info_reason_when_summary_empty():
    """AC-4 (unit): client.complete() returning ok=True with an
    empty/whitespace-only summary is a different failure mode than a failed
    call and must state a distinct reason, routed through `log`."""
    client = MagicMock(spec=["complete"])
    client.complete.return_value = (True, "   ")
    logged: list = []

    result = _detect_document_context(client, "some sample text", log=logged.append)

    assert result == ""
    assert any("empty summary" in m for m in logged), (
        f"expected log(...) to receive a distinct 'empty summary' reason, got: {logged!r}"
    )


def test_process_files_logs_info_reason_when_sample_empty(tmp_path, caplog):
    """AC-4 (unit): when the context-detection gates are open but the sample
    is falsy, process_files emits an INFO skip line naming the file and
    reason through the SAME channel JobLogger.log() uses in production
    (JobLogger.log -> the "TranslateTool" logger -> the RotatingFileHandler
    that writes translator.log; app/backend/utils/logging_utils.py). This
    test wires `log` to that real logger (mirroring JobLogger.log's own
    behavior) and asserts via caplog against the "TranslateTool" logger name
    -- proving the message reaches the actual production log channel, not
    merely a module-level logger with no attached handler.
    _detect_document_context must NOT be called."""
    from app.backend.utils import logging_utils

    src = tmp_path / "unreadable.docx"
    src.write_bytes(b"not-a-real-docx-file-corrupt-bytes")
    out_dir = tmp_path / "out"
    out_dir.mkdir()

    def _log_via_translate_tool_logger(message: str) -> None:
        logging_utils.logger.info(message)

    with patch("app.backend.processors.orchestrator.translate_docx", return_value=False), \
         patch("app.backend.processors.orchestrator._detect_document_context") as mock_detect, \
         caplog.at_level(logging.INFO, logger="TranslateTool"):
        result = process_files(
            files=[src],
            output_dir=out_dir,
            targets=["French"],
            src_lang="English",
            include_headers_shapes_via_com=False,
            ollama_model="qwen3.5:9b",
            model_type="general",
            system_prompt="",
            profile_id="general",
            provider_id=None,
            log=_log_via_translate_tool_logger,
        )

    _, _, stopped, _client, _term, _winning = result
    assert stopped is False
    mock_detect.assert_not_called()
    assert any(
        record.name == "TranslateTool"
        and record.levelno == logging.INFO
        and "unreadable.docx" in record.message
        and "empty sample" in record.message
        for record in caplog.records
    ), "expected the skip line naming the file and reason to reach the TranslateTool logger channel"


# ---------------------------------------------------------------------------
# AC-5: successful detection is visible at INFO level
# ---------------------------------------------------------------------------

def test_detect_document_context_logs_info_on_success():
    """AC-5 (unit): a successful detection routes [CONTEXT] Detected: through
    the `log` callback -- the channel that reaches translator.log in
    production -- not only a module-level logger call invisible to any real
    handler."""
    client = MagicMock(spec=["complete"])
    client.complete.return_value = (True, "A purchase order document.")
    logged: list = []

    result = _detect_document_context(client, "some sample text", log=logged.append)

    assert result == "A purchase order document."
    assert any("[CONTEXT] Detected:" in m for m in logged), (
        f"expected the log(...) callback to receive the detection line, got: {logged!r}"
    )


# ---------------------------------------------------------------------------
# AC-6: graceful degradation -- sampler failure never aborts the job and
# never injects a "Document context:" preamble
# ---------------------------------------------------------------------------

@pytest.mark.parametrize(
    "ext,translate_attr",
    [
        (".docx", "translate_docx"),
        (".pptx", "translate_pptx"),
        (".xls", "translate_xlsx_xls"),
    ],
)
def test_sampling_exception_degrades_to_no_preamble_job_completes(tmp_path, ext, translate_attr):
    """AC-6 (data-boundary): a sampler exception for each of .docx/.pptx/.xls
    degrades to an empty sample; the job still completes (stopped is False)
    with no "Document context:" preamble reaching the outgoing request."""
    src = tmp_path / f"corrupt{ext}"
    src.write_bytes(b"\x00\x01garbage-not-a-real-office-file-content")
    out_dir = tmp_path / "out"
    out_dir.mkdir()

    captured: dict = {}

    def _stub(src_path, out_path, targets, src_lang, client, **kwargs):
        captured["system_prompt"] = client.system_prompt
        return False

    with patch(f"app.backend.processors.orchestrator.{translate_attr}", side_effect=_stub), \
         patch("app.backend.processors.orchestrator.is_libreoffice_available", return_value=False), \
         patch("app.backend.processors.orchestrator.DYNAMIC_SCENARIO_STRATEGY_ENABLED", False):
        result = process_files(
            files=[src],
            output_dir=out_dir,
            targets=["French"],
            src_lang="English",
            include_headers_shapes_via_com=False,
            ollama_model="qwen3.5:9b",
            model_type="general",
            system_prompt="Base prompt.",
            profile_id="general",
            provider_id=None,
        )

    _, _, stopped, _client, _term, _winning = result
    assert stopped is False
    assert "Document context:" not in captured.get("system_prompt", "")


# ---------------------------------------------------------------------------
# AC-7: sampler-side .xls conversion does not double-invoke LibreOffice
# ---------------------------------------------------------------------------

def test_xls_sampling_does_not_double_convert_via_libreoffice(tmp_path):
    """AC-7 (integration): with translate_xlsx_xls stubbed (so the
    processor-side conversion never fires), a single .xls file through one
    process_files() run must invoke subprocess.Popen exactly once -- the
    sampler-side conversion, and only that."""
    src = tmp_path / "legacy.xls"
    src.write_bytes(b"dummy-xls-bytes-never-read-popen-is-mocked")
    out_dir = tmp_path / "out"
    out_dir.mkdir()

    fake_popen_cls = _make_xlsx_writing_popen("PANJIT-XLS-TOKEN-909")
    popen_calls: list = []

    def _counting_popen(cmd, **kwargs):
        popen_calls.append(cmd)
        return fake_popen_cls(cmd, **kwargs)

    with patch("shutil.which", side_effect=_make_soffice_available_patch()), \
         patch("subprocess.Popen", side_effect=_counting_popen), \
         patch("app.backend.processors.orchestrator.translate_xlsx_xls", return_value=False), \
         patch.object(OllamaClient, "_call_ollama", return_value=(True, "A spreadsheet document.")), \
         patch("app.backend.processors.orchestrator.DYNAMIC_SCENARIO_STRATEGY_ENABLED", False):
        result = process_files(
            files=[src],
            output_dir=out_dir,
            targets=["French"],
            src_lang="English",
            include_headers_shapes_via_com=False,
            ollama_model="qwen3.5:9b",
            model_type="general",
            system_prompt="",
            profile_id="general",
            provider_id=None,
        )

    _, _, stopped, _client, _term, _winning = result
    assert stopped is False
    assert len(popen_calls) == 1, (
        f"expected exactly one sampler-side LibreOffice conversion, got {len(popen_calls)}"
    )


# ---------------------------------------------------------------------------
# AC-8: both a legacy .xls and a table-only .docx emit [CONTEXT] Detected:
# ---------------------------------------------------------------------------

def test_legacy_xls_and_table_only_docx_both_emit_context_detected(tmp_path, caplog):
    """AC-8 (integration): a job processing both a legacy .xls and a
    table-only .docx emits [CONTEXT] Detected: for each file, through the
    SAME channel JobLogger.log() uses in production (JobLogger.log -> the
    "TranslateTool" logger -> the RotatingFileHandler that writes
    translator.log). This is the user's own success condition -- the line
    they will look for lives in translator.log, not in a module-level
    logger with no attached handler -- so `log` is wired to the real
    "TranslateTool" logger here and the assertion targets that logger name."""
    import docx as _docx

    from app.backend.utils import logging_utils

    xls_src = tmp_path / "legacy.xls"
    xls_src.write_bytes(b"dummy-xls-bytes-never-read-popen-is-mocked")

    doc = _docx.Document()
    table = doc.add_table(rows=1, cols=1)
    table.cell(0, 0).text = "PANJIT-TABLE-TOKEN-321"
    docx_src = tmp_path / "table_only.docx"
    doc.save(str(docx_src))

    out_dir = tmp_path / "out"
    out_dir.mkdir()

    def _log_via_translate_tool_logger(message: str) -> None:
        logging_utils.logger.info(message)

    with patch("shutil.which", side_effect=_make_soffice_available_patch()), \
         patch("subprocess.Popen", side_effect=_make_xlsx_writing_popen("PANJIT-XLS-TOKEN-654")), \
         patch("app.backend.processors.orchestrator.translate_xlsx_xls", return_value=False), \
         patch("app.backend.processors.orchestrator.translate_docx", return_value=False), \
         patch.object(OllamaClient, "_call_ollama", return_value=(True, "A document.")), \
         patch("app.backend.processors.orchestrator.DYNAMIC_SCENARIO_STRATEGY_ENABLED", False), \
         caplog.at_level(logging.INFO, logger="TranslateTool"):
        result = process_files(
            files=[xls_src, docx_src],
            output_dir=out_dir,
            targets=["French"],
            src_lang="English",
            include_headers_shapes_via_com=False,
            ollama_model="qwen3.5:9b",
            model_type="general",
            system_prompt="",
            profile_id="general",
            provider_id=None,
            log=_log_via_translate_tool_logger,
        )

    _, _, stopped, _client, _term, _winning = result
    assert stopped is False
    # Count only records on the "TranslateTool" logger — the one that owns the
    # RotatingFileHandler writing translator.log. This is the channel the user
    # actually reads, and AC-8 is stated in terms of it. Counting unfiltered
    # caplog.records would also count the orchestrator module logger's additive
    # `logger.info(...)`, which reaches no production handler, and AC-8 would
    # pass even if translator.log never showed the line again.
    detected_count = sum(
        1
        for r in caplog.records
        if r.name == "TranslateTool" and "[CONTEXT] Detected:" in r.message
    )
    assert detected_count == 2, (
        f"expected both files to emit [CONTEXT] Detected:, got {detected_count}"
    )
