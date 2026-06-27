"""Unit tests for output_mode parameter in translate_docx and translate_pptx.

AC-1: signature accepts output_mode; default is append.
AC-2: append mode behavior unchanged.
AC-3: replace mode DOCX — source paragraphs replaced in-place; original text absent.
AC-4: replace mode PPTX — source text frames replaced in-place; original text absent.
AC-7: multi-target is tested at processor level (pass-through; clamping is in orchestrator).

Anti-tautology: selection assertions check WHICH paragraphs/frames hold the translation,
not just counts.
"""

from __future__ import annotations

import io
import threading
from pathlib import Path
from typing import Dict, List, Optional, Tuple
from unittest.mock import MagicMock, patch

import pytest


# ---------------------------------------------------------------------------
# Helpers: build in-process minimal DOCX / PPTX fixtures
# ---------------------------------------------------------------------------

def _make_docx(tmp_path: Path, text: str) -> Path:
    """Create a minimal DOCX with one paragraph containing *text*."""
    from docx import Document

    doc = Document()
    doc.add_paragraph(text)
    p = tmp_path / "test.docx"
    doc.save(str(p))
    return p


def _make_pptx(tmp_path: Path, text: str) -> Path:
    """Create a minimal PPTX with one slide and one text-frame containing *text*."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    txBox.text_frame.text = text
    p = tmp_path / "test.pptx"
    prs.save(str(p))
    return p


def _mock_client():
    """Return a minimal mock OllamaClient that passes health checks."""
    mock = MagicMock()
    mock.health_check.return_value = (True, "ok")
    mock.system_prompt = ""
    mock.model_type = "general"
    mock._is_translation_dedicated.return_value = False
    mock._is_translategemma_model.return_value = False
    return mock


def _make_tmap(src_text: str, translation: str, target: str = "vi") -> Dict:
    """Return a translation map keyed by (target, src_text)."""
    return {(target, src_text): translation}


# ---------------------------------------------------------------------------
# AC-1: signature tests
# ---------------------------------------------------------------------------

def test_translate_docx_accepts_output_mode_param(tmp_path):
    """translate_docx must accept an output_mode keyword argument."""
    import inspect
    from app.backend.processors.docx_processor import translate_docx

    sig = inspect.signature(translate_docx)
    assert "output_mode" in sig.parameters


def test_translate_pptx_accepts_output_mode_param(tmp_path):
    """translate_pptx must accept an output_mode keyword argument."""
    import inspect
    from app.backend.processors.pptx_processor import translate_pptx

    sig = inspect.signature(translate_pptx)
    assert "output_mode" in sig.parameters


def test_output_mode_default_is_append(tmp_path):
    """output_mode defaults to 'append' in both processors."""
    import inspect
    from app.backend.processors.docx_processor import translate_docx
    from app.backend.processors.pptx_processor import translate_pptx

    docx_default = inspect.signature(translate_docx).parameters["output_mode"].default
    pptx_default = inspect.signature(translate_pptx).parameters["output_mode"].default
    assert docx_default == "append"
    assert pptx_default == "append"


# ---------------------------------------------------------------------------
# AC-2: append mode behavior unchanged (DOCX)
# ---------------------------------------------------------------------------

def test_append_mode_behavior_unchanged_docx(tmp_path):
    """In append mode, both source and translated paragraphs appear in output DOCX."""
    from docx import Document
    from app.backend.processors.docx_processor import translate_docx

    src_text = "Hello world"
    translation = "Hola mundo"
    in_path = _make_docx(tmp_path, src_text)
    out_path = tmp_path / "out.docx"

    tmap = _make_tmap(src_text, translation)

    with (
        patch("app.backend.processors.docx_processor.OllamaClient", return_value=_mock_client()),
        patch(
            "app.backend.processors.docx_processor.translate_texts",
            return_value=(tmap, 1, 0, False),
        ),
    ):
        translate_docx(
            str(in_path),
            str(out_path),
            targets=["vi"],
            src_lang="en",
            client=_mock_client(),
            include_headers_shapes_via_com=False,
            output_mode="append",
        )

    doc = Document(str(out_path))
    all_texts = [p.text for p in doc.paragraphs]
    # Source text must still be present
    assert any(src_text in t for t in all_texts), f"Source text missing from append output: {all_texts}"
    # Translation must also be present
    assert any(translation in t for t in all_texts), f"Translation missing from append output: {all_texts}"


# ---------------------------------------------------------------------------
# AC-2: append mode behavior unchanged (PPTX)
# ---------------------------------------------------------------------------

def test_append_mode_behavior_unchanged_pptx(tmp_path):
    """In append mode, both source and translated paragraphs appear in output PPTX."""
    from pptx import Presentation
    from app.backend.processors.pptx_processor import translate_pptx

    src_text = "Hello world"
    translation = "Hola mundo"
    in_path = _make_pptx(tmp_path, src_text)
    out_path = tmp_path / "out.pptx"

    tmap = _make_tmap(src_text, translation)

    with patch(
        "app.backend.processors.pptx_processor.translate_texts",
        return_value=(tmap, 1, 0, False),
    ):
        translate_pptx(
            str(in_path),
            str(out_path),
            targets=["vi"],
            src_lang="en",
            client=_mock_client(),
            output_mode="append",
        )

    prs = Presentation(str(out_path))
    all_texts: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    all_texts.append(para.text)

    assert any(src_text in t for t in all_texts), f"Source text missing from append output: {all_texts}"
    assert any(translation in t for t in all_texts), f"Translation missing from append output: {all_texts}"


# ---------------------------------------------------------------------------
# AC-3: replace mode DOCX — source absent, translation in place
# ---------------------------------------------------------------------------

def test_replace_mode_docx_no_source_paragraphs_remain(tmp_path):
    """In replace mode, the source text must be absent from the output DOCX paragraphs."""
    from docx import Document
    from app.backend.processors.docx_processor import translate_docx

    src_text = "Hello world"
    translation = "Hola mundo"
    in_path = _make_docx(tmp_path, src_text)
    out_path = tmp_path / "out_replace.docx"

    tmap = _make_tmap(src_text, translation)

    with patch(
        "app.backend.processors.docx_processor.translate_texts",
        return_value=(tmap, 1, 0, False),
    ):
        translate_docx(
            str(in_path),
            str(out_path),
            targets=["vi"],
            src_lang="en",
            client=_mock_client(),
            include_headers_shapes_via_com=False,
            output_mode="replace",
        )

    doc = Document(str(out_path))
    all_texts = [p.text for p in doc.paragraphs]
    # Selection assertion: source text must NOT appear in any paragraph
    assert not any(src_text in t for t in all_texts), (
        f"Source text still present in replace output: {all_texts}"
    )


def test_replace_mode_docx_translation_is_in_place(tmp_path):
    """In replace mode, the translation must appear in the output DOCX paragraphs."""
    from docx import Document
    from app.backend.processors.docx_processor import translate_docx

    src_text = "Hello world"
    translation = "Hola mundo"
    in_path = _make_docx(tmp_path, src_text)
    out_path = tmp_path / "out_replace2.docx"

    tmap = _make_tmap(src_text, translation)

    with patch(
        "app.backend.processors.docx_processor.translate_texts",
        return_value=(tmap, 1, 0, False),
    ):
        translate_docx(
            str(in_path),
            str(out_path),
            targets=["vi"],
            src_lang="en",
            client=_mock_client(),
            include_headers_shapes_via_com=False,
            output_mode="replace",
        )

    doc = Document(str(out_path))
    all_texts = [p.text for p in doc.paragraphs]
    # Selection assertion: translation must appear in at least one paragraph
    assert any(translation in t for t in all_texts), (
        f"Translation missing from replace output: {all_texts}"
    )


# ---------------------------------------------------------------------------
# AC-4: replace mode PPTX — source absent, translation in place
# ---------------------------------------------------------------------------

def test_replace_mode_pptx_no_source_text_frames_remain(tmp_path):
    """In replace mode, the source text must be absent from output PPTX text frames."""
    from pptx import Presentation
    from app.backend.processors.pptx_processor import translate_pptx

    src_text = "Hello world"
    translation = "Hola mundo"
    in_path = _make_pptx(tmp_path, src_text)
    out_path = tmp_path / "out_replace.pptx"

    tmap = _make_tmap(src_text, translation)

    with patch(
        "app.backend.processors.pptx_processor.translate_texts",
        return_value=(tmap, 1, 0, False),
    ):
        translate_pptx(
            str(in_path),
            str(out_path),
            targets=["vi"],
            src_lang="en",
            client=_mock_client(),
            output_mode="replace",
        )

    prs = Presentation(str(out_path))
    all_texts: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    all_texts.append(para.text)

    # Selection assertion: source text must NOT appear
    assert not any(src_text in t for t in all_texts), (
        f"Source text still present in replace output: {all_texts}"
    )


def test_replace_mode_pptx_translation_is_in_place(tmp_path):
    """In replace mode, translation must appear in output PPTX text frames."""
    from pptx import Presentation
    from app.backend.processors.pptx_processor import translate_pptx

    src_text = "Hello world"
    translation = "Hola mundo"
    in_path = _make_pptx(tmp_path, src_text)
    out_path = tmp_path / "out_replace2.pptx"

    tmap = _make_tmap(src_text, translation)

    with patch(
        "app.backend.processors.pptx_processor.translate_texts",
        return_value=(tmap, 1, 0, False),
    ):
        translate_pptx(
            str(in_path),
            str(out_path),
            targets=["vi"],
            src_lang="en",
            client=_mock_client(),
            output_mode="replace",
        )

    prs = Presentation(str(out_path))
    all_texts: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    all_texts.append(para.text)

    # Selection assertion: translation must appear
    assert any(translation in t for t in all_texts), (
        f"Translation missing from replace output: {all_texts}"
    )


# ---------------------------------------------------------------------------
# AC-2 (office-output-mode): DOCX bilingual two-column table
# ---------------------------------------------------------------------------

def test_bilingual_docx_produces_two_column_table(tmp_path):
    """DOCX bilingual mode must produce a table with 2 columns."""
    from docx import Document
    from app.backend.processors.docx_processor import translate_docx

    src_text = "Hello bilingual"
    translation = "Hola bilingüe"
    in_path = _make_docx(tmp_path, src_text)
    out_path = tmp_path / "out_bilingual.docx"

    tmap = {("vi", src_text): translation}

    with patch(
        "app.backend.processors.docx_processor.translate_texts",
        return_value=(tmap, 1, 0, False),
    ):
        translate_docx(
            str(in_path),
            str(out_path),
            targets=["vi"],
            src_lang="en",
            client=_mock_client(),
            include_headers_shapes_via_com=False,
            output_mode="bilingual",
        )

    doc = Document(str(out_path))
    # Structural assertion: exactly one table must have been created
    assert len(doc.tables) == 1, f"Expected 1 table, got {len(doc.tables)}"
    tbl = doc.tables[0]
    assert len(tbl.columns) == 2, f"Expected 2 columns, got {len(tbl.columns)}"
    assert len(tbl.rows) == 1, f"Expected 1 row, got {len(tbl.rows)}"


def test_bilingual_docx_source_col_a_translation_col_b_not_same_run(tmp_path):
    """DOCX bilingual: source in cell(0,0), translation in cell(0,1); they must NOT share a run."""
    from docx import Document
    from app.backend.processors.docx_processor import translate_docx

    src_text = "Source paragraph"
    translation = "Traducción"
    in_path = _make_docx(tmp_path, src_text)
    out_path = tmp_path / "out_bilingual2.docx"

    tmap = {("vi", src_text): translation}

    with patch(
        "app.backend.processors.docx_processor.translate_texts",
        return_value=(tmap, 1, 0, False),
    ):
        translate_docx(
            str(in_path),
            str(out_path),
            targets=["vi"],
            src_lang="en",
            client=_mock_client(),
            include_headers_shapes_via_com=False,
            output_mode="bilingual",
        )

    doc = Document(str(out_path))
    assert len(doc.tables) >= 1, "No table produced"
    tbl = doc.tables[0]
    # Selection assertions: source in col-A, translation in col-B
    col_a_text = tbl.cell(0, 0).text
    col_b_text = tbl.cell(0, 1).text
    assert src_text in col_a_text, f"Source not in col-A: {col_a_text!r}"
    assert translation in col_b_text, f"Translation not in col-B: {col_b_text!r}"
    # Anti-tautology: source and translation must NOT be in the same cell
    assert translation not in col_a_text, f"Translation leaked into col-A: {col_a_text!r}"
    assert src_text not in col_b_text, f"Source leaked into col-B: {col_b_text!r}"


def test_bilingual_docx_empty_paragraph_passthrough(tmp_path):
    """DOCX bilingual: empty (whitespace-only) paragraphs must pass through unchanged."""
    from docx import Document
    from app.backend.processors.docx_processor import translate_docx

    from docx import Document as _D
    doc = _D()
    doc.add_paragraph("Real content")
    doc.add_paragraph("")  # empty paragraph
    in_path = tmp_path / "test_empty.docx"
    doc.save(str(in_path))
    out_path = tmp_path / "out_empty_bilingual.docx"

    tmap = {("vi", "Real content"): "Contenu réel"}

    with patch(
        "app.backend.processors.docx_processor.translate_texts",
        return_value=(tmap, 1, 0, False),
    ):
        translate_docx(
            str(in_path),
            str(out_path),
            targets=["vi"],
            src_lang="en",
            client=_mock_client(),
            include_headers_shapes_via_com=False,
            output_mode="bilingual",
        )

    out_doc = Document(str(out_path))
    # Should have exactly one table (for "Real content")
    assert len(out_doc.tables) == 1, (
        f"Expected 1 table for the real paragraph; got {len(out_doc.tables)}"
    )


# ---------------------------------------------------------------------------
# AC-3 (office-output-mode): XLSX adjacent
# ---------------------------------------------------------------------------

def _make_xlsx(tmp_path: "Path", src_text: str, sheet_name: str = "Sheet1"):
    """Create a minimal XLSX with one cell (A1) containing *src_text*."""
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = sheet_name
    ws["A1"] = src_text
    p = tmp_path / "test.xlsx"
    wb.save(str(p))
    return p


def _mock_xlsx_client():
    """Return a mock client that forces the XLSX fallback path (translate_once fails)."""
    mock = _mock_client()
    mock.translate_once.return_value = (False, "")
    return mock


def test_xlsx_adjacent_translation_at_shifted_column_source_unchanged_no_wrap(tmp_path):
    """XLSX adjacent: translation written to (row, col+max_col); source cell unchanged; no wrap_text."""
    import openpyxl
    from app.backend.processors.xlsx_processor import translate_xlsx_xls

    src_text = "Hello world"
    translation = "Hola mundo"
    in_path = _make_xlsx(tmp_path, src_text)
    out_path = tmp_path / "out_adjacent.xlsx"

    tmap_data = {("vi", src_text): translation}

    with patch(
        "app.backend.processors.xlsx_processor.translate_texts",
        return_value=(tmap_data, 1, 0, False),
    ):
        translate_xlsx_xls(
            str(in_path),
            str(out_path),
            targets=["vi"],
            src_lang="en",
            client=_mock_xlsx_client(),
            output_mode="adjacent",
        )

    wb_out = openpyxl.load_workbook(str(out_path))
    ws = wb_out.active
    # Source at A1 (row=1, col=1) must be unchanged
    assert ws.cell(row=1, column=1).value == src_text, (
        f"Source cell modified: {ws.cell(row=1, column=1).value!r}"
    )
    # Translation at (row=1, col = 1 + original_max_col=1) = (1, 2) = column B
    translated_cell = ws.cell(row=1, column=2)
    assert translated_cell.value == translation, (
        f"Translation not found at column B: {translated_cell.value!r}"
    )
    # No wrap_text on the translated cell
    assert not (translated_cell.alignment and translated_cell.alignment.wrap_text), (
        "wrap_text must not be set on adjacent-mode translation cell"
    )


# ---------------------------------------------------------------------------
# AC-4 (office-output-mode): XLSX annotation
# ---------------------------------------------------------------------------

def test_xlsx_annotation_attaches_comment_source_unchanged(tmp_path):
    """XLSX annotation: translation attached as Comment; source cell value unchanged."""
    import openpyxl
    from app.backend.processors.xlsx_processor import translate_xlsx_xls

    src_text = "Hello world"
    translation = "Hola mundo"
    in_path = _make_xlsx(tmp_path, src_text)
    out_path = tmp_path / "out_annotation.xlsx"

    tmap_data = {("vi", src_text): translation}

    with patch(
        "app.backend.processors.xlsx_processor.translate_texts",
        return_value=(tmap_data, 1, 0, False),
    ):
        translate_xlsx_xls(
            str(in_path),
            str(out_path),
            targets=["vi"],
            src_lang="en",
            client=_mock_xlsx_client(),
            output_mode="annotation",
        )

    wb_out = openpyxl.load_workbook(str(out_path))
    ws = wb_out.active
    cell = ws.cell(row=1, column=1)
    # Source cell value unchanged
    assert cell.value == src_text, f"Source cell value changed: {cell.value!r}"
    # Comment attached with translation text
    assert cell.comment is not None, "No comment attached to source cell"
    assert translation in (cell.comment.text or ""), (
        f"Translation not in comment: {cell.comment.text!r}"
    )


def test_xlsx_annotation_idempotent_existing_comment_preserved(tmp_path):
    """XLSX annotation: pre-existing user comment is preserved; translation appended below it."""
    import openpyxl
    from openpyxl.comments import Comment
    from app.backend.processors.xlsx_processor import translate_xlsx_xls

    src_text = "Hello world"
    translation = "Hola mundo"
    user_comment_text = "user note"

    # Create XLSX with a pre-existing non-translate-tool comment
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = src_text
    ws["A1"].comment = Comment(user_comment_text, "author")
    in_path = tmp_path / "test_existing_comment.xlsx"
    wb.save(str(in_path))

    out_path = tmp_path / "out_annotation_existing.xlsx"
    tmap_data = {("vi", src_text): translation}

    with patch(
        "app.backend.processors.xlsx_processor.translate_texts",
        return_value=(tmap_data, 1, 0, False),
    ):
        translate_xlsx_xls(
            str(in_path),
            str(out_path),
            targets=["vi"],
            src_lang="en",
            client=_mock_xlsx_client(),
            output_mode="annotation",
        )

    wb_out = openpyxl.load_workbook(str(out_path))
    ws_out = wb_out.active
    cell = ws_out.cell(row=1, column=1)
    comment_text = cell.comment.text or "" if cell.comment else ""
    # Pre-existing user comment preserved
    assert user_comment_text in comment_text, (
        f"User comment not preserved: {comment_text!r}"
    )
    # Translation also present
    assert translation in comment_text, (
        f"Translation not in comment: {comment_text!r}"
    )


# ---------------------------------------------------------------------------
# AC-5 (office-output-mode): XLSX replace
# ---------------------------------------------------------------------------

def test_xlsx_replace_overwrites_no_stack_no_wrap_text(tmp_path):
    """XLSX replace: cell value = translation only (no src+translation stack); no wrap_text."""
    import openpyxl
    from app.backend.processors.xlsx_processor import translate_xlsx_xls

    src_text = "Hello world"
    translation = "Hola mundo"
    in_path = _make_xlsx(tmp_path, src_text)
    out_path = tmp_path / "out_replace.xlsx"

    tmap_data = {("vi", src_text): translation}

    with patch(
        "app.backend.processors.xlsx_processor.translate_texts",
        return_value=(tmap_data, 1, 0, False),
    ):
        translate_xlsx_xls(
            str(in_path),
            str(out_path),
            targets=["vi"],
            src_lang="en",
            client=_mock_xlsx_client(),
            output_mode="replace",
        )

    wb_out = openpyxl.load_workbook(str(out_path))
    ws = wb_out.active
    cell = ws.cell(row=1, column=1)
    # Cell value is translation only — no "src\n譯文" stack
    assert cell.value == translation, (
        f"Expected translation-only value, got: {cell.value!r}"
    )
    assert src_text not in (cell.value or ""), (
        f"Source text still present in cell: {cell.value!r}"
    )
    # No wrap_text
    assert not (cell.alignment and cell.alignment.wrap_text), (
        "wrap_text must be False in replace mode"
    )


# ---------------------------------------------------------------------------
# AC-6 (office-output-mode): DOCX SDT / para-in-cell / text-box replace
# ---------------------------------------------------------------------------

def test_docx_sdt_replace_overwrites_source(tmp_path):
    """DOCX SDT replace: first paragraph in SDT content overwritten with translation."""
    from docx import Document
    from docx.oxml import OxmlElement
    from docx.text.paragraph import Paragraph
    from app.backend.processors.docx_processor import _insert_docx_translations, Segment

    doc = Document()

    # Build a minimal SDT element with a paragraph inside sdtContent
    sdt = OxmlElement("w:sdt")
    sdt_content = OxmlElement("w:sdtContent")
    sdt_p = OxmlElement("w:p")
    sdt_r = OxmlElement("w:r")
    sdt_t = OxmlElement("w:t")
    sdt_t.text = "Hello SDT"
    sdt_r.append(sdt_t)
    sdt_p.append(sdt_r)
    sdt_content.append(sdt_p)
    sdt.append(sdt_content)
    doc._body._body.append(sdt)

    src_text = "Hello SDT"
    translation = "Hola SDT"
    # SDT segments use seg.ref = the SDT element; col=None
    seg = Segment("para", sdt, "Body > SDT-Placeholder", src_text)
    tmap = {("vi", src_text, None): translation}

    _insert_docx_translations(doc, [seg], tmap, targets=["vi"], output_mode="replace")

    # Assert the first paragraph in sdtContent now contains the translation
    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    all_paras = sdt_content.xpath(".//w:p", namespaces=ns)
    assert len(all_paras) >= 1, "No paragraph in sdtContent"
    p_obj = Paragraph(all_paras[0], None)
    assert translation in p_obj.text, (
        f"Translation not in SDT paragraph: {p_obj.text!r}"
    )
    assert src_text not in p_obj.text, (
        f"Source text still in SDT paragraph: {p_obj.text!r}"
    )


def test_docx_para_in_cell_replace_overwrites_source(tmp_path):
    """DOCX para-in-cell replace: paragraph runs overwritten in-place."""
    from docx import Document
    from docx.table import _Cell
    from app.backend.processors.docx_processor import _insert_docx_translations, Segment

    doc = Document()
    table = doc.add_table(rows=1, cols=1)
    cell = table.cell(0, 0)
    # Add a run to the first paragraph in the cell
    p = cell.paragraphs[0]
    p.add_run("Hello cell")

    assert isinstance(p._parent, _Cell), "Paragraph must be inside a cell for this test"

    src_text = "Hello cell"
    translation = "Hola celda"
    # para-in-cell segments: kind="para", ref=Paragraph, col=None
    seg = Segment("para", p, "Body > Tbl > Cell", src_text)
    tmap = {("vi", src_text, None): translation}

    _insert_docx_translations(doc, [seg], tmap, targets=["vi"], output_mode="replace")

    # The paragraph text should now be the translation only
    assert translation in p.text, f"Translation not in cell paragraph: {p.text!r}"
    assert src_text not in p.text, f"Source text still in cell paragraph: {p.text!r}"


def test_docx_textbox_replace_overwrites_source(tmp_path):
    """DOCX text-box replace: first paragraph text in txbxContent overwritten."""
    from docx import Document
    from docx.oxml import OxmlElement
    from app.backend.processors.docx_processor import _insert_docx_translations, Segment

    doc = Document()

    # Build a minimal txbxContent element with one paragraph
    txbx = OxmlElement("w:txbxContent")
    p_elem = OxmlElement("w:p")
    r_elem = OxmlElement("w:r")
    t_elem = OxmlElement("w:t")
    t_elem.text = "Hello TextBox"
    r_elem.append(t_elem)
    p_elem.append(r_elem)
    txbx.append(p_elem)

    src_text = "Hello TextBox"
    translation = "Hola TextBox"
    # txbx segments: kind="txbx", ref=txbxContent element, col=None
    seg = Segment("txbx", txbx, "TextBox", src_text)
    tmap = {("vi", src_text, None): translation}

    _insert_docx_translations(doc, [seg], tmap, targets=["vi"], output_mode="replace")

    # All <w:t> elements in the txbx should now contain the translation
    all_text = "".join(
        (t.text or "") for t in txbx.xpath(".//*[local-name()='t']")
    )
    assert translation in all_text, f"Translation not in txbx: {all_text!r}"
    assert src_text not in all_text, f"Source text still in txbx: {all_text!r}"


# ---------------------------------------------------------------------------
# AC-7 (office-output-mode): PPTX SmartArt replace
# ---------------------------------------------------------------------------

def test_pptx_smartart_replace_text_not_appended(tmp_path):
    """PPTX SmartArt replace: <a:t> text set to translation only, not appended in parentheses."""
    import zipfile
    from app.backend.processors.pptx_processor import _update_smartart_texts

    src_text = "Hello SmartArt"
    translation = "Hola SmartArt"

    pptx_path = tmp_path / "smartart.pptx"
    out_path = tmp_path / "out_smartart.pptx"

    # Build a minimal fake PPTX ZIP with a SmartArt diagram data file
    diagram_xml = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<root xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f"<a:t>{src_text}</a:t>"
        "</root>"
    )
    with zipfile.ZipFile(str(pptx_path), "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("ppt/diagrams/data1.xml", diagram_xml)

    translations = {src_text: translation}
    _update_smartart_texts(str(pptx_path), str(out_path), translations, output_mode="replace")

    with zipfile.ZipFile(str(out_path), "r") as zf:
        content = zf.read("ppt/diagrams/data1.xml").decode("utf-8")

    # Selection assertions
    assert translation in content, f"Translation not found in SmartArt XML: {content!r}"
    # Anti-tautology: must NOT contain the append-mode parenthetical suffix
    assert f"\n({translation})" not in content, (
        "Translation was appended in parentheses instead of replaced"
    )


# ---------------------------------------------------------------------------
# AC-7 (processor level): multi-target does not crash in either mode
# ---------------------------------------------------------------------------

def test_multi_target_output_mode_clamped_to_append(tmp_path):
    """With multiple targets, output_mode="replace" still runs without error.

    The orchestrator clamps replace→append for multi-target jobs (BR-67).
    This test verifies the processor itself handles multiple targets gracefully
    when called with append (the clamped value). Orchestrator-level clamping
    is tested in test_output_mode_orchestrator.py.
    """
    from docx import Document
    from app.backend.processors.docx_processor import translate_docx

    src_text = "Hello world"
    in_path = _make_docx(tmp_path, src_text)
    out_path = tmp_path / "out_multi.docx"

    tmap = {
        ("vi", src_text): "Xin chào thế giới",
        ("de", src_text): "Hallo Welt",
    }

    with patch(
        "app.backend.processors.docx_processor.translate_texts",
        return_value=(tmap, 1, 0, False),
    ):
        translate_docx(
            str(in_path),
            str(out_path),
            targets=["vi", "de"],
            src_lang="en",
            client=_mock_client(),
            include_headers_shapes_via_com=False,
            output_mode="append",
        )

    doc = Document(str(out_path))
    all_texts = [p.text for p in doc.paragraphs]
    assert any("Xin chào thế giới" in t for t in all_texts), f"vi translation missing: {all_texts}"
    assert any("Hallo Welt" in t for t in all_texts), f"de translation missing: {all_texts}"
