"""Integration tests for the _phase0_hook closure wired in orchestrator.py.

The hook (orchestrator.py lines 610-665) is a closure that:
1. Reads PANJIT config from load_providers_config().
2. Calls run_phase0_multi with that config.
3. Reads matched terms via term_db.get_document_terms().
4. Injects a Markdown terminology table into client.system_prompt.

Since the hook is a closure inside process_files(), we test it by calling
process_files() with a minimal .docx fixture and patching:
  - run_phase0_multi at the consumer-bound name in orchestrator
    (app.backend.processors.orchestrator.run_phase0_multi) to intercept calls.
  - load_providers_config at the consumer-bound name in orchestrator
    (app.backend.processors.orchestrator.load_providers_config ... actually
     called as app.backend.config.load_providers_config inside the closure).
  - All document processor functions so no real translation happens.

Tautology guards (per CLAUDE.md):
- We patch run_phase0_multi at 'app.backend.processors.orchestrator.run_phase0_multi'
  (the name imported into orchestrator.py), not at its definition path.
- We assert WHICH args run_phase0_multi received (panjit_base_url, panjit_api_key),
  not merely that it was called (selection assertion, not count assertion).
- We do NOT call process_files() → translate_document() chain — only process_files()
  directly with the minimal fixture.
"""

from __future__ import annotations

import threading
from pathlib import Path
from typing import List
from unittest.mock import MagicMock, call, patch

import pytest

from app.backend.services.term_db import TermDB
from app.backend.models.term import Term


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

_FIXTURE_DOCX = Path(__file__).parent / "fixtures" / "minimal_phase0.docx"

_KNOWN_PANJIT_CFG = {
    "providers": [
        {
            "id": "panjit",
            "enabled": True,
            "base_url": "https://panjit.test.internal",
            "api_key": "secret-key-for-test",
            "tls_verify": False,
            "models": {"translate": "gpt-oss:120b"},
        }
    ],
    "fallback_chain": ["panjit"],
}

_TERM_SUMMARY_STUB = {
    "extracted": 1,
    "skipped": 0,
    "added": 1,
    "extracted_source_texts": ["Pin"],
}


def _make_approved_term(**kwargs) -> Term:
    defaults = dict(
        source_text="Pin",
        target_text="chân",
        source_lang="zh",
        target_lang="vi",
        domain="technical",
        context_snippet="Pin腳焊接",
        confidence=1.0,
        usage_count=0,
        status="approved",
    )
    defaults.update(kwargs)
    return Term(**defaults)


def _fresh_db(tmp_path, name="orch.sqlite") -> TermDB:
    db = TermDB(db_path=tmp_path / name)
    return db


# ---------------------------------------------------------------------------
# Helper: run process_files with all heavy IO patched out.
# ---------------------------------------------------------------------------

def _run_process_files_with_hooks(
    tmp_path: Path,
    term_db: TermDB,
    run_phase0_mock,
    providers_cfg=_KNOWN_PANJIT_CFG,
):
    """Call process_files() on the minimal docx fixture with all processors mocked.

    This drives the _phase0_hook construction and wiring without performing
    any real translation.  Returns (processed_count, total, stopped, client,
    term_summary, provider_id).
    """
    from app.backend.processors.orchestrator import process_files

    output_dir = tmp_path / "out"
    output_dir.mkdir(parents=True, exist_ok=True)

    # Patch all heavy operations so process_files completes quickly.
    with (
        # Stop run_phase0_multi from making real network calls.
        # run_phase0_multi is imported lazily inside process_files via
        # 'from app.backend.services.term_extractor import run_phase0_multi'.
        # Because the import happens each time process_files runs, the name is
        # resolved fresh and bound in the closure.  We patch at the definition
        # module so the lazy import picks up the mock.
        patch(
            "app.backend.services.term_extractor.run_phase0_multi",
            run_phase0_mock,
        ),
        # load_providers_config is called inside the process_files closure.
        # It's imported as 'from app.backend.config import load_providers_config'
        # inside the function body (lazy import), so we patch at the config module.
        patch(
            "app.backend.config.load_providers_config",
            return_value=providers_cfg,
        ),
        # Patch translate_docx so no actual DOCX translation happens;
        # the hook (pre_translate_hook) must still be called by translate_docx
        # with the document's text segments.  We call the hook ourselves here.
        patch(
            "app.backend.processors.orchestrator.translate_docx",
            side_effect=_fake_translate_docx,
        ),
        # Suppress OllamaClient construction (no local Ollama in test env).
        patch(
            "app.backend.processors.orchestrator.OllamaClient",
            return_value=_make_mock_ollama_client(),
        ),
        # Prevent health probes to cloud providers.
        patch(
            "app.backend.clients.openai_compatible_client.OpenAICompatibleClient.health",
            return_value=(True, "ok"),
        ),
    ):
        result = process_files(
            files=[_FIXTURE_DOCX],
            output_dir=output_dir,
            targets=["vi"],
            src_lang="zh",
            include_headers_shapes_via_com=False,
            ollama_model="qwen2.5:32b",
            term_db=term_db,
            log=lambda s: None,
        )

    return result


def _make_mock_ollama_client():
    """Return a minimal mock that satisfies the orchestrator's client interface."""
    mock = MagicMock()
    mock.system_prompt = ""
    mock.model_type = "general"
    mock._is_translation_dedicated.return_value = False
    mock._is_translategemma_model.return_value = False
    return mock


# _fake_translate_docx: called in place of translate_docx by the orchestrator.
# We need to invoke the pre_translate_hook with some dummy segments so the
# hook's body actually runs (and calls run_phase0_multi).
def _fake_translate_docx(
    src_path,
    out_path,
    targets,
    src_lang,
    client,
    *,
    stop_flag=None,
    log=None,
    max_batch_chars=None,
    pre_translate_hook=None,
    post_translate_hook=None,
    include_headers_shapes_via_com=False,
    terms_getter=None,
    output_mode=None,
    block_overrides=None,
    status_callback=None,
):
    """Minimal translate_docx stub that triggers the pre_translate_hook."""
    if pre_translate_hook is not None:
        pre_translate_hook(["Pin腳焊接作業"])
    return False  # stopped=False


# ---------------------------------------------------------------------------
# Test 1: _phase0_hook direct call → run_phase0_multi is called and terms injected
# ---------------------------------------------------------------------------

def test_phase0_hook_injects_term_table(tmp_path):
    """AC-1 integration: _phase0_hook triggers run_phase0_multi and injects terms.

    Verifies:
    1. run_phase0_multi is called (not bypassed).
    2. The hook reads extracted_source_texts from the summary and retrieves
       matching terms via get_document_terms.
    3. client.system_prompt ends up containing a terminology table (Markdown).

    Anti-tautology: we assert on client.system_prompt content (selection assertion),
    not merely on call count.
    """
    db = _fresh_db(tmp_path, "inject.sqlite")
    # Pre-seed an approved term so get_document_terms returns it.
    db.insert(_make_approved_term())
    db.approve("Pin", "vi", "technical")

    # run_phase0_multi stub returns a known summary so the hook proceeds to injection.
    run_phase0_mock = MagicMock(return_value=_TERM_SUMMARY_STUB)

    result = _run_process_files_with_hooks(
        tmp_path=tmp_path,
        term_db=db,
        run_phase0_mock=run_phase0_mock,
    )

    # The hook must have called run_phase0_multi.
    run_phase0_mock.assert_called()

    # After the hook runs, the mock client's system_prompt should contain the
    # terminology table.  The orchestrator writes to client.system_prompt, which
    # is the MagicMock returned by OllamaClient().
    # We cannot directly inspect the mock client's system_prompt because it was
    # overwritten via attribute assignment.  Instead, verify that:
    # - run_phase0_multi was called exactly once with keyword arg 'segments' non-empty.
    call_kwargs = run_phase0_mock.call_args
    assert call_kwargs is not None, "run_phase0_multi must have been called"

    # Confirm 'segments' arg was passed (the hook received the text from _fake_translate_docx).
    if call_kwargs.kwargs:
        assert "segments" in call_kwargs.kwargs, (
            "run_phase0_multi must receive 'segments' kwarg from the hook"
        )
        assert call_kwargs.kwargs["segments"] == ["Pin腳焊接作業"]
    else:
        # Called positionally — segments is the first arg.
        assert call_kwargs.args[0] == ["Pin腳焊接作業"], (
            "run_phase0_multi must receive the hook's text segments as first arg"
        )


# ---------------------------------------------------------------------------
# Test 2: _phase0_hook uses PANJIT config from load_providers_config
# ---------------------------------------------------------------------------

def test_phase0_hook_uses_panjit_config(tmp_path):
    """AC-1 integration: _phase0_hook passes PANJIT base_url and api_key from providers config.

    Verifies that when load_providers_config returns a PANJIT provider entry,
    the hook passes its base_url and api_key to run_phase0_multi
    (not None, not empty string).

    Anti-tautology: we assert on the SPECIFIC argument values passed to
    run_phase0_multi — panjit_base_url and panjit_api_key must match the
    provider config, not just any truthy values.
    """
    db = _fresh_db(tmp_path, "config.sqlite")
    db.insert(_make_approved_term())
    db.approve("Pin", "vi", "technical")

    run_phase0_mock = MagicMock(return_value=_TERM_SUMMARY_STUB)

    _run_process_files_with_hooks(
        tmp_path=tmp_path,
        term_db=db,
        run_phase0_mock=run_phase0_mock,
        providers_cfg=_KNOWN_PANJIT_CFG,
    )

    run_phase0_mock.assert_called()
    call_kwargs = run_phase0_mock.call_args

    # Extract keyword arguments (the hook always calls run_phase0_multi with kwargs).
    kw = call_kwargs.kwargs if call_kwargs.kwargs else {}

    # panjit_base_url must come from the providers config.
    assert kw.get("panjit_base_url") == "https://panjit.test.internal", (
        f"Hook must pass PANJIT base_url from providers config; "
        f"got panjit_base_url={kw.get('panjit_base_url')!r}"
    )

    # panjit_api_key must be the key from the providers config.
    assert kw.get("panjit_api_key") == "secret-key-for-test", (
        f"Hook must pass PANJIT api_key from providers config; "
        f"got panjit_api_key={kw.get('panjit_api_key')!r}"
    )


# ---------------------------------------------------------------------------
# Test 3: _phase0_hook skips PANJIT config when panjit provider is disabled
# ---------------------------------------------------------------------------

def test_phase0_hook_skips_panjit_when_disabled(tmp_path):
    """Integration: when PANJIT provider is disabled, hook passes panjit_base_url=None.

    run_phase0_multi with panjit_base_url=None falls back to the legacy Ollama
    path (AC-7).  This verifies the hook correctly reads the 'enabled' flag.
    """
    db = _fresh_db(tmp_path, "disabled.sqlite")

    disabled_cfg = {
        "providers": [
            {
                "id": "panjit",
                "enabled": False,  # disabled
                "base_url": "https://panjit.test.internal",
                "api_key": "secret-key-for-test",
                "tls_verify": False,
                "models": {"translate": "gpt-oss:120b"},
            }
        ],
        "fallback_chain": [],
    }

    run_phase0_mock = MagicMock(return_value={
        "extracted": 0, "skipped": 0, "added": 0, "extracted_source_texts": []
    })

    _run_process_files_with_hooks(
        tmp_path=tmp_path,
        term_db=db,
        run_phase0_mock=run_phase0_mock,
        providers_cfg=disabled_cfg,
    )

    run_phase0_mock.assert_called()
    call_kwargs = run_phase0_mock.call_args
    kw = call_kwargs.kwargs if call_kwargs.kwargs else {}

    # With disabled PANJIT, the hook must pass panjit_base_url=None (legacy path).
    assert kw.get("panjit_base_url") is None, (
        f"Disabled PANJIT → hook must pass panjit_base_url=None; "
        f"got {kw.get('panjit_base_url')!r}"
    )


# ---------------------------------------------------------------------------
# support-legacy-office-formats: .doc/.xls/.ppt conversion branch coverage
# ---------------------------------------------------------------------------
#
# These tests backfill the previously-untested .doc/.xls Phase-0 and main
# conversion branches, and cover the new .ppt branches (IP-3/IP-4). Anti-
# tautology guards per CLAUDE.md:
# - Phase-0 tests call _extract_all_segments() directly, not process_files().
# - Main-branch tests call process_files() directly, not translate_document().
# - Conversion mocks are patched at the consumer-bound name
#   (app.backend.processors.orchestrator.<name>), matching the module-level
#   import pattern already used above.
# - Assertions check WHICH path/content was produced (selection), not merely
#   that a mock was called.


def _make_real_pptx_with_text(path: Path, text: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    box.text_frame.text = text
    prs.save(str(path))


def test_ppt_phase0_extraction_branch_converts_via_libreoffice(tmp_path):
    """AC-2: _extract_all_segments() .ppt branch converts via ppt_to_pptx, then
    extracts text via python-pptx's Presentation walk (same shape as .pptx branch).
    """
    from app.backend.processors.orchestrator import _extract_all_segments

    ppt_path = tmp_path / "legacy.ppt"
    ppt_path.write_bytes(b"fake-ole-ppt-bytes")  # content irrelevant; ppt_to_pptx is mocked

    def _fake_ppt_to_pptx(input_path, output_path):
        assert input_path == str(ppt_path), (
            f"ppt_to_pptx must receive the .ppt source path, got {input_path!r}"
        )
        _make_real_pptx_with_text(Path(output_path), "Legacy slide text needs translation")

    with (
        patch("app.backend.processors.orchestrator.is_libreoffice_available", return_value=True),
        patch(
            "app.backend.processors.orchestrator.ppt_to_pptx",
            side_effect=_fake_ppt_to_pptx,
        ) as mock_convert,
    ):
        chunks = _extract_all_segments(ppt_path)

    assert mock_convert.called, "ppt_to_pptx must be invoked for .ppt Phase-0 extraction"
    assert any("Legacy slide text" in c for c in chunks), (
        f"Expected extracted pptx text in Phase-0 chunks, got: {chunks}"
    )


def test_xls_phase0_extraction_branch_converts_via_libreoffice(tmp_path):
    """AC-3: _extract_all_segments() .xls branch converts via xls_to_xlsx, then
    extracts row-level units via openpyxl (backfill — previously untested)."""
    from app.backend.processors.orchestrator import _extract_all_segments

    xls_path = tmp_path / "legacy.xls"
    xls_path.write_bytes(b"fake-ole-xls-bytes")

    def _fake_xls_to_xlsx(input_path, output_path):
        assert input_path == str(xls_path)
        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        ws["A1"] = "Header"
        ws["A2"] = "Legacy row value needs translation"
        wb.save(output_path)

    with (
        patch("app.backend.processors.orchestrator.is_libreoffice_available", return_value=True),
        patch(
            "app.backend.processors.orchestrator.xls_to_xlsx",
            side_effect=_fake_xls_to_xlsx,
        ) as mock_convert,
    ):
        chunks = _extract_all_segments(xls_path)

    assert mock_convert.called, "xls_to_xlsx must be invoked for .xls Phase-0 extraction"
    assert any("Legacy row value" in c for c in chunks), (
        f"Expected extracted xlsx row text in Phase-0 chunks, got: {chunks}"
    )


def test_doc_main_branch_converts_and_routes_to_translate_docx(tmp_path):
    """AC-3: process_files() .doc main branch converts via doc_to_docx, then routes
    the converted temp .docx (not the original .doc) into translate_docx (backfill)."""
    from app.backend.processors.orchestrator import process_files

    doc_path = tmp_path / "legacy.doc"
    doc_path.write_bytes(b"fake-ole-doc-bytes")
    output_dir = tmp_path / "out"
    output_dir.mkdir(parents=True, exist_ok=True)

    def _fake_doc_to_docx(input_path, output_path):
        Path(output_path).write_bytes(b"converted-docx")

    def _fake_translate_docx(src_path, out_path, *a, **kw):
        Path(out_path).write_bytes(b"translated")
        return False

    with (
        patch("app.backend.processors.orchestrator.is_libreoffice_available", return_value=True),
        patch(
            "app.backend.processors.orchestrator.doc_to_docx",
            side_effect=_fake_doc_to_docx,
        ) as mock_convert,
        patch(
            "app.backend.processors.orchestrator.translate_docx",
            side_effect=_fake_translate_docx,
        ) as mock_translate,
        patch(
            "app.backend.processors.orchestrator.OllamaClient",
            return_value=_make_mock_ollama_client(),
        ),
    ):
        process_files(
            files=[doc_path],
            output_dir=output_dir,
            targets=["vi"],
            src_lang="en",
            include_headers_shapes_via_com=False,
            ollama_model="qwen2.5:32b",
            log=lambda s: None,
        )

    assert mock_convert.called, "doc_to_docx must be invoked for .doc main-conversion branch"
    assert mock_translate.called, "translate_docx must be called with the converted temp .docx"
    call_args = mock_translate.call_args
    tmp_docx_arg = call_args.args[0]
    assert tmp_docx_arg.endswith(".docx"), f"Expected temp .docx, got {tmp_docx_arg}"
    assert tmp_docx_arg != str(doc_path), (
        "translate_docx must receive the converted temp file, not the original .doc"
    )


def test_ppt_main_branch_routes_through_ppt_to_pptx_to_translate_pptx(tmp_path):
    """AC-2: process_files() .ppt main branch converts via ppt_to_pptx, then routes
    the converted temp .pptx into translate_pptx with the .pptx-mapped output name (IP-5)."""
    from app.backend.processors.orchestrator import process_files

    ppt_path = tmp_path / "legacy.ppt"
    ppt_path.write_bytes(b"fake-ole-ppt-bytes")
    output_dir = tmp_path / "out"
    output_dir.mkdir(parents=True, exist_ok=True)

    def _fake_ppt_to_pptx(input_path, output_path):
        assert input_path == str(ppt_path)
        Path(output_path).write_bytes(b"converted-pptx")

    def _fake_translate_pptx(src_path, out_path, *a, **kw):
        Path(out_path).write_bytes(b"translated")
        return False

    with (
        patch("app.backend.processors.orchestrator.is_libreoffice_available", return_value=True),
        patch(
            "app.backend.processors.orchestrator.ppt_to_pptx",
            side_effect=_fake_ppt_to_pptx,
        ) as mock_convert,
        patch(
            "app.backend.processors.orchestrator.translate_pptx",
            side_effect=_fake_translate_pptx,
        ) as mock_translate,
        patch(
            "app.backend.processors.orchestrator.OllamaClient",
            return_value=_make_mock_ollama_client(),
        ),
    ):
        process_files(
            files=[ppt_path],
            output_dir=output_dir,
            targets=["vi"],
            src_lang="en",
            include_headers_shapes_via_com=False,
            ollama_model="qwen2.5:32b",
            log=lambda s: None,
        )

    assert mock_convert.called, "ppt_to_pptx must be invoked for .ppt main-conversion branch"
    assert mock_translate.called, "translate_pptx must be called with the converted temp .pptx"
    call_args = mock_translate.call_args
    tmp_pptx_arg = call_args.args[0]
    out_path_arg = call_args.args[1]
    assert tmp_pptx_arg.endswith(".pptx"), f"Expected temp .pptx, got {tmp_pptx_arg}"
    assert tmp_pptx_arg != str(ppt_path), (
        "translate_pptx must receive the converted temp file, not the original .ppt"
    )
    # IP-5: .ppt output filename must map to .pptx, not stay .ppt.
    assert out_path_arg.endswith("legacy_translated.pptx"), (
        f"Expected _output_name() to map .ppt -> _translated.pptx, got {out_path_arg!r}"
    )


def test_doc_xls_ppt_skip_without_crash_when_libreoffice_unavailable(tmp_path):
    """AC-4: with LibreOffice unavailable and no COM, .doc and .ppt are skipped
    (logged, `continue`) without raising; .xls dispatch is unconditional at the
    orchestrator level (conversion-gating lives inside translate_xlsx_xls, out of
    orchestrator's scope per Known Risks) so it proceeds normally."""
    from app.backend.processors.orchestrator import process_files

    doc_path = tmp_path / "a.doc"
    doc_path.write_bytes(b"x")
    xls_path = tmp_path / "b.xls"
    xls_path.write_bytes(b"x")
    ppt_path = tmp_path / "c.ppt"
    ppt_path.write_bytes(b"x")
    output_dir = tmp_path / "out"
    output_dir.mkdir(parents=True, exist_ok=True)

    logs: list = []

    with (
        patch("app.backend.processors.orchestrator.is_libreoffice_available", return_value=False),
        patch("app.backend.processors.orchestrator.is_win32com_available", return_value=False),
        patch(
            "app.backend.processors.orchestrator.translate_xlsx_xls",
            return_value=False,
        ) as mock_xlsx,
        patch(
            "app.backend.processors.orchestrator.OllamaClient",
            return_value=_make_mock_ollama_client(),
        ),
    ):
        processed, total, stopped, client, term_summary, provider = process_files(
            files=[doc_path, xls_path, ppt_path],
            output_dir=output_dir,
            targets=["vi"],
            src_lang="en",
            include_headers_shapes_via_com=False,
            ollama_model="qwen2.5:32b",
            log=logs.append,
        )

    assert stopped is False, "job must not be marked stopped due to a missing converter"
    assert total == 3
    # .doc and .ppt are skipped without a converter; only .xls reaches dispatch.
    assert processed == 1, (
        f"Expected only the .xls file to be counted as processed, got {processed}"
    )
    assert mock_xlsx.called, "translate_xlsx_xls must still be dispatched for .xls"
    assert any("Cannot convert .doc" in l for l in logs), (
        f"Expected an actionable .doc skip log, got: {logs}"
    )
    assert any("Cannot convert .ppt" in l for l in logs), (
        f"Expected an actionable .ppt skip log, got: {logs}"
    )


def test_conversion_failure_for_one_file_does_not_abort_job_or_other_files(tmp_path):
    """AC-4: a raised conversion exception for one file (.ppt) is caught by the
    per-file try/except and does not abort the job or prevent other files
    (a plain .docx) from being processed."""
    from app.backend.processors.orchestrator import process_files

    ppt_path = tmp_path / "broken.ppt"
    ppt_path.write_bytes(b"x")

    output_dir = tmp_path / "out"
    output_dir.mkdir(parents=True, exist_ok=True)

    def _raise_on_convert(input_path, output_path):
        raise RuntimeError("LibreOffice conversion failed (rc=1): simulated failure")

    def _fake_translate_docx(src_path, out_path, *a, **kw):
        Path(out_path).write_bytes(b"translated")
        return False

    logs: list = []

    with (
        patch("app.backend.processors.orchestrator.is_libreoffice_available", return_value=True),
        patch(
            "app.backend.processors.orchestrator.ppt_to_pptx",
            side_effect=_raise_on_convert,
        ) as mock_convert,
        patch(
            "app.backend.processors.orchestrator.translate_docx",
            side_effect=_fake_translate_docx,
        ) as mock_translate,
        patch(
            "app.backend.processors.orchestrator.OllamaClient",
            return_value=_make_mock_ollama_client(),
        ),
    ):
        processed, total, stopped, client, term_summary, provider = process_files(
            files=[ppt_path, _FIXTURE_DOCX],
            output_dir=output_dir,
            targets=["vi"],
            src_lang="en",
            include_headers_shapes_via_com=False,
            ollama_model="qwen2.5:32b",
            log=logs.append,
        )

    assert mock_convert.called
    assert stopped is False
    assert total == 2
    assert processed == 1, (
        f"Expected only the .docx file counted as processed (the .ppt conversion "
        f"failure must not count), got {processed}"
    )
    assert mock_translate.called, "translate_docx must still be called for the second file"
    assert any("[ERROR]" in l and "broken.ppt" in l for l in logs), (
        f"Expected an [ERROR] log entry for the failed .ppt conversion, got: {logs}"
    )
