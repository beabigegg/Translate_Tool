"""Regression tests for pptx-group-shape-collection (BR-116).

Bug-fix lane. Symptom: text and tables nested inside a PowerPoint
``GroupShape`` were never collected/translated/written back because
``GroupShape`` reports ``has_table=False`` / ``has_text_frame=False`` and the
flat ``for shape in slide.shapes`` loop silently skipped it.

AC-1 (``TestGroupTextCollection::test_grouped_textbox_reaches_translate_texts_payload``)
is the named RED reproduction: pre-fix it fails with a behavioral assertion
error (grouped text absent from the captured ``translate_texts`` ``uniq``
payload), not an import/collection error.

Anti-tautology: assertions check the captured outgoing ``uniq`` payload, the
written-back saved PPTX cell/paragraph text at specific (row, col) positions,
or a real call-count spy — never an internal helper's call-wiring alone.
"""

from __future__ import annotations

import gc
import logging
from typing import Dict, List
from unittest.mock import MagicMock, patch

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from app.backend import config
from app.backend.processors.pptx_processor import translate_pptx


# ---------------------------------------------------------------------------
# Fixture helpers (in-test python-pptx construction; no docs/TEST_DOC/ reads)
# ---------------------------------------------------------------------------

def _add_textbox(shapes, text: str, left=1, top=1, width=2, height=1):
    tb = shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tb.text_frame.text = text
    return tb


def _add_table(shapes, rows: int, cols: int, cell_texts: Dict, left=1, top=1, width=3, height=2):
    tbl_shape = shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    for (r, c), text in cell_texts.items():
        tbl_shape.table.cell(r, c).text = text
    return tbl_shape


def _capturing_translate_texts():
    """Return (fake_fn, captured) where captured['uniq'] holds the UNION of
    every outgoing batch payload passed to translate_texts across all calls
    (the text-frame batch AND any per-cell table fallback batch), so a
    fallback call never hides an earlier call's payload."""
    captured: Dict[str, List[str]] = {"uniq": []}

    def _fake(uniq, targets, src_lang, client, **kwargs):
        for s in uniq:
            if s not in captured["uniq"]:
                captured["uniq"].append(s)
        tmap = {(t, s): f"TR:{s}" for t in targets for s in uniq}
        return tmap, len(uniq), 0, False

    return _fake, captured


def _mock_client_for_table_translation():
    """MagicMock client whose flag-off table path uppercases cell content,
    round-tripping through the real table_serializer.serialize()/parse()."""
    client = MagicMock()
    client._build_table_translate_prompt.side_effect = lambda serialized, src, tgt: serialized
    client.translate_once.side_effect = lambda prompt, tgt, src_lang: (True, prompt.upper())
    return client


def _all_text_frame_texts(prs: Presentation) -> List[str]:
    """Flatten every text-frame paragraph text across the whole deck,
    descending into groups (for reading back what was actually written)."""
    out: List[str] = []

    def _walk(shapes):
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                _walk(shape.shapes)
            elif getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        out.append(cell.text)
            elif getattr(shape, "has_text_frame", False):
                out.append(shape.text_frame.text)

    for slide in prs.slides:
        _walk(slide.shapes)
    return out


def _grouped_table_cell(prs: Presentation, r: int, c: int, marker: str):
    """Find the (row, col) cell of the (first) grouped table whose (0,0)
    cell text starts with `marker`, descending into groups."""
    found = []

    def _walk(shapes):
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                _walk(shape.shapes)
            elif getattr(shape, "has_table", False):
                cell00 = shape.table.cell(0, 0).text
                if marker in cell00:
                    found.append(shape.table.cell(r, c).text)

    for slide in prs.slides:
        _walk(slide.shapes)
    return found


# ---------------------------------------------------------------------------
# AC-1 / AC-2: grouped text reaches the translate_texts outgoing payload
# ---------------------------------------------------------------------------

class TestGroupTextCollection:
    def test_grouped_textbox_reaches_translate_texts_payload(self, tmp_path):
        """Named RED reproduction. A 2-shape group's textboxes must reach the
        captured translate_texts `uniq` payload. Pre-fix: GroupShape reports
        has_table=False/has_text_frame=False and is skipped by the flat loop,
        so this FAILS with grouped text absent (behavioral, not import error)."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_textbox(slide.shapes, "Flat sibling text", left=0.2, top=0.2)
        tb1 = _add_textbox(slide.shapes, "Grouped alpha text", left=1, top=2)
        tb2 = _add_textbox(slide.shapes, "Grouped beta text", left=3, top=2)
        slide.shapes.add_group_shape([tb1, tb2])
        in_path = tmp_path / "in.pptx"
        prs.save(str(in_path))
        out_path = tmp_path / "out.pptx"

        fake_translate_texts, captured = _capturing_translate_texts()
        with patch("app.backend.processors.pptx_processor.translate_texts", side_effect=fake_translate_texts):
            translate_pptx(
                str(in_path), str(out_path), targets=["vi"], src_lang="en",
                client=MagicMock(),
            )

        uniq = captured["uniq"]
        assert "Flat sibling text" in uniq
        assert "Grouped alpha text" in uniq, f"grouped text missing from outgoing payload: {uniq}"
        assert "Grouped beta text" in uniq, f"grouped text missing from outgoing payload: {uniq}"

    def test_nested_group_text_reaches_translate_texts_payload(self, tmp_path):
        """AC-2: text inside a group-within-a-group also reaches the payload."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        inner_tb1 = _add_textbox(slide.shapes, "Inner nested text one", left=1, top=1)
        inner_tb2 = _add_textbox(slide.shapes, "Inner nested text two", left=3, top=1)
        inner_group = slide.shapes.add_group_shape([inner_tb1, inner_tb2])
        outer_sibling = _add_textbox(slide.shapes, "Outer sibling text", left=5, top=1)
        slide.shapes.add_group_shape([inner_group, outer_sibling])
        in_path = tmp_path / "in.pptx"
        prs.save(str(in_path))
        out_path = tmp_path / "out.pptx"

        fake_translate_texts, captured = _capturing_translate_texts()
        with patch("app.backend.processors.pptx_processor.translate_texts", side_effect=fake_translate_texts):
            translate_pptx(
                str(in_path), str(out_path), targets=["vi"], src_lang="en",
                client=MagicMock(),
            )

        uniq = captured["uniq"]
        assert "Inner nested text one" in uniq, f"nested-group text missing: {uniq}"
        assert "Inner nested text two" in uniq, f"nested-group text missing: {uniq}"
        assert "Outer sibling text" in uniq, f"outer-group sibling text missing: {uniq}"


# ---------------------------------------------------------------------------
# AC-3: grouped table cells map to the correct (row, col)
# ---------------------------------------------------------------------------

class TestGroupedTableCoordinates:
    def test_grouped_table_cells_map_to_correct_row_col(self, tmp_path, monkeypatch):
        monkeypatch.setattr(config, "JSON_STRUCTURED_TRANSLATION_ENABLED", False)

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        cell_texts = {(0, 0): "r0c0", (0, 1): "r0c1", (1, 0): "r1c0", (1, 1): "r1c1"}
        tbl_shape = _add_table(slide.shapes, 2, 2, cell_texts)
        sibling = _add_textbox(slide.shapes, "sibling", left=5, top=1)
        slide.shapes.add_group_shape([tbl_shape, sibling])
        in_path = tmp_path / "in.pptx"
        prs.save(str(in_path))
        out_path = tmp_path / "out.pptx"

        fake_translate_texts, _ = _capturing_translate_texts()
        with patch("app.backend.processors.pptx_processor.translate_texts", side_effect=fake_translate_texts):
            translate_pptx(
                str(in_path), str(out_path), targets=["vi"], src_lang="en",
                client=_mock_client_for_table_translation(),
            )

        out_prs = Presentation(str(out_path))
        r0c1 = _grouped_table_cell(out_prs, 0, 1, "r0c0")
        r1c0 = _grouped_table_cell(out_prs, 1, 0, "r0c0")
        assert r0c1 and "R0C1" in r0c1[0], f"(0,1) did not receive its own translation: {r0c1}"
        assert r1c0 and "R1C0" in r1c0[0], f"(1,0) did not receive its own translation: {r1c0}"
        # Selection, not count: (0,1) must NOT carry (1,0)'s translation or vice versa.
        assert "R1C0" not in r0c1[0]
        assert "R0C1" not in r1c0[0]


# ---------------------------------------------------------------------------
# AC-4: document-order counter replaces id(shape); no cross-table collision
# ---------------------------------------------------------------------------

class TestTableIdCounterNoCollision:
    def test_many_grouped_tables_no_shared_key_under_forced_gc(self, tmp_path, monkeypatch):
        """Build many grouped tables (each with a uniquely-identifying cell),
        forcing gc.collect() between constructions to exercise proxy-address
        reuse. Every table's cell must map to its OWN translation, never
        another table's, under the new document-order counter.

        Note (verified by snapshot-sabotage during development, never via
        git checkout): reverting the grouping key to `id(shape)` did NOT
        reproduce a collision here, because every populated cell is stored
        as a segment that transitively retains its parent shape (`_Cell.
        _parent` -> `Table._graphic_frame` -> the shape), keeping the
        address alive for the whole run — the same "masked in the main
        loop" finding BR-81/BR-113 document for DOCX's analogous id()
        keys (see specs/archive/2026/docx-nested-table-collection/
        evidence/id-key-hazard.md). The document-order counter is still the
        right fix: it removes reliance on an unstated, untested retention
        invariant (BR-116), independent of whether today's exact code shape
        happens to mask the hazard."""
        monkeypatch.setattr(config, "JSON_STRUCTURED_TRANSLATION_ENABLED", False)

        # NOTE: each table has 2 cells (not 1) because the legacy pipe-grid
        # serializer/parser round-trip (table_serializer.serialize()/parse())
        # requires at least one "|" delimiter to succeed; a genuine 1-cell
        # grid has no delimiter and always falls back to a per-cell batch
        # translation (a pre-existing table_serializer quirk, out of scope
        # for this change) — which would route through translate_texts
        # instead of exercising the translate_once uppercase transform this
        # test relies on to detect cross-table collisions.
        n_tables = 40
        prs = Presentation()
        for i in range(n_tables):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            tbl_shape = _add_table(slide.shapes, 1, 2, {(0, 0): f"cell_{i:03d}", (0, 1): "filler"})
            sibling = _add_textbox(slide.shapes, f"sibling_{i:03d}", left=5, top=1)
            slide.shapes.add_group_shape([tbl_shape, sibling])
            gc.collect()
        in_path = tmp_path / "in.pptx"
        prs.save(str(in_path))
        out_path = tmp_path / "out.pptx"

        fake_translate_texts, _ = _capturing_translate_texts()
        gc.collect()
        with patch("app.backend.processors.pptx_processor.translate_texts", side_effect=fake_translate_texts):
            translate_pptx(
                str(in_path), str(out_path), targets=["vi"], src_lang="en",
                client=_mock_client_for_table_translation(),
            )

        out_prs = Presentation(str(out_path))
        mismatches = []
        for i in range(n_tables):
            marker = f"cell_{i:03d}"
            hits = _grouped_table_cell(out_prs, 0, 0, marker)
            assert hits, f"table {i} cell missing entirely from output"
            expected = f"CELL_{i:03d}"
            if expected not in hits[0]:
                mismatches.append((i, hits[0]))
        assert not mismatches, (
            f"{len(mismatches)} of {n_tables} grouped tables received a "
            f"translation that was not their own (cross-table collision): {mismatches[:5]}"
        )


# ---------------------------------------------------------------------------
# AC-5: flat (non-grouped) shape output unchanged
# ---------------------------------------------------------------------------

class TestFlatShapeRegression:
    def test_flat_textbox_and_table_output_unchanged(self, tmp_path, monkeypatch):
        monkeypatch.setattr(config, "JSON_STRUCTURED_TRANSLATION_ENABLED", False)

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_textbox(slide.shapes, "Hello world", left=1, top=1)
        # 2 cells (not 1): a genuine 1-cell grid has no "|" delimiter and
        # always falls back to a per-cell translate_texts batch (pre-existing
        # table_serializer quirk, out of scope) instead of exercising the
        # translate_once uppercase transform asserted below.
        _add_table(slide.shapes, 1, 2, {(0, 0): "flat_cell", (0, 1): "filler"}, left=1, top=3)
        in_path = tmp_path / "in.pptx"
        prs.save(str(in_path))
        out_path = tmp_path / "out.pptx"

        fake_translate_texts, captured = _capturing_translate_texts()
        with patch("app.backend.processors.pptx_processor.translate_texts", side_effect=fake_translate_texts):
            translate_pptx(
                str(in_path), str(out_path), targets=["vi"], src_lang="en",
                client=_mock_client_for_table_translation(),
            )

        assert "Hello world" in captured["uniq"]

        out_prs = Presentation(str(out_path))
        all_texts = _all_text_frame_texts(out_prs)
        assert any("Hello world" in t for t in all_texts)
        assert any("TR:Hello world" in t for t in all_texts)
        assert any("flat_cell" in t for t in all_texts)
        assert any("FLAT_CELL" in t for t in all_texts)


# ---------------------------------------------------------------------------
# AC-6: bounded group-nesting depth — never-drop + exactly one WARNING
# ---------------------------------------------------------------------------

class TestGroupNestingDepthGuard:
    def test_over_limit_group_still_collected_with_one_warning(self, tmp_path, caplog):
        """4 levels of nested groups exceeds MAX_GROUP_NESTING_DEPTH=3. The
        deepest content must still be collected (never dropped) and exactly
        one WARNING must be logged via the "TranslateTool" logger."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        deep_tb = _add_textbox(slide.shapes, "deep_text_at_depth4", left=1, top=1)
        level4 = slide.shapes.add_group_shape([deep_tb])
        sibling3 = _add_textbox(slide.shapes, "level3_sibling", left=3, top=1)
        level3 = slide.shapes.add_group_shape([level4, sibling3])
        sibling2 = _add_textbox(slide.shapes, "level2_sibling", left=5, top=1)
        level2 = slide.shapes.add_group_shape([level3, sibling2])
        sibling1 = _add_textbox(slide.shapes, "level1_sibling", left=7, top=1)
        slide.shapes.add_group_shape([level2, sibling1])
        in_path = tmp_path / "in.pptx"
        prs.save(str(in_path))
        out_path = tmp_path / "out.pptx"

        fake_translate_texts, captured = _capturing_translate_texts()
        with caplog.at_level(logging.WARNING):
            with patch("app.backend.processors.pptx_processor.translate_texts", side_effect=fake_translate_texts):
                translate_pptx(
                    str(in_path), str(out_path), targets=["vi"], src_lang="en",
                    client=MagicMock(),
                )

        uniq = captured["uniq"]
        assert "deep_text_at_depth4" in uniq, f"over-limit group content was dropped: {uniq}"

        tool_warnings = [r for r in caplog.records if r.name == "TranslateTool" and r.levelno == logging.WARNING]
        assert len(tool_warnings) == 1, (
            f"expected exactly one TranslateTool WARNING for the over-limit group, got "
            f"{len(tool_warnings)}: {[r.message for r in tool_warnings]}"
        )

    def test_depth_within_limit_no_warning(self, tmp_path, caplog):
        """3 levels of nested groups is within MAX_GROUP_NESTING_DEPTH=3 (fully
        recursed); no depth-limit WARNING should be logged."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        deep_tb = _add_textbox(slide.shapes, "within_bound_text", left=1, top=1)
        level3 = slide.shapes.add_group_shape([deep_tb])
        level2 = slide.shapes.add_group_shape([level3])
        slide.shapes.add_group_shape([level2])
        in_path = tmp_path / "in.pptx"
        prs.save(str(in_path))
        out_path = tmp_path / "out.pptx"

        fake_translate_texts, captured = _capturing_translate_texts()
        with caplog.at_level(logging.WARNING):
            with patch("app.backend.processors.pptx_processor.translate_texts", side_effect=fake_translate_texts):
                translate_pptx(
                    str(in_path), str(out_path), targets=["vi"], src_lang="en",
                    client=MagicMock(),
                )

        assert "within_bound_text" in captured["uniq"]
        tool_warnings = [r for r in caplog.records if r.name == "TranslateTool" and r.levelno == logging.WARNING]
        assert not tool_warnings, f"unexpected depth WARNING within the bound: {[r.message for r in tool_warnings]}"


# ---------------------------------------------------------------------------
# AC-7: SmartArt path untouched by group collection
# ---------------------------------------------------------------------------

class TestSmartArtUntouched:
    def test_smartart_path_not_invoked_for_group_collection(self, tmp_path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb1 = _add_textbox(slide.shapes, "Grouped text", left=1, top=1)
        tb2 = _add_textbox(slide.shapes, "Another grouped text", left=3, top=1)
        slide.shapes.add_group_shape([tb1, tb2])
        in_path = tmp_path / "in.pptx"
        prs.save(str(in_path))
        out_path = tmp_path / "out.pptx"

        fake_translate_texts, _ = _capturing_translate_texts()
        with (
            patch("app.backend.processors.pptx_processor.translate_texts", side_effect=fake_translate_texts),
            patch(
                "app.backend.processors.pptx_processor._extract_smartart_texts",
                wraps=__import__(
                    "app.backend.processors.pptx_processor", fromlist=["_extract_smartart_texts"]
                )._extract_smartart_texts,
            ) as spy_smartart,
        ):
            translate_pptx(
                str(in_path), str(out_path), targets=["vi"], src_lang="en",
                client=MagicMock(),
            )

        assert spy_smartart.call_count == 1, (
            f"_extract_smartart_texts call count changed for group collection: {spy_smartart.call_count}"
        )
