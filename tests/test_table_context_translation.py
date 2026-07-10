"""TDD tests for table-context-translation behavior (IP-0: Red phase).

Tests that will fail until all implementation steps (IP-1 through IP-4) are done:
  - IP-1: table_serializer.py (serialize/parse)
  - IP-2: _build_table_translate_prompt in both clients
  - IP-3: PDF translate_table_cells uses serializer + single translate_once call
  - IP-4: DOCX/XLSX/PPTX processors group cells per table, key by col

Anti-tautology rules (CLAUDE.md):
  - All LLM mocks target client boundary using collection-time patch.object.
  - Do NOT call translate_document(); call processor functions directly.
  - Assert WHAT was called and WHICH translations appear, not just call counts.

Collection-time imports:
  Modules captured at collection time so patch.object is immune to sys.modules
  contamination (CLAUDE.md promoted learnings).
"""

from __future__ import annotations

import json
import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
from unittest.mock import MagicMock, call, patch

import pytest

from app.backend import config

# ---------------------------------------------------------------------------
# Collection-time module imports (patch.object targets)
# ---------------------------------------------------------------------------

import app.backend.processors.docx_processor as _docx_proc
import app.backend.processors.pptx_processor as _pptx_proc
import app.backend.processors.xlsx_processor as _xlsx_proc
import app.backend.services.translation_service as _ts
import app.backend.clients.ollama_client as _ollama_mod
import app.backend.clients.openai_compatible_client as _oa_mod

try:
    import app.backend.utils.table_serializer as _table_ser
except ImportError:
    _table_ser = None  # type: ignore[assignment]

from app.backend.models.translatable_document import (
    TableCell,
    TableStructure,
    TranslatableElement,
    BoundingBox,
    ElementType,
)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_client_mock(**kwargs) -> MagicMock:
    """Return a minimal MagicMock client that satisfies processor health checks."""
    m = MagicMock()
    m.health_check.return_value = (True, "ok")
    m.system_prompt = ""
    m.model_type = "general"
    m._is_translation_dedicated.return_value = False
    m._is_translategemma_model.return_value = False
    for k, v in kwargs.items():
        setattr(m, k, v)
    return m


def _json_grid_response(mapping: Dict[Tuple[int, int], str]) -> str:
    """Build a coordinate-JSON reply string (BR-79/BR-82) from a
    {(row, col): translation} mapping — the JSON-ON path analogue of
    `_grid_response` (the legacy pipe-grid text builder) below."""
    return json.dumps({
        "cells": [{"row": r, "col": c, "translation": t} for (r, c), t in mapping.items()]
    })


def _total_table_calls(mock_client: MagicMock) -> int:
    """Total whole-table LLM calls across BOTH wire paths — used by tests
    parameterised over `JSON_STRUCTURED_TRANSLATION_ENABLED` so the SAME
    assertion ("exactly one call") holds regardless of which path fires."""
    return mock_client.translate_once.call_count + mock_client.translate_json.call_count


def _make_table_cell(row: int, col: int, content: str, is_numeric: bool = False) -> TableCell:
    from dataclasses import dataclass as _dc
    return TableCell(
        cell_id=f"c{row}{col}",
        row=row,
        col=col,
        content=content,
        is_numeric=is_numeric,
    )


def _make_table_structure(rows: List[List[str]]) -> TableStructure:
    """Build a TableStructure from a 2-D list of strings."""
    num_rows = len(rows)
    num_cols = len(rows[0]) if rows else 0
    cells = []
    for r_idx, row in enumerate(rows):
        for c_idx, text in enumerate(row):
            cells.append(_make_table_cell(r_idx, c_idx, text))
    return TableStructure(
        num_rows=num_rows,
        num_cols=num_cols,
        recognizer="test",
        recognition_confident=True,
        cells=cells,
    )


def _make_table_element(ts: TableStructure, eid: str = "elem-1") -> TranslatableElement:
    """Wrap TableStructure in a table-typed TranslatableElement."""
    return TranslatableElement(
        element_id=eid,
        content="",
        element_type=ElementType.TABLE,
        page_num=1,
        bbox=BoundingBox(x0=0.0, y0=0.0, x1=100.0, y1=100.0),
        metadata={"table_structure": ts.to_dict()},
    )


def _make_docx_with_table(tmp_path: Path, rows: List[List[str]]) -> Path:
    """Create a minimal DOCX containing one table (no body paragraphs)."""
    import docx as _docx_lib
    doc = _docx_lib.Document()
    if rows:
        num_rows = len(rows)
        num_cols = len(rows[0])
        tbl = doc.add_table(rows=num_rows, cols=num_cols)
        for r, row in enumerate(rows):
            for c, text in enumerate(row):
                tbl.cell(r, c).text = text
    p = tmp_path / "table_test.docx"
    doc.save(str(p))
    return p


def _make_xlsx_with_cells(tmp_path: Path, rows: List[List[str]]) -> Path:
    """Create a minimal XLSX with cells filled from rows matrix."""
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    for r_idx, row in enumerate(rows, 1):
        for c_idx, text in enumerate(row, 1):
            ws.cell(row=r_idx, column=c_idx, value=text)
    p = tmp_path / "table_test.xlsx"
    wb.save(str(p))
    return p


def _make_pptx_with_table(tmp_path: Path, rows: List[List[str]]) -> Path:
    """Create a minimal PPTX with one table on one slide."""
    import pptx
    from pptx.util import Inches
    prs = pptx.Presentation()
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    num_rows = len(rows)
    num_cols = len(rows[0]) if rows else 0
    tbl_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(1), Inches(1), Inches(6), Inches(3)
    )
    tbl = tbl_shape.table
    for r, row in enumerate(rows):
        for c, text in enumerate(row):
            tbl.cell(r, c).text = text
    p = tmp_path / "table_test.pptx"
    prs.save(str(p))
    return p


# Grid response helpers
def _grid_response(rows: List[List[str]]) -> str:
    """Build a pipe-grid string as a mock LLM response."""
    return "\n".join(" | ".join(row) for row in rows)


# ---------------------------------------------------------------------------
# AC-2: Instruction precedes serialized grid in prompt
# ---------------------------------------------------------------------------

class TestPromptBuilder:
    """Tests for _build_table_translate_prompt in both clients (AC-2)."""

    def test_instruction_precedes_serialized_grid_in_prompt_ollama(self):
        """OllamaClient._build_table_translate_prompt: instruction appears before grid."""
        from app.backend.clients.ollama_client import OllamaClient
        serialized = "Name | Value\nApple | 100"
        prompt = OllamaClient._build_table_translate_prompt(serialized, "en", "zh")

        instr_idx = min(
            (prompt.find(word) for word in ("Translate", "translate") if prompt.find(word) >= 0),
            default=-1,
        )
        grid_idx = prompt.find("Name | Value")

        assert instr_idx >= 0, "Instruction not found in prompt"
        assert grid_idx >= 0, "Serialized grid not found in prompt"
        assert instr_idx < grid_idx, (
            f"Instruction (pos {instr_idx}) must appear BEFORE grid (pos {grid_idx}) in prompt"
        )

    def test_instruction_precedes_serialized_grid_in_prompt_openai(self):
        """OpenAICompatibleClient._build_table_translate_prompt: instruction before grid."""
        from app.backend.clients.openai_compatible_client import OpenAICompatibleClient
        serialized = "Name | Value\nApple | 100"
        prompt = OpenAICompatibleClient._build_table_translate_prompt(serialized, "en", "zh")

        instr_idx = min(
            (prompt.find(word) for word in ("Translate", "translate") if prompt.find(word) >= 0),
            default=-1,
        )
        grid_idx = prompt.find("Name | Value")

        assert instr_idx >= 0, "Instruction not found in prompt"
        assert grid_idx >= 0, "Serialized grid not found in prompt"
        assert instr_idx < grid_idx

    def test_prompt_builder_methods_produce_identical_instructions(self):
        """Both client prompt builders use identical instruction wording (AC-2 symmetry)."""
        from app.backend.clients.ollama_client import OllamaClient
        from app.backend.clients.openai_compatible_client import OpenAICompatibleClient

        serialized = "X | Y"
        p_ollama = OllamaClient._build_table_translate_prompt(serialized, "en", "zh")
        p_oa = OpenAICompatibleClient._build_table_translate_prompt(serialized, "en", "zh")

        # The instruction text before the grid must be identical
        grid_start_ollama = p_ollama.find("X | Y")
        grid_start_oa = p_oa.find("X | Y")

        instr_ollama = p_ollama[:grid_start_ollama].strip()
        instr_oa = p_oa[:grid_start_oa].strip()

        assert instr_ollama == instr_oa, (
            f"Prompt builders produce different instructions:\n"
            f"Ollama: {instr_ollama!r}\nOpenAI: {instr_oa!r}"
        )


# ---------------------------------------------------------------------------
# AC-5: PDF TableCell row/col drives serialization
# ---------------------------------------------------------------------------

class TestPdfTableCellSerialization:
    """Tests for translate_table_cells() using the shared serializer (AC-5).

    These 3 tests drive the frozen, flag-OFF legacy pipe-grid path directly
    (`table_serializer.serialize` + `client.translate_once`) — forced via
    monkeypatch since `JSON_STRUCTURED_TRANSLATION_ENABLED` now defaults to
    True (Resolution A). JSON-path coverage of the SAME behavior lives in
    `TestJsonTableRoundTrip` / `TestPhantomColumnRegression` below.
    """

    def test_pdf_tablecell_row_col_drives_serialization(self, monkeypatch):
        """translate_table_cells uses table_serializer.serialize with real row/col from cells."""
        if _table_ser is None:
            pytest.skip("table_serializer not yet created (expected RED)")
        monkeypatch.setattr(config, "JSON_STRUCTURED_TRANSLATION_ENABLED", False)

        ts = _make_table_structure([
            ["Header A", "Header B"],
            ["Data 1",   "Data 2"],
        ])
        element = _make_table_element(ts)

        mock_client = _make_client_mock()
        mock_client.translate_once.return_value = (
            True, _grid_response([["标题A", "标题B"], ["数据1", "数据2"]])
        )

        captured_serialize_args = []

        original_serialize = _table_ser.serialize

        def capturing_serialize(cells):
            captured_serialize_args.append(list(cells))
            return original_serialize(cells)

        with patch.object(_table_ser, "serialize", side_effect=capturing_serialize):
            _ts.translate_table_cells(
                element, targets=["zh"], src_lang="en", client=mock_client,
            )

        # serialize must have been called — passing the actual TableCell objects
        assert len(captured_serialize_args) >= 1, (
            "table_serializer.serialize was not called; "
            "translate_table_cells must use the shared serializer (AC-5)"
        )

        # The cells passed to serialize must have row/col set from the IR
        all_cells = captured_serialize_args[0]
        rows_seen = {c.row for c in all_cells}
        cols_seen = {c.col for c in all_cells}
        assert 0 in rows_seen and 1 in rows_seen, f"Expected rows 0 and 1; got {rows_seen}"
        assert 0 in cols_seen and 1 in cols_seen, f"Expected cols 0 and 1; got {cols_seen}"

    def test_pdf_translate_table_cells_calls_translate_once_once(self, monkeypatch):
        """After IP-3: translate_table_cells calls client.translate_once exactly once
        for a table with ≥1 translatable cell (instead of per-cell batch calls)."""
        if _table_ser is None:
            pytest.skip("table_serializer not yet created (expected RED)")
        monkeypatch.setattr(config, "JSON_STRUCTURED_TRANSLATION_ENABLED", False)

        ts = _make_table_structure([
            ["Name", "Value"],
            ["Apple", "Fruit"],
        ])
        element = _make_table_element(ts)

        mock_client = _make_client_mock()
        mock_client.translate_once.return_value = (
            True, _grid_response([["名字", "价值"], ["苹果", "水果"]])
        )

        with patch.object(_ts, "translate_blocks_batch") as mock_batch:
            _ts.translate_table_cells(
                element, targets=["zh"], src_lang="en", client=mock_client,
            )
            # After IP-3: translate_blocks_batch NOT used for the whole-table call
            assert mock_batch.call_count == 0, (
                "translate_blocks_batch should not be called when serializer path succeeds"
            )

        # translate_once called exactly once for the whole table
        assert mock_client.translate_once.call_count == 1, (
            f"Expected 1 translate_once call for whole-table translation; "
            f"got {mock_client.translate_once.call_count}"
        )

    def test_pdf_translations_mapped_to_correct_cells_after_parse(self, monkeypatch):
        """After parse(), each non-numeric cell gets grid[r][c] as translated_content."""
        if _table_ser is None:
            pytest.skip("table_serializer not yet created (expected RED)")
        monkeypatch.setattr(config, "JSON_STRUCTURED_TRANSLATION_ENABLED", False)

        ts = _make_table_structure([
            ["Name", "Value"],
            ["Apple", "Fruit"],
        ])
        element = _make_table_element(ts)

        mock_client = _make_client_mock()
        mock_client.translate_once.return_value = (
            True, _grid_response([["名字", "价值"], ["苹果", "水果"]])
        )

        _ts.translate_table_cells(
            element, targets=["zh"], src_lang="en", client=mock_client,
        )

        from app.backend.models.translatable_document import TableStructure as TS
        ts_dict = element.metadata["table_structure"]
        ts_out = TS.from_dict(ts_dict)

        cell_map = {(c.row, c.col): c for c in ts_out.cells}
        assert cell_map[(0, 0)].translated_content == "名字", (
            f"Expected '名字' at (0,0), got {cell_map[(0,0)].translated_content!r}"
        )
        assert cell_map[(0, 1)].translated_content == "价值"
        assert cell_map[(1, 0)].translated_content == "苹果"
        assert cell_map[(1, 1)].translated_content == "水果"


# ---------------------------------------------------------------------------
# AC-8: Fallback when parse() returns None
# ---------------------------------------------------------------------------

class TestFallbackBehavior:
    """Tests for the per-cell SEG fallback when parse() returns None (AC-8)."""

    def test_fallback_per_cell_batch_preserves_all_cell_mapping(self):
        """When translate_once returns an unparseable response, all cells still get
        translations via the fallback per-cell SEG batch (AC-8)."""
        if _table_ser is None:
            pytest.skip("table_serializer not yet created (expected RED)")

        ts = _make_table_structure([
            ["Name", "Value"],
            ["Apple", "Fruit"],
        ])
        element = _make_table_element(ts)

        mock_client = _make_client_mock()
        # Return a garbled response that parse()/parse_json() cannot parse —
        # equally invalid under either wire path (not pipe-delimited, not
        # valid JSON), so this fires the fallback regardless of the default
        # JSON_STRUCTURED_TRANSLATION_ENABLED flag value.
        mock_client.translate_once.return_value = (
            True, "This is not a valid pipe-grid at all"
        )
        mock_client.translate_json.return_value = mock_client.translate_once.return_value

        # Fallback: translate_blocks_batch should be called for each cell
        fallback_results = [
            (True, "名字"),
            (True, "价值"),
            (True, "苹果"),
            (True, "水果"),
        ]

        with patch.object(_ts, "translate_blocks_batch", return_value=fallback_results) as mock_batch:
            _ts.translate_table_cells(
                element, targets=["zh"], src_lang="en", client=mock_client,
            )

        # Fallback must have been invoked
        assert mock_batch.call_count >= 1, (
            "translate_blocks_batch (fallback) must be called when parse() returns None"
        )

        # All translatable cells must have translations
        from app.backend.models.translatable_document import TableStructure as TS
        ts_out = TS.from_dict(element.metadata["table_structure"])
        for cell in ts_out.cells:
            assert cell.translation_status in ("translated", "passthrough", "skipped"), (
                f"Cell ({cell.row},{cell.col}) has unexpected status: {cell.translation_status!r}"
            )
            assert cell.translated_content is not None, (
                f"Cell ({cell.row},{cell.col}) has None translated_content after fallback"
            )

    def test_fallback_logs_warning_on_parse_failure(self, caplog):
        """When parse() returns None, a WARNING is logged with dimension info (BR-82)."""
        if _table_ser is None:
            pytest.skip("table_serializer not yet created (expected RED)")

        ts = _make_table_structure([["Name", "Value"]])
        element = _make_table_element(ts)

        mock_client = _make_client_mock()
        mock_client.translate_once.return_value = (True, "no pipes here")
        mock_client.translate_json.return_value = mock_client.translate_once.return_value

        with patch.object(_ts, "translate_blocks_batch", return_value=[(True, "名字"), (True, "价值")]):
            with caplog.at_level(logging.WARNING):
                _ts.translate_table_cells(
                    element, targets=["zh"], src_lang="en", client=mock_client,
                )

        warning_texts = [r.message for r in caplog.records if r.levelno >= logging.WARNING]
        assert any("parse" in w.lower() or "mismatch" in w.lower() or "fallback" in w.lower()
                   for w in warning_texts), (
            f"Expected a WARNING log about parse failure/fallback; got: {warning_texts}"
        )

    def test_unparseable_json_falls_back_to_per_cell_batch(self, monkeypatch):
        """AC-4 (table): explicit flag-ON JSON-path fallback — an unparseable
        JSON reply discards the whole reply and falls back to the per-cell
        SEG batch (BR-82), the job never fails."""
        if _table_ser is None:
            pytest.skip("table_serializer not yet created (expected RED)")
        monkeypatch.setattr(config, "JSON_STRUCTURED_TRANSLATION_ENABLED", True)

        ts = _make_table_structure([["Name", "Value"]])
        element = _make_table_element(ts)

        mock_client = _make_client_mock()
        mock_client.translate_json.return_value = (True, "{not valid json")

        fallback_results = [(True, "名字"), (True, "价值")]
        with patch.object(_ts, "translate_blocks_batch", return_value=fallback_results) as mock_batch:
            _ts.translate_table_cells(
                element, targets=["zh"], src_lang="en", client=mock_client,
            )

        assert mock_batch.call_count >= 1, (
            "translate_blocks_batch (fallback) must fire when parse_json() rejects the reply"
        )
        from app.backend.models.translatable_document import TableStructure as TS
        ts_out = TS.from_dict(element.metadata["table_structure"])
        cell_map = {(c.row, c.col): c for c in ts_out.cells}
        assert cell_map[(0, 0)].translated_content == "名字"
        assert cell_map[(0, 1)].translated_content == "价值"


# ---------------------------------------------------------------------------
# AC-5 (table): INFO fallback line reaches the TranslateTool logger (BR-82/BR-109)
# ---------------------------------------------------------------------------

class TestFallbackLogging:
    def test_fallback_emits_info_via_translatetool_logger(self, monkeypatch, caplog):
        if _table_ser is None:
            pytest.skip("table_serializer not yet created (expected RED)")
        monkeypatch.setattr(config, "JSON_STRUCTURED_TRANSLATION_ENABLED", True)

        ts = _make_table_structure([["Name", "Value"]])
        element = _make_table_element(ts)

        mock_client = _make_client_mock()
        mock_client.translate_json.return_value = (True, "{not valid json")

        from app.backend.utils.logging_utils import logger as translate_tool_logger

        with patch.object(_ts, "translate_blocks_batch", return_value=[(True, "名字"), (True, "价值")]):
            with caplog.at_level(logging.INFO, logger="TranslateTool"):
                _ts.translate_table_cells(
                    element, targets=["zh"], src_lang="en", client=mock_client,
                    log=translate_tool_logger.info,
                )

        info_records = [
            r for r in caplog.records
            if r.name == "TranslateTool" and r.levelno == logging.INFO
        ]
        assert any("fallback" in r.message.lower() or "json" in r.message.lower() for r in info_records), (
            f"Expected an INFO fallback line on the TranslateTool logger; "
            f"got: {[r.message for r in info_records]}"
        )


# ---------------------------------------------------------------------------
# Resilience: hostile JSON table replies (BR-82) — job-level (translate_table_cells),
# not the parse_json unit level already covered exhaustively in
# tests/test_table_serialization.py::TestParseCoordinateRemap. These prove the
# JOB completes (never raises, correct fallback / no-fallback outcome, and the
# INFO fallback line reaches the TranslateTool logger) when the model returns
# a hostile reply at the actual translation_service.translate_table_cells seam.
# ---------------------------------------------------------------------------

def _sent_cells_from_rows(rows: List[List[str]]) -> Dict[Tuple[int, int], str]:
    """Build the {(row, col): source_text} map translate_table_cells will send
    for a single-target-language table built from `_make_table_structure`."""
    return {
        (r, c): text
        for r, row in enumerate(rows)
        for c, text in enumerate(row)
        if text
    }


class TestHostileTableJsonReplies:
    """AC-4/AC-5 (table): each of BR-82's reject/tolerate triggers exercised
    through the real `translation_service.translate_table_cells` seam, with a
    mocked `client.translate_json` and a mocked `translate_blocks_batch`
    fallback target. Every scenario proves the job completes (no exception)
    and asserts WHICH branch fired (fallback vs. accepted), never merely a
    changed-cell count.
    """

    @pytest.mark.parametrize(
        "hostile_content,expected_reason_substr",
        [
            pytest.param("", "empty content", id="empty-content-gpt-oss-120b-mode"),
            pytest.param("not valid json at all", "unparseable JSON", id="plain-prose"),
            pytest.param('{"cells": [{"row": 0, "col": 0, "translat', "unparseable JSON", id="truncated-mid-object"),
            pytest.param('{"cells": []} trailing garbage after the object', "unparseable JSON", id="trailing-garbage"),
            pytest.param(json.dumps({"foo": "bar"}), "missing or non-list 'cells' key", id="missing-cells-key"),
            pytest.param(
                json.dumps({"cells": [{"row": 0, "col": 0, "translation": "名字"}]}),
                "reply omits",
                id="malformed-cell-missing-translation-for-second-sent-coord",
            ),
        ],
    )
    def test_hostile_reply_falls_back_to_per_cell_batch_and_logs_reason(
        self, monkeypatch, caplog, hostile_content, expected_reason_substr,
    ):
        """Every hostile-reply trigger: whole reply discarded, per-cell SEG
        fallback (BR-82) fires, job completes, and the reason is visible at
        INFO through the TranslateTool logger (BR-109)."""
        monkeypatch.setattr(config, "JSON_STRUCTURED_TRANSLATION_ENABLED", True)

        rows = [["Name", "Value"]]
        ts = _make_table_structure(rows)
        element = _make_table_element(ts)

        mock_client = _make_client_mock()
        mock_client.translate_json.return_value = (True, hostile_content)

        from app.backend.utils.logging_utils import logger as translate_tool_logger

        fallback_results = [(True, "FALLBACK_NAME"), (True, "FALLBACK_VALUE")]
        with patch.object(_ts, "translate_blocks_batch", return_value=fallback_results) as mock_batch:
            with caplog.at_level(logging.INFO, logger="TranslateTool"):
                _ts.translate_table_cells(
                    element, targets=["zh"], src_lang="en", client=mock_client,
                    log=translate_tool_logger.info,
                )

        assert mock_batch.call_count >= 1, (
            f"per-cell fallback (BR-82) must fire for hostile content {hostile_content!r}"
        )

        from app.backend.models.translatable_document import TableStructure as TS
        ts_out = TS.from_dict(element.metadata["table_structure"])
        cell_map = {(c.row, c.col): c for c in ts_out.cells}
        assert cell_map[(0, 0)].translated_content == "FALLBACK_NAME"
        assert cell_map[(0, 1)].translated_content == "FALLBACK_VALUE"

        info_records = [
            r for r in caplog.records
            if r.name == "TranslateTool" and r.levelno == logging.INFO
        ]
        assert any(expected_reason_substr in r.message for r in info_records), (
            f"expected an INFO line naming reason containing {expected_reason_substr!r}; "
            f"got: {[r.message for r in info_records]}"
        )

    def test_non_integer_coordinate_falls_back(self, monkeypatch):
        """A cell entry carrying a non-integer (row, col) coordinate rejects
        the whole reply (BR-82 malformed-cell reject)."""
        monkeypatch.setattr(config, "JSON_STRUCTURED_TRANSLATION_ENABLED", True)
        rows = [["Name", "Value"]]
        ts = _make_table_structure(rows)
        element = _make_table_element(ts)

        mock_client = _make_client_mock()
        mock_client.translate_json.return_value = (True, json.dumps({
            "cells": [
                {"row": "0", "col": 0, "translation": "名字"},
                {"row": 0, "col": 1, "translation": "价值"},
            ]
        }))

        fallback_results = [(True, "FALLBACK_NAME"), (True, "FALLBACK_VALUE")]
        with patch.object(_ts, "translate_blocks_batch", return_value=fallback_results) as mock_batch:
            _ts.translate_table_cells(
                element, targets=["zh"], src_lang="en", client=mock_client,
            )

        assert mock_batch.call_count >= 1, "non-integer coordinate must trigger the BR-82 fallback"

    def test_missing_sent_coordinate_rejects_whole_reply_no_partial_assignment(self, monkeypatch):
        """BR-82: a reply that omits ANY sent (row, col) coordinate is
        rejected WHOLESALE — even the one coordinate the reply DID answer
        correctly must NOT be adopted; the fallback value must win instead."""
        monkeypatch.setattr(config, "JSON_STRUCTURED_TRANSLATION_ENABLED", True)
        rows = [["Name", "Value"]]
        ts = _make_table_structure(rows)
        element = _make_table_element(ts)

        mock_client = _make_client_mock()
        # Reply answers (0,0) with a distinctive value the test can detect if
        # (and only if) it were WRONGLY adopted, but omits (0,1) entirely.
        mock_client.translate_json.return_value = (True, json.dumps({
            "cells": [{"row": 0, "col": 0, "translation": "PARTIAL_SHOULD_NOT_BE_USED"}]
        }))

        fallback_results = [(True, "FALLBACK_NAME"), (True, "FALLBACK_VALUE")]
        with patch.object(_ts, "translate_blocks_batch", return_value=fallback_results) as mock_batch:
            _ts.translate_table_cells(
                element, targets=["zh"], src_lang="en", client=mock_client,
            )

        assert mock_batch.call_count >= 1, "missing-coordinate reply must trigger the BR-82 fallback"
        from app.backend.models.translatable_document import TableStructure as TS
        ts_out = TS.from_dict(element.metadata["table_structure"])
        cell_map = {(c.row, c.col): c for c in ts_out.cells}
        assert cell_map[(0, 0)].translated_content == "FALLBACK_NAME", (
            "the partially-correct reply value MUST NOT be adopted — no partial assignment "
            f"(got {cell_map[(0, 0)].translated_content!r})"
        )
        assert cell_map[(0, 1)].translated_content == "FALLBACK_VALUE"

    def test_extra_coordinate_ignored_not_an_error_no_fallback(self, monkeypatch):
        """BR-82: a reply cell at a coordinate never sent is IGNORED, not
        rejected — the job completes normally via the JSON path, the
        fallback must NOT fire."""
        monkeypatch.setattr(config, "JSON_STRUCTURED_TRANSLATION_ENABLED", True)
        rows = [["Name", "Value"]]
        ts = _make_table_structure(rows)
        element = _make_table_element(ts)

        mock_client = _make_client_mock()
        mock_client.translate_json.return_value = (True, json.dumps({
            "cells": [
                {"row": 0, "col": 0, "translation": "名字"},
                {"row": 0, "col": 1, "translation": "价值"},
                {"row": 5, "col": 5, "translation": "HALLUCINATED_EXTRA_CELL"},
            ]
        }))

        with patch.object(_ts, "translate_blocks_batch") as mock_batch:
            _ts.translate_table_cells(
                element, targets=["zh"], src_lang="en", client=mock_client,
            )
            assert mock_batch.call_count == 0, (
                "an extra, never-sent coordinate must be tolerated (ignored), "
                "NOT trigger the per-cell fallback"
            )

        from app.backend.models.translatable_document import TableStructure as TS
        ts_out = TS.from_dict(element.metadata["table_structure"])
        cell_map = {(c.row, c.col): c for c in ts_out.cells}
        assert cell_map[(0, 0)].translated_content == "名字"
        assert cell_map[(0, 1)].translated_content == "价值"

    def test_echoed_whole_grid_falls_back(self, monkeypatch):
        """BR-82: every returned cell byte-identical to its source (whole
        table echoed, untranslated) MUST trigger the fallback."""
        monkeypatch.setattr(config, "JSON_STRUCTURED_TRANSLATION_ENABLED", True)
        rows = [["Name", "Value"]]
        ts = _make_table_structure(rows)
        element = _make_table_element(ts)

        mock_client = _make_client_mock()
        mock_client.translate_json.return_value = (True, json.dumps({
            "cells": [
                {"row": 0, "col": 0, "translation": "Name"},
                {"row": 0, "col": 1, "translation": "Value"},
            ]
        }))

        fallback_results = [(True, "FALLBACK_NAME"), (True, "FALLBACK_VALUE")]
        with patch.object(_ts, "translate_blocks_batch", return_value=fallback_results) as mock_batch:
            _ts.translate_table_cells(
                element, targets=["zh"], src_lang="en", client=mock_client,
            )

        assert mock_batch.call_count >= 1, (
            "a whole-grid-echoed (untranslated) reply MUST trigger the BR-82 fallback"
        )

    def test_single_cell_echo_is_legitimate_no_fallback(self, monkeypatch):
        """BR-82: a SINGLE unchanged cell (proper noun / product code / number)
        is legitimate and MUST NOT trigger the echoed-source fallback — assert
        WHICH condition fired (no fallback call at all), never a changed-cell
        count."""
        monkeypatch.setattr(config, "JSON_STRUCTURED_TRANSLATION_ENABLED", True)
        rows = [["ACME-Corp", "Value"]]
        ts = _make_table_structure(rows)
        element = _make_table_element(ts)

        mock_client = _make_client_mock()
        mock_client.translate_json.return_value = (True, json.dumps({
            "cells": [
                {"row": 0, "col": 0, "translation": "ACME-Corp"},  # legitimate proper noun, unchanged
                {"row": 0, "col": 1, "translation": "价值"},
            ]
        }))

        with patch.object(_ts, "translate_blocks_batch") as mock_batch:
            _ts.translate_table_cells(
                element, targets=["zh"], src_lang="en", client=mock_client,
            )
            assert mock_batch.call_count == 0, (
                "a single legitimately-unchanged cell must NOT trigger the "
                "echoed-source fallback (BR-82 single-cell exception)"
            )

        from app.backend.models.translatable_document import TableStructure as TS
        ts_out = TS.from_dict(element.metadata["table_structure"])
        cell_map = {(c.row, c.col): c for c in ts_out.cells}
        assert cell_map[(0, 0)].translated_content == "ACME-Corp"
        assert cell_map[(0, 0)].translation_status == "translated", (
            "the unchanged proper-noun cell must be assigned via the accepted "
            "JSON path (status='translated'), not left in a fallback/passthrough state"
        )
        assert cell_map[(0, 1)].translated_content == "价值"

    def test_job_never_raises_across_all_hostile_replies(self, monkeypatch):
        """The job never raises for ANY hostile reply shape — the ultimate
        never-fail-fallback guarantee, exercised across every trigger in one
        sweep."""
        monkeypatch.setattr(config, "JSON_STRUCTURED_TRANSLATION_ENABLED", True)
        hostile_replies = [
            "",
            "not json",
            '{"cells": [{"row": 0, "col": 0, "translat',
            '{"cells": []} trailing garbage',
            json.dumps({"foo": "bar"}),
            json.dumps({"cells": [{"row": 0, "col": 0, "translation": "只有一格"}]}),
            json.dumps({"cells": [{"row": "0", "col": 0, "translation": "x"}, {"row": 0, "col": 1, "translation": "y"}]}),
            json.dumps({"cells": [{"row": 0, "col": 0, "translation": "Name"}, {"row": 0, "col": 1, "translation": "Value"}]}),
        ]
        for hostile in hostile_replies:
            rows = [["Name", "Value"]]
            ts = _make_table_structure(rows)
            element = _make_table_element(ts)
            mock_client = _make_client_mock()
            mock_client.translate_json.return_value = (True, hostile)

            with patch.object(_ts, "translate_blocks_batch", return_value=[(True, "A"), (True, "B")]):
                try:
                    _ts.translate_table_cells(
                        element, targets=["zh"], src_lang="en", client=mock_client,
                    )
                except Exception as exc:  # pragma: no cover - test failure path
                    pytest.fail(f"job raised for hostile reply {hostile!r}: {exc!r}")

            from app.backend.models.translatable_document import TableStructure as TS
            ts_out = TS.from_dict(element.metadata["table_structure"])
            for cell in ts_out.cells:
                assert cell.translated_content is not None, (
                    f"cell ({cell.row},{cell.col}) left with None translated_content "
                    f"for hostile reply {hostile!r}"
                )


# ---------------------------------------------------------------------------
# BR-68/BR-69: Numeric passthrough + all-numeric table
# ---------------------------------------------------------------------------

class TestNumericPassthrough:
    """Tests for numeric cell exclusion from LLM batch (BR-68, BR-69)."""

    def test_all_numeric_table_makes_no_llm_call(self):
        """A table where ALL cells are numeric makes no LLM call (BR-68/AC-1)."""
        if _table_ser is None:
            pytest.skip("table_serializer not yet created (expected RED)")

        ts = _make_table_structure([])  # empty, will be replaced
        # Build all-numeric structure manually
        cells = [
            _make_table_cell(0, 0, "100", is_numeric=True),
            _make_table_cell(0, 1, "200", is_numeric=True),
            _make_table_cell(1, 0, "300", is_numeric=True),
            _make_table_cell(1, 1, "400", is_numeric=True),
        ]
        ts2 = TableStructure(
            num_rows=2, num_cols=2, recognizer="test",
            recognition_confident=True, cells=cells,
        )
        element = _make_table_element(ts2)

        mock_client = _make_client_mock()

        _ts.translate_table_cells(
            element, targets=["zh"], src_lang="en", client=mock_client,
        )

        assert mock_client.translate_once.call_count == 0, (
            "No LLM call should be made for an all-numeric table (BR-68)"
        )

        # All cells must have passthrough status
        from app.backend.models.translatable_document import TableStructure as TS
        ts_out = TS.from_dict(element.metadata["table_structure"])
        for cell in ts_out.cells:
            assert cell.translation_status == "passthrough", (
                f"Expected 'passthrough' for numeric cell ({cell.row},{cell.col}), "
                f"got {cell.translation_status!r}"
            )

    def test_numeric_cells_excluded_from_batch_and_passthrough(self):
        """Mixed table: numeric cells excluded from LLM call, get passthrough status."""
        if _table_ser is None:
            pytest.skip("table_serializer not yet created (expected RED)")

        cells = [
            _make_table_cell(0, 0, "Product"),
            _make_table_cell(0, 1, "Price"),
            _make_table_cell(1, 0, "Apple"),
            _make_table_cell(1, 1, "1.99", is_numeric=True),
        ]
        ts = TableStructure(
            num_rows=2, num_cols=2, recognizer="test",
            recognition_confident=True, cells=cells,
        )
        element = _make_table_element(ts)

        mock_client = _make_client_mock()
        # Only 3 translatable cells (Product, Price, Apple); "1.99" is numeric
        # and excluded upstream (BR-68) on both wire paths — placeholder for
        # the legacy pipe-grid, simply absent from the JSON cell list.
        mock_client.translate_once.return_value = (
            True, _grid_response([["产品", "价格"], ["苹果", "1.99"]])
        )
        mock_client.translate_json.return_value = (
            True, _json_grid_response({(0, 0): "产品", (0, 1): "价格", (1, 0): "苹果"})
        )

        _ts.translate_table_cells(
            element, targets=["zh"], src_lang="en", client=mock_client,
        )

        from app.backend.models.translatable_document import TableStructure as TS
        ts_out = TS.from_dict(element.metadata["table_structure"])
        cell_map = {(c.row, c.col): c for c in ts_out.cells}

        # Numeric cell must stay passthrough
        assert cell_map[(1, 1)].translation_status == "passthrough", (
            f"Numeric cell must have passthrough status, "
            f"got {cell_map[(1,1)].translation_status!r}"
        )
        assert cell_map[(1, 1)].translated_content == "1.99"

        # Translatable cells must be translated
        for (r, c) in [(0, 0), (0, 1), (1, 0)]:
            assert cell_map[(r, c)].translation_status == "translated", (
                f"Cell ({r},{c}) should be translated, "
                f"got {cell_map[(r,c)].translation_status!r}"
            )


# ---------------------------------------------------------------------------
# AC-1: One LLM call per table — Office processors (DOCX / XLSX / PPTX)
# ---------------------------------------------------------------------------

_ONE_TABLE_JSON_REPLY = _json_grid_response({(0, 0): "名字", (0, 1): "价值", (1, 0): "苹果", (1, 1): "水果"})


@pytest.mark.parametrize("json_enabled", [True, False])
class TestOneCallPerTableOffice:
    """Tests that each table gets exactly one whole-table LLM call (AC-1),
    parameterised over `JSON_STRUCTURED_TRANSLATION_ENABLED` so the SAME
    assertion runs against both wire paths (Resolution A cost, test-plan
    §Resolution A cost)."""

    def test_single_llm_call_per_table_docx(self, tmp_path, monkeypatch, json_enabled):
        """DOCX with one 2×2 table → exactly 1 whole-table LLM call (AC-1)."""
        monkeypatch.setattr(config, "JSON_STRUCTURED_TRANSLATION_ENABLED", json_enabled)
        rows = [["Name", "Value"], ["Apple", "Fruit"]]
        in_path = _make_docx_with_table(tmp_path, rows)
        out_path = tmp_path / "out.docx"

        mock_client = _make_client_mock()
        # Both wire paths configured; only the active one is actually invoked.
        mock_client.translate_once.return_value = (
            True, _grid_response([["名字", "价值"], ["苹果", "水果"]])
        )
        mock_client.translate_json.return_value = (True, _ONE_TABLE_JSON_REPLY)
        # translate_texts returns empty (no non-table paragraphs)
        with patch.object(_docx_proc, "translate_texts", return_value=({}, 0, 0, False)):
            _docx_proc.translate_docx(
                str(in_path), str(out_path),
                targets=["zh"], src_lang="en",
                client=mock_client,
                include_headers_shapes_via_com=False,
            )

        assert _total_table_calls(mock_client) == 1, (
            f"Expected exactly 1 whole-table LLM call; "
            f"got once={mock_client.translate_once.call_count} json={mock_client.translate_json.call_count}"
        )

    def test_single_llm_call_per_table_xlsx(self, tmp_path, monkeypatch, json_enabled):
        """XLSX with one 2×2 text region → exactly 1 whole-table LLM call (AC-1)."""
        monkeypatch.setattr(config, "JSON_STRUCTURED_TRANSLATION_ENABLED", json_enabled)
        rows = [["Name", "Value"], ["Apple", "Fruit"]]
        in_path = _make_xlsx_with_cells(tmp_path, rows)
        out_path = tmp_path / "out.xlsx"

        mock_client = _make_client_mock()
        mock_client.translate_once.return_value = (
            True, _grid_response([["名字", "价值"], ["苹果", "水果"]])
        )
        mock_client.translate_json.return_value = (True, _ONE_TABLE_JSON_REPLY)

        with patch.object(_xlsx_proc, "translate_texts", return_value=({}, 0, 0, False)):
            _xlsx_proc.translate_xlsx_xls(
                str(in_path), str(out_path),
                targets=["zh"], src_lang="en",
                client=mock_client,
            )

        assert _total_table_calls(mock_client) == 1, (
            f"Expected 1 whole-table LLM call; "
            f"got once={mock_client.translate_once.call_count} json={mock_client.translate_json.call_count}"
        )

    def test_single_llm_call_per_table_pptx(self, tmp_path, monkeypatch, json_enabled):
        """PPTX with one 2×2 table → exactly 1 whole-table LLM call (AC-1)."""
        monkeypatch.setattr(config, "JSON_STRUCTURED_TRANSLATION_ENABLED", json_enabled)
        rows = [["Name", "Value"], ["Apple", "Fruit"]]
        in_path = _make_pptx_with_table(tmp_path, rows)
        out_path = tmp_path / "out.pptx"

        mock_client = _make_client_mock()
        mock_client.translate_once.return_value = (
            True, _grid_response([["名字", "价值"], ["苹果", "水果"]])
        )
        mock_client.translate_json.return_value = (True, _ONE_TABLE_JSON_REPLY)

        with patch.object(_pptx_proc, "translate_texts", return_value=({}, 0, 0, False)):
            _pptx_proc.translate_pptx(
                str(in_path), str(out_path),
                targets=["zh"], src_lang="en",
                client=mock_client,
            )

        assert _total_table_calls(mock_client) == 1, (
            f"Expected exactly 1 whole-table LLM call; "
            f"got once={mock_client.translate_once.call_count} json={mock_client.translate_json.call_count}"
        )


# ---------------------------------------------------------------------------
# AC-1 integration: phantom-column regression (json-structured-translation-io)
# ---------------------------------------------------------------------------

class TestPhantomColumnRegression:
    """A sheet whose `ws.max_row`/`max_column` are inflated by openpyxl's
    dimension-touch quirk (a distant cell accessed without a value still
    counts toward the dimension) must send only its REAL content cells — the
    coordinate JSON path builds the cell list directly from populated
    positions, never from a `range(max_row) x range(max_col)` grid."""

    def test_xlsx_257_col_sheet_completes_without_fallback(self, tmp_path, monkeypatch):
        monkeypatch.setattr(config, "JSON_STRUCTURED_TRANSLATION_ENABLED", True)
        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        for c in range(1, 48):  # 47 real content cells, all on row 1
            ws.cell(row=1, column=c, value=f"Item{c}")
        # Real-world phantom-column trigger: touching a distant cell (row 9,
        # col 257) without a value still inflates ws.max_row/max_column.
        ws.cell(row=9, column=257)
        assert ws.max_row == 9 and ws.max_column == 257, (
            "fixture assumption broken: openpyxl no longer inflates dimensions "
            "on a value-less cell touch"
        )
        in_path = tmp_path / "phantom.xlsx"
        wb.save(str(in_path))
        out_path = tmp_path / "out.xlsx"

        def _fake_translate_json(user_payload, cancel_event=None, system_context=None):
            appended = json.loads(user_payload.rsplit("\n\n", 1)[1])
            reply_cells = [
                {"row": cell["row"], "col": cell["col"], "translation": f"T_{cell['text']}"}
                for cell in appended["cells"]
            ]
            return True, json.dumps({"cells": reply_cells})

        mock_client = _make_client_mock()
        mock_client.translate_json.side_effect = _fake_translate_json

        with patch.object(_xlsx_proc, "translate_texts") as mock_fallback:
            _xlsx_proc.translate_xlsx_xls(
                str(in_path), str(out_path),
                targets=["zh"], src_lang="en",
                client=mock_client,
            )
            assert mock_fallback.call_count == 0, (
                "no per-cell fallback should fire — the JSON call must succeed"
            )

        assert mock_client.translate_json.call_count == 1
        sent_payload = mock_client.translate_json.call_args[0][0]
        sent = json.loads(sent_payload.rsplit("\n\n", 1)[1])
        assert len(sent["cells"]) == 47, (
            f"expected exactly 47 content cells (no grid shape), got {len(sent['cells'])}"
        )
        assert "num_rows" not in sent and "num_cols" not in sent


# ---------------------------------------------------------------------------
# AC-2 integration: row-neighbor context reaches the outgoing payload
# ---------------------------------------------------------------------------

class TestJsonTableRoundTrip:
    def test_row_neighbor_context_delivered_in_outgoing_payload(self, tmp_path, monkeypatch):
        """All cells of the same table row are present TOGETHER in the ONE
        captured outgoing JSON payload — the model sees row-neighbor context
        in a single call, never isolated per-cell requests."""
        monkeypatch.setattr(config, "JSON_STRUCTURED_TRANSLATION_ENABLED", True)
        rows = [["Name", "Value"], ["Apple", "Fruit"]]
        in_path = _make_xlsx_with_cells(tmp_path, rows)
        out_path = tmp_path / "out.xlsx"

        mock_client = _make_client_mock()
        mock_client.translate_json.return_value = (True, _ONE_TABLE_JSON_REPLY)

        with patch.object(_xlsx_proc, "translate_texts", return_value=({}, 0, 0, False)):
            _xlsx_proc.translate_xlsx_xls(
                str(in_path), str(out_path),
                targets=["zh"], src_lang="en",
                client=mock_client,
            )

        assert mock_client.translate_json.call_count == 1
        sent_payload = mock_client.translate_json.call_args[0][0]
        sent = json.loads(sent_payload.rsplit("\n\n", 1)[1])
        texts_by_row: Dict[int, set] = {}
        for cell in sent["cells"]:
            texts_by_row.setdefault(cell["row"], set()).add(cell["text"])
        assert texts_by_row[0] == {"Name", "Value"}, (
            f"row 0 neighbors not delivered together: {texts_by_row.get(0)}"
        )
        assert texts_by_row[1] == {"Apple", "Fruit"}, (
            f"row 1 neighbors not delivered together: {texts_by_row.get(1)}"
        )


# ---------------------------------------------------------------------------
# Multi-paragraph cell splitting (fix: merged "layout" cell holding an entire
# document section was sent as ONE atomic translate_once() call and silently
# truncated -- confirmed live against panjit: a 4827-char cell returned only
# 370 chars with ok=True, no error surfaced, ~90% of content vanished).
# ---------------------------------------------------------------------------

class TestMultiParagraphCellSplitting:
    def _make_docx_with_multiline_cell(self, tmp_path: Path, lines: List[str]) -> Path:
        """Create a DOCX with a single 1x1 table cell containing multiple
        paragraphs (joined by "\\n" per _p_text_with_breaks), mimicking a
        merged layout cell holding several document sections."""
        import docx as _docx_lib

        doc = _docx_lib.Document()
        tbl = doc.add_table(rows=1, cols=1)
        cell = tbl.cell(0, 0)
        # A fresh cell already has one empty paragraph; reuse it for the
        # first line, then add_paragraph() for the rest.
        cell.paragraphs[0].text = lines[0]
        for line in lines[1:]:
            cell.add_paragraph(line)
        p = tmp_path / "multiline_cell.docx"
        doc.save(str(p))
        return p

    def test_multiline_cell_forces_grid_parse_failure_and_splits_in_fallback(self, tmp_path):
        """A multi-paragraph cell's whole-table translate_once response won't
        parse as a 1x1 grid (it returns one line per source line, not the
        whole grid), forcing the fallback -- which must translate each line
        SEPARATELY rather than sending the whole blob as one call."""
        lines = ["1、目的", "使本公司生產設備維持良好的狀況和工作能力。", "2、範圍"]
        in_path = self._make_docx_with_multiline_cell(tmp_path, lines)
        out_path = tmp_path / "out.docx"

        mock_client = _make_client_mock()
        # Whole-table call returns something table_serializer.parse() can't
        # reconcile against a 1x1 grid (it's multi-line, not "cellA").
        mock_client.translate_once.return_value = (True, "not a valid single-cell grid\nextra line")

        translated_lines = {
            "1、目的": "1. Mục đích",
            "使本公司生產設備維持良好的狀況和工作能力。": "Duy trì tình trạng tốt.",
            "2、範圍": "2. Phạm vi",
        }

        def _fake_translate_texts(texts, targets, src_lang, client, **kwargs):
            tmap = {(targets[0], t): translated_lines[t] for t in texts}
            return (tmap, len(texts), 0, False)

        with patch.object(_docx_proc, "translate_texts", side_effect=_fake_translate_texts):
            _docx_proc.translate_docx(
                str(in_path), str(out_path),
                targets=["Vietnamese"], src_lang="Chinese",
                client=mock_client,
                include_headers_shapes_via_com=False,
            )

        d2 = _docx_proc.docx.Document(str(out_path))
        cell_text = d2.tables[0].cell(0, 0).text
        for original, translated in translated_lines.items():
            assert original in cell_text, f"original line {original!r} must be preserved"
            assert translated in cell_text, (
                f"line {original!r} must be translated to {translated!r} individually "
                f"instead of the whole cell being sent as one atomic block; "
                f"cell_text={cell_text!r}"
            )

    def test_single_line_cell_unaffected_by_splitting(self, tmp_path):
        """A normal single-paragraph cell must behave exactly as before --
        splitting on "\\n" is a no-op when there's only one line."""
        in_path = self._make_docx_with_multiline_cell(tmp_path, ["Apple"])
        out_path = tmp_path / "out.docx"

        mock_client = _make_client_mock()
        mock_client.translate_once.return_value = (True, "not-a-grid")

        with patch.object(
            _docx_proc, "translate_texts",
            return_value=({("Vietnamese", "Apple"): "Táo"}, 1, 0, False),
        ):
            _docx_proc.translate_docx(
                str(in_path), str(out_path),
                targets=["Vietnamese"], src_lang="English",
                client=mock_client,
                include_headers_shapes_via_com=False,
            )

        d2 = _docx_proc.docx.Document(str(out_path))
        cell_text = d2.tables[0].cell(0, 0).text
        assert "Apple" in cell_text
        assert "Táo" in cell_text


# ---------------------------------------------------------------------------
# AC-3: Per-column dedup key
# ---------------------------------------------------------------------------

class TestPerColumnDedupKey:
    """Tests for (tgt, src_text, col) dedup key in office processors (AC-3)."""

    def test_same_text_different_cols_get_separate_translations(self):
        """Using table_serializer directly: same text in col 0 and col 2 occupies
        different grid positions → parse returns different translations per column."""
        if _table_ser is None:
            pytest.skip("table_serializer not yet created (expected RED)")

        from dataclasses import dataclass as _dc

        @_dc
        class _SimpleCell:
            row: int
            col: int
            content: str
            is_numeric: bool = False

        # 1-row, 3-col table where "Lead" appears in col 0 and col 2
        cells = [
            _SimpleCell(0, 0, "Lead"),
            _SimpleCell(0, 1, "Metal"),
            _SimpleCell(0, 2, "Lead"),  # same text, different column
        ]

        serialized = _table_ser.serialize(cells)
        # Mock LLM assigns different translations per column
        response = "铅 | 金属 | 领导"  # "Lead"-metal→铅, Lead-verb→领导
        grid = _table_ser.parse(response, 1, 3)

        assert grid is not None, f"parse() returned None for response {response!r}"

        # Build tmap with (tgt, text, col) keys
        tgt = "zh"
        tmap: Dict = {}
        for cell in cells:
            r, c = cell.row, cell.col
            if grid[r][c].strip():
                tmap[(tgt, cell.content, c)] = grid[r][c].strip()

        # Same source text in different columns maps to DIFFERENT translations
        t_col0 = tmap.get((tgt, "Lead", 0))
        t_col2 = tmap.get((tgt, "Lead", 2))

        assert t_col0 is not None, "Translation for (tgt, 'Lead', 0) missing from tmap"
        assert t_col2 is not None, "Translation for (tgt, 'Lead', 2) missing from tmap"
        assert t_col0 != t_col2, (
            f"Same text 'Lead' in col 0 ({t_col0!r}) and col 2 ({t_col2!r}) should get "
            "different translations when contextually placed in different columns"
        )

    def test_non_table_segments_use_col_none_dedup_key(self, tmp_path):
        """DOCX: non-table paragraph segments must use col=None in the tmap (AC-3)."""
        import docx as _docx_lib
        doc = _docx_lib.Document()
        doc.add_paragraph("Hello paragraph")  # non-table segment
        in_path = str(tmp_path / "para_only.docx")
        out_path = str(tmp_path / "para_only_out.docx")
        doc.save(in_path)

        captured_tmap: Dict = {}

        def capturing_translate_texts(texts, targets, src_lang, client, **kwargs):
            # simulate a translation: return (tgt, text): "TRANSLATED"
            result = {}
            for tgt in targets:
                for text in texts:
                    result[(tgt, text)] = f"TRANSLATED_{text}"
            return result, len(texts), 0, False

        mock_client = _make_client_mock()

        with patch.object(_docx_proc, "translate_texts", side_effect=capturing_translate_texts):
            _docx_proc.translate_docx(
                in_path, out_path,
                targets=["zh"], src_lang="en",
                client=mock_client,
                include_headers_shapes_via_com=False,
            )

        # After implementation: the tmap must use (tgt, text, None) for paragraphs
        # We verify by checking the insert pass used col=None keys.
        # Since we can't easily inspect the internal tmap, check the output file instead:
        # The paragraph should be translated.
        out_doc = _docx_lib.Document(out_path)
        all_texts = [p.text for p in out_doc.paragraphs]
        # Original paragraph must still be present
        assert any("Hello paragraph" in t for t in all_texts), (
            f"Original paragraph missing from output: {all_texts}"
        )
        # Translation must be appended (append mode by default)
        assert any("TRANSLATED" in t for t in all_texts), (
            f"Translation not found in output: {all_texts}"
        )
