"""PPTX parser for extracting translatable content.

This module parses PPTX files and extracts text content with structure
information into a TranslatableDocument format.
"""

from __future__ import annotations

import logging
import uuid
from pathlib import Path
from typing import Any, List, Optional, Set

import pptx
from pptx.shapes.base import BaseShape
from pptx.slide import Slide

from app.backend.config import MAX_SEGMENTS, MAX_TEXT_LENGTH
from app.backend.models.translatable_document import (
    BoundingBox,
    DocumentMetadata,
    ElementType,
    PageInfo,
    TranslatableDocument,
    TranslatableElement,
)
from app.backend.parsers.base import BaseParser
from app.backend.utils.exceptions import check_document_size_limits
from app.backend.utils.text_utils import should_translate

logger = logging.getLogger(__name__)


class PptxParser(BaseParser):
    """Parser for PPTX presentations.

    Extracts text from slides, shapes, and text frames into
    a unified TranslatableDocument format.
    """

    def __init__(
        self,
        max_segments: int = MAX_SEGMENTS,
        max_text_length: int = MAX_TEXT_LENGTH,
    ):
        """Initialize the parser.

        Args:
            max_segments: Maximum number of segments to extract.
            max_text_length: Maximum total text length.
        """
        self.max_segments = max_segments
        self.max_text_length = max_text_length

    @property
    def supported_extensions(self) -> list[str]:
        """Supported file extensions."""
        return [".pptx"]

    def parse(self, file_path: str) -> TranslatableDocument:
        """Parse a PPTX file.

        Args:
            file_path: Path to the PPTX file.

        Returns:
            TranslatableDocument with extracted elements.

        Raises:
            FileNotFoundError: If file does not exist.
            ValueError: If file is not a PPTX.
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        if path.suffix.lower() != ".pptx":
            raise ValueError(f"Not a PPTX file: {file_path}")

        prs = pptx.Presentation(file_path)

        elements: List[TranslatableElement] = []
        seen_keys: Set[str] = set()
        total_text_length = 0
        pages: List[PageInfo] = []

        # Get slide dimensions
        slide_width = prs.slide_width.pt if prs.slide_width else 960
        slide_height = prs.slide_height.pt if prs.slide_height else 540

        for slide_num, slide in enumerate(prs.slides, 1):
            # Add page info
            pages.append(
                PageInfo(
                    page_num=slide_num,
                    width=slide_width,
                    height=slide_height,
                )
            )

            # Extract from slide shapes
            slide_elements = self._extract_from_slide(
                slide, slide_num, seen_keys, slide_width, slide_height
            )
            elements.extend(slide_elements)
            total_text_length += sum(len(e.content) for e in slide_elements)

        # Validate document size
        check_document_size_limits(
            segment_count=len(elements),
            total_text_length=total_text_length,
            max_segments=self.max_segments,
            max_text_length=self.max_text_length,
            document_type="PowerPoint presentation",
        )

        # Build metadata
        core_props = prs.core_properties
        metadata = DocumentMetadata(
            title=core_props.title,
            author=core_props.author,
            subject=core_props.subject,
            page_count=len(prs.slides),
            has_text_layer=True,
        )

        return TranslatableDocument(
            source_path=file_path,
            source_type="pptx",
            elements=elements,
            pages=pages,
            metadata=metadata,
        )

    def _extract_from_slide(
        self,
        slide: Slide,
        slide_num: int,
        seen_keys: Set[str],
        slide_width: float,
        slide_height: float,
    ) -> List[TranslatableElement]:
        """Extract elements from a slide.

        Args:
            slide: Slide object.
            slide_num: Slide number (1-indexed).
            seen_keys: Set of already-seen text keys.
            slide_width: Slide width in points.
            slide_height: Slide height in points.

        Returns:
            List of extracted elements.
        """
        elements: List[TranslatableElement] = []

        for shape in slide.shapes:
            shape_elements = self._extract_from_shape(
                shape, slide_num, seen_keys, slide_width, slide_height
            )
            elements.extend(shape_elements)

        return elements

    def _extract_from_shape(
        self,
        shape: BaseShape,
        slide_num: int,
        seen_keys: Set[str],
        slide_width: float,
        slide_height: float,
    ) -> List[TranslatableElement]:
        """Extract elements from a shape.

        Args:
            shape: Shape object.
            slide_num: Slide number.
            seen_keys: Set of already-seen text keys.
            slide_width: Slide width in points.
            slide_height: Slide height in points.

        Returns:
            List of extracted elements.
        """
        elements: List[TranslatableElement] = []

        # Handle grouped shapes
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            if hasattr(shape, "shapes"):
                for child_shape in shape.shapes:
                    child_elements = self._extract_from_shape(
                        child_shape, slide_num, seen_keys, slide_width, slide_height
                    )
                    elements.extend(child_elements)
            return elements

        # Handle tables
        if hasattr(shape, "table"):
            table_elements = self._extract_from_table(
                shape, slide_num, seen_keys
            )
            elements.extend(table_elements)
            return elements

        # Handle text frames
        if not getattr(shape, "has_text_frame", False):
            return elements

        tf = shape.text_frame
        text = self._get_text_frame_text(tf)

        if not text.strip():
            return elements

        # Check for should_translate
        if not should_translate(text, "auto"):
            return elements

        # Generate key for deduplication
        key = f"pptx_{slide_num}_{hash(text)}_{len(text)}"
        if key in seen_keys:
            return elements
        seen_keys.add(key)

        # Get bounding box from shape position
        bbox = self._get_shape_bbox(shape)

        # Determine element type
        element_type = self._classify_shape_type(shape)

        elements.append(
            TranslatableElement(
                element_id=f"pptx_{uuid.uuid4().hex[:8]}",
                content=text,
                element_type=element_type,
                page_num=slide_num,
                bbox=bbox,
                should_translate=True,
                metadata={
                    "shape_type": shape.shape_type,
                    "shape_name": shape.name if hasattr(shape, "name") else None,
                    "text_frame_ref": tf,  # Keep reference for rendering
                },
            )
        )

        return elements

    def _extract_from_table(
        self,
        shape: BaseShape,
        slide_num: int,
        seen_keys: Set[str],
    ) -> List[TranslatableElement]:
        """Extract elements from a table shape.

        Args:
            shape: Table shape object.
            slide_num: Slide number.
            seen_keys: Set of already-seen text keys.

        Returns:
            List of extracted elements.
        """
        elements: List[TranslatableElement] = []
        table = shape.table

        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                text = cell.text.strip()
                if not text:
                    continue

                if not should_translate(text, "auto"):
                    continue

                key = f"pptx_tbl_{slide_num}_{row_idx}_{col_idx}_{hash(text)}"
                if key in seen_keys:
                    continue
                seen_keys.add(key)

                elements.append(
                    TranslatableElement(
                        element_id=f"pptx_cell_{uuid.uuid4().hex[:8]}",
                        content=text,
                        element_type=ElementType.TABLE_CELL,
                        page_num=slide_num,
                        should_translate=True,
                        metadata={
                            "in_table": True,
                            "row": row_idx + 1,
                            "col": col_idx + 1,
                            "cell_ref": cell,  # Keep reference for rendering
                        },
                    )
                )

        return elements

    def _get_text_frame_text(self, tf: Any) -> str:
        """Get text from a text frame.

        Args:
            tf: Text frame object.

        Returns:
            Text content with paragraph breaks.
        """
        return "\n".join(p.text for p in tf.paragraphs)

    def _get_shape_bbox(self, shape: BaseShape) -> Optional[BoundingBox]:
        """Get bounding box from shape position.

        Args:
            shape: Shape object.

        Returns:
            BoundingBox or None if position unavailable.
        """
        try:
            # Convert EMU to points (1 inch = 914400 EMU, 1 inch = 72 points)
            emu_to_pt = 72.0 / 914400.0

            left = shape.left * emu_to_pt if shape.left else 0
            top = shape.top * emu_to_pt if shape.top else 0
            width = shape.width * emu_to_pt if shape.width else 0
            height = shape.height * emu_to_pt if shape.height else 0

            return BoundingBox(
                x0=left,
                y0=top,
                x1=left + width,
                y1=top + height,
            )
        except (AttributeError, TypeError):
            return None

    def _classify_shape_type(self, shape: BaseShape) -> ElementType:
        """Classify shape type based on placeholder type.

        Args:
            shape: Shape to classify.

        Returns:
            ElementType for the shape.
        """
        # Check if it's a placeholder
        try:
            if hasattr(shape, "placeholder_format"):
                ph_format = shape.placeholder_format
                if ph_format:
                    ph_type = ph_format.type
                    # Title placeholders
                    if ph_type in (1, 3):  # TITLE, CENTER_TITLE
                        return ElementType.TITLE
                    # Subtitle
                    if ph_type == 4:  # SUBTITLE
                        return ElementType.TEXT
                    # Body text
                    if ph_type == 2:  # BODY
                        return ElementType.TEXT
        except ValueError:
            # shape.placeholder_format raises ValueError for non-placeholder shapes
            pass

        # Check shape name for hints
        if hasattr(shape, "name"):
            name_lower = shape.name.lower()
            if "title" in name_lower:
                return ElementType.TITLE
            if "footer" in name_lower:
                return ElementType.FOOTER
            if "header" in name_lower:
                return ElementType.HEADER

        return ElementType.TEXT
