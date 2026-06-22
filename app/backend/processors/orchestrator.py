"""Translation job orchestrator."""

from __future__ import annotations

import logging
import os
import threading
from pathlib import Path
from typing import Callable, Dict, List, Optional, Tuple

from app.backend.clients.ollama_client import OllamaClient
from app.backend.config import (
    CONTEXT_DETECTION_ENABLED,
    CONTEXT_SAMPLE_CHARS,
    DEFAULT_MAX_BATCH_CHARS,
    DEFAULT_MODEL,
    DYNAMIC_SCENARIO_STRATEGY_ENABLED,
    GENERAL_NUM_CTX,
    LAYOUT_PRESERVATION_MODE,
    OLLAMA_BASE_URL,
    PDF_SKIP_HEADER_FOOTER,
    QWEN_CONTEXT_FLOW_ENABLED,
    SCENARIO_CACHE_VARIANT_ENABLED,
    SUPPORTED_EXTENSIONS,
    TRANSLATION_NUM_CTX,
    TimeoutConfig,
)
from app.backend.processors.com_helpers import is_win32com_available, word_convert
from app.backend.processors.docx_processor import translate_docx
from app.backend.processors.libreoffice_helpers import doc_to_docx, is_libreoffice_available, xls_to_xlsx
from app.backend.processors.pdf_processor import translate_pdf
from app.backend.processors.pptx_processor import translate_pptx
from app.backend.processors.xlsx_processor import translate_xlsx_xls
from app.backend.services.context_prompts import (
    _CONTEXT_DETECTION_PROMPTS,
    _get_context_detection_prompt,
)
from app.backend.services.translation_strategy import (
    build_strategy,
    build_terminology_block,
    detect_translation_scenario,
    scenario_from_profile,
)

logger = logging.getLogger(__name__)


def _cap_terms_by_budget(terms, num_ctx, existing_prompt):
    """Limit terms to fit within ≤25% of num_ctx (conservative: 2.5 chars/token)."""
    max_term_tokens = int(num_ctx * 0.25)
    existing_tokens = int(len(existing_prompt) / 2.5)
    available_tokens = max(0, max_term_tokens - existing_tokens)
    header_chars = len("Terminology constraints:\n")
    used_chars = header_chars
    capped = []
    for t in terms:
        line_chars = len(f"- {t.source_text} => {t.target_text}\n")
        if int((used_chars + line_chars) / 2.5) > available_tokens:
            break
        used_chars += line_chars
        capped.append(t)
    return capped


def _sample_file_text(file_path: Path, max_chars: int = CONTEXT_SAMPLE_CHARS) -> str:
    """Extract the first ~max_chars of text from a file for context detection."""
    ext = file_path.suffix.lower()
    try:
        if ext == ".docx":
            from docx import Document
            doc = Document(str(file_path))
            parts: List[str] = []
            total = 0
            for para in doc.paragraphs:
                t = para.text.strip()
                if not t:
                    continue
                parts.append(t)
                total += len(t)
                if total >= max_chars:
                    break
            return "\n".join(parts)[:max_chars]
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(file_path))
            parts = []
            total = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if not t:
                                continue
                            parts.append(t)
                            total += len(t)
                            if total >= max_chars:
                                break
                    if total >= max_chars:
                        break
                if total >= max_chars:
                    break
            return "\n".join(parts)[:max_chars]
        elif ext in (".xlsx", ".xls"):
            from openpyxl import load_workbook
            wb = load_workbook(str(file_path), read_only=True, data_only=True)
            parts = []
            total = 0
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    for cell in row:
                        if cell is not None:
                            t = str(cell).strip()
                            if t:
                                parts.append(t)
                                total += len(t)
                                if total >= max_chars:
                                    break
                    if total >= max_chars:
                        break
                if total >= max_chars:
                    break
            wb.close()
            return "\n".join(parts)[:max_chars]
        elif ext == ".pdf":
            import fitz
            doc = fitz.open(str(file_path))
            text = ""
            for page in doc:
                text += page.get_text()
                if len(text) >= max_chars:
                    break
            doc.close()
            return text[:max_chars] if text.strip() else file_path.stem.replace("_", " ").replace("-", " ")
        elif ext == ".doc":
            # .doc gets converted to .docx later; use filename
            return file_path.stem.replace("_", " ").replace("-", " ")
    except Exception as exc:
        logger.debug(f"Context sampling failed for {file_path.name}: {exc}")
    return ""


_PHASE0_CHUNK_SIZE = 2000  # Max chars per segment sent to Phase 0 extraction


def _group_into_chunks(units: List[str], max_size: int) -> List[str]:
    """Group semantic text units into chunks not exceeding *max_size* characters."""
    chunks: List[str] = []
    current: List[str] = []
    current_len = 0
    for unit in units:
        unit = unit.strip()
        if not unit:
            continue
        unit_len = len(unit)
        if current_len + unit_len + 1 > max_size and current:
            chunks.append("\n".join(current))
            current = [unit]
            current_len = unit_len
        else:
            current.append(unit)
            current_len += unit_len + 1
    if current:
        chunks.append("\n".join(current))
    return chunks


def _split_paragraphs_to_sentences(paragraphs: List[str]) -> List[str]:
    """Split paragraph-level text into sentences using blingfire."""
    try:
        from blingfire import text_to_sentences
    except ImportError:
        return paragraphs  # fallback: use paragraphs as-is
    sentences: List[str] = []
    for para in paragraphs:
        sents = text_to_sentences(para).split("\n")
        sentences.extend(s.strip() for s in sents if s.strip())
    return sentences


def _extract_spreadsheet_rows(wb) -> List[str]:
    """Extract rows from an openpyxl workbook, preserving row structure with header labels."""
    parts: List[str] = []
    for ws in wb.worksheets:
        header_row: Optional[List[str]] = None
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() if c is not None else "" for c in row]
            non_empty = [c for c in cells if c]
            if not non_empty:
                continue
            if header_row is None:
                header_row = cells
                parts.append(" | ".join(non_empty))
                continue
            # Build row with header context
            labeled: List[str] = []
            for i, val in enumerate(cells):
                if not val:
                    continue
                hdr = header_row[i] if i < len(header_row) and header_row[i] else ""
                if hdr and hdr != val:
                    labeled.append(f"{hdr}: {val}")
                else:
                    labeled.append(val)
            if labeled:
                parts.append(" | ".join(labeled))
    return parts


def _extract_all_segments(file_path: Path, chunk_size: int = _PHASE0_CHUNK_SIZE) -> List[str]:
    """Extract text from a document as semantic units, then group into chunks for Phase 0.

    - xlsx/xls: row-level units with column header prefixes
    - docx/doc: sentence-level units via blingfire
    - pptx: sentence-level units from slide text
    - pdf: filename stem only (placeholder)
    """
    ext = file_path.suffix.lower()
    units: List[str] = []  # semantic units (sentences or rows)
    try:
        if ext == ".docx":
            from docx import Document
            doc = Document(str(file_path))
            paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            units = _split_paragraphs_to_sentences(paragraphs)
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(file_path))
            paragraphs: List[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                paragraphs.append(t)
            units = _split_paragraphs_to_sentences(paragraphs)
        elif ext == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(str(file_path), read_only=True, data_only=True)
            units = _extract_spreadsheet_rows(wb)
            wb.close()
        elif ext == ".xls":
            if is_libreoffice_available():
                import tempfile
                tmp_xlsx = tempfile.mktemp(suffix=".xlsx")
                try:
                    xls_to_xlsx(str(file_path), tmp_xlsx)
                    from openpyxl import load_workbook
                    wb = load_workbook(tmp_xlsx, read_only=True, data_only=True)
                    units = _extract_spreadsheet_rows(wb)
                    wb.close()
                finally:
                    import os as _os
                    _os.unlink(tmp_xlsx) if _os.path.exists(tmp_xlsx) else None
            else:
                units.append(file_path.stem.replace("_", " ").replace("-", " "))
        elif ext == ".pdf":
            import fitz
            doc = fitz.open(str(file_path))
            for page in doc:
                page_text = page.get_text().strip()
                if page_text:
                    units.append(page_text)
            doc.close()
            if units:
                sentences = _split_paragraphs_to_sentences(units)
                units = sentences
            else:
                units.append(file_path.stem.replace("_", " ").replace("-", " "))
        elif ext == ".doc":
            if is_libreoffice_available():
                import tempfile
                tmp_docx = tempfile.mktemp(suffix=".docx")
                try:
                    doc_to_docx(str(file_path), tmp_docx)
                    from docx import Document
                    doc = Document(tmp_docx)
                    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
                    units = _split_paragraphs_to_sentences(paragraphs)
                finally:
                    import os as _os
                    _os.unlink(tmp_docx) if _os.path.exists(tmp_docx) else None
            else:
                units.append(file_path.stem.replace("_", " ").replace("-", " "))
    except Exception as exc:
        logger.debug("Phase 0 text extraction failed for %s: %s", file_path.name, exc)

    if not units:
        return []
    return _group_into_chunks(units, chunk_size)


def _detect_document_context(
    client: OllamaClient,
    sample: str,
    log: Callable[[str], None],
    target_lang: str = "",
) -> str:
    """Ask LLM to describe the document in one sentence (localized prompt)."""
    prompt = _get_context_detection_prompt(target_lang).format(sample=sample)
    payload = client._build_no_system_payload(prompt)
    try:
        ok, result = client._call_ollama(payload)
        if ok and result.strip():
            context = result.strip()[:200]
            log(f"[CONTEXT] Detected: {context}")
            return context
    except Exception as exc:
        logger.debug(f"Context detection failed: {exc}")
    return ""


def _output_name(src: Path, output_format: Optional[str] = None, output_suffix: str = "") -> str:
    """Generate output filename for translated file.

    Args:
        src: Source file path.
        output_format: Optional output format override (e.g., 'pdf' for PDF output).
        output_suffix: Optional suffix inserted before the extension (e.g., '_en', '_vi').

    Returns:
        Output filename with _translated suffix.
    """
    ext = src.suffix.lower()
    stem = src.stem
    tag = f"_{output_suffix}" if output_suffix else ""
    if ext in (".docx", ".pptx", ".xlsx"):
        return f"{stem}_translated{tag}{ext}"
    if ext == ".pdf":
        if output_format == "pdf":
            return f"{stem}_translated{tag}.pdf"
        return f"{stem}_translated{tag}.docx"
    if ext in (".doc", ".xls"):
        return f"{stem}_translated{tag}.docx" if ext == ".doc" else f"{stem}_translated{tag}.xlsx"
    return f"{stem}_translated{tag}{ext}"


def process_files(
    files: List[Path],
    output_dir: Path,
    targets: List[str],
    src_lang: Optional[str],
    include_headers_shapes_via_com: bool,
    ollama_model: str,
    model_type: str = "general",
    system_prompt: str = "",
    profile_id: str = "general",
    num_ctx_override: Optional[int] = None,
    timeout_config: Optional[TimeoutConfig] = None,
    stop_flag: Optional[threading.Event] = None,
    log: Callable[[str], None] = lambda s: None,
    max_batch_chars: int = DEFAULT_MAX_BATCH_CHARS,
    layout_mode: Optional[str] = None,
    output_format: Optional[str] = None,
    output_suffix: str = "",
    mode: str = "translation",
    term_db=None,
    provider_id: Optional[str] = None,
    post_translate_hook: Optional[Callable[[List[Tuple[str, str, str]]], None]] = None,
    output_mode: str = "append",
    block_overrides: Optional[Dict[str, str]] = None,
) -> Tuple[int, int, bool, Optional[OllamaClient], Dict, Optional[str]]:
    """Process files for translation.

    Args:
        files: List of files to process.
        output_dir: Output directory for translated files.
        targets: Target languages.
        src_lang: Source language (or None for auto-detect).
        include_headers_shapes_via_com: Use COM for headers/shapes (Windows).
        ollama_model: Ollama model name.
        model_type: Profile-resolved model type.
        system_prompt: Domain-specific system prompt.
        profile_id: Resolved profile id.
        num_ctx_override: Optional per-job num_ctx override.
        timeout_config: Optional timeout configuration.
        stop_flag: Optional stop flag for cancellation.
        log: Logging callback.
        max_batch_chars: Maximum characters per batch.
        layout_mode: Layout preservation mode (inline|overlay|side_by_side).
        output_format: Output format for PDF (docx|pdf).
        mode: Job mode: 'translation' (default) or 'extraction_only'.
        term_db: Optional TermDB instance for terminology extraction and injection.
        provider_id: Optional provider ID from routing config (p1-cloud-providers).

    Returns:
        Tuple of (processed_count, total_count, stopped, client, term_summary,
        winning_provider_id).
    """
    # Use defaults from config if not specified
    if layout_mode is None:
        layout_mode = LAYOUT_PRESERVATION_MODE
    output_dir.mkdir(parents=True, exist_ok=True)

    # BR-67: multi-target jobs cannot replace in-place (ambiguous which translation owns
    # the paragraph); clamp silently to append so existing callers never break.
    effective_output_mode = "append" if len(targets) > 1 else output_mode

    extraction_only = mode == "extraction_only"
    aggregate_term_summary: Dict = {"extracted": 0, "skipped": 0, "added": 0}

    # ── p1-cloud-providers: build primary client from config or fall back to Ollama ──
    winning_provider: Optional[str] = None
    _provider_id = provider_id  # may be None → Ollama path

    # Try to build an OpenAICompatibleClient if the provider_id is not "ollama-local"
    # and providers.yml has config for it.
    _cloud_client = None
    if _provider_id and _provider_id != "ollama-local":
        try:
            from app.backend.config import load_providers_config
            from app.backend.clients.openai_compatible_client import OpenAICompatibleClient

            _cfg = load_providers_config()
            if _cfg:
                _providers = {p["id"]: p for p in _cfg.get("providers", [])}
                _prov = _providers.get(_provider_id)
                if _prov and _prov.get("enabled") is True:
                    _models = _prov.get("models", {})
                    _model_name = _models.get("translate") or ollama_model
                    _cloud_client = OpenAICompatibleClient(
                        base_url=_prov["base_url"],
                        api_key=_prov["api_key"],
                        model=_model_name,
                        provider_id=_provider_id,
                        verify_ssl=_prov.get("tls_verify", True),
                    )
                    log(f"[PROVIDER] Using cloud provider: {_provider_id} model={_model_name}")
        except Exception as _exc:
            log(f"[PROVIDER] Failed to build cloud client for {_provider_id}: {_exc}; falling back to Ollama")
            _cloud_client = None

    # If cloud client could not be built, walk fallback_chain from config
    if _cloud_client is None and _provider_id and _provider_id != "ollama-local":
        try:
            from app.backend.config import load_providers_config
            from app.backend.clients.openai_compatible_client import OpenAICompatibleClient

            _cfg = load_providers_config()
            if _cfg:
                _chain = _cfg.get("fallback_chain", [])
                _providers = {p["id"]: p for p in _cfg.get("providers", [])}
                for _fb_id in _chain:
                    _prov = _providers.get(_fb_id)
                    if _prov and _prov.get("enabled") is True:
                        _models = _prov.get("models", {})
                        _model_name = _models.get("translate") or ollama_model
                        try:
                            _fb_client = OpenAICompatibleClient(
                                base_url=_prov["base_url"],
                                api_key=_prov["api_key"],
                                model=_model_name,
                                provider_id=_fb_id,
                                verify_ssl=_prov.get("tls_verify", True),
                            )
                            # Quick health probe to verify reachability
                            _fb_ok, _ = _fb_client.health()
                            if _fb_ok:
                                _cloud_client = _fb_client
                                _provider_id = _fb_id
                                log(f"[PROVIDER] Fallback to: {_fb_id}")
                                break
                        except Exception:
                            continue
        except Exception:
            pass

    # Build the primary translation client (OllamaClient or cloud passthrough).
    # Always build OllamaClient as ollama_client so the orchestrator scenario
    # machinery (system_prompt, model_type, _is_translation_dedicated, etc.) has
    # a valid Ollama handle.  When a cloud client was resolved, use it as the
    # primary ``client`` passed to translate_docx/pptx/etc. so translation
    # requests actually reach the cloud endpoint (AC-5, AC-6 — BR-16).
    ollama_client = OllamaClient(
        model=ollama_model,
        model_type=model_type,
        system_prompt=system_prompt,
        profile_id=profile_id,
        num_ctx_override=num_ctx_override,
        timeout=timeout_config,
        log=log,
    )
    if _cloud_client is not None:
        # Cloud provider selected: use it for translation dispatch.
        # OpenAICompatibleClient carries compatibility stubs for the
        # orchestrator attributes (system_prompt, model_type, etc.).
        client = _cloud_client
        log(f"[PROVIDER] Primary translation client: {_provider_id} (cloud)")
    else:
        client = ollama_client
        _provider_id = "ollama-local"  # BR-16: record the provider that will actually process the job
        log("[PROVIDER] Primary translation client: ollama-local")

    processed_count = 0
    total_count = len(files)
    stopped = False
    base_system_prompt = client.system_prompt

    forced_scenario = scenario_from_profile(profile_id)
    if DYNAMIC_SCENARIO_STRATEGY_ENABLED and forced_scenario:
        log(f"[STRATEGY] fixed scenario from profile={profile_id}: {forced_scenario.value}")

    for src in files:
        if stop_flag and stop_flag.is_set():
            log(f"[STOP] stopped at {processed_count}/{total_count} files")
            stopped = True
            break
        ext = src.suffix.lower()
        if ext not in SUPPORTED_EXTENSIONS:
            log(f"[SKIP] Unsupported file: {src.name}")
            continue
        # Determine output format for PDF files
        pdf_output_fmt = output_format if ext == ".pdf" else None
        out_path = output_dir / _output_name(src, output_format=pdf_output_fmt, output_suffix=output_suffix)
        log("=" * 24)
        log(f"Processing: {src.name} ({processed_count + 1}/{total_count})")

        sample = _sample_file_text(src)
        doc_context = ""
        if (
            CONTEXT_DETECTION_ENABLED
            and QWEN_CONTEXT_FLOW_ENABLED
            and not client._is_translation_dedicated()
            and sample
        ):
            doc_context = _detect_document_context(ollama_client, sample, log, target_lang=targets[0] if targets else "")

        if DYNAMIC_SCENARIO_STRATEGY_ENABLED:
            scenario = forced_scenario or detect_translation_scenario(src.name, sample_text=sample, detected_context=doc_context)
            decision = build_strategy(
                base_system_prompt=base_system_prompt,
                model_type=client.model_type,
                scenario=scenario,
                detected_context=doc_context,
                enable_context_flow=QWEN_CONTEXT_FLOW_ENABLED,
            )
            client.system_prompt = decision.system_prompt
            client.set_runtime_options_override(decision.options_override)
            if SCENARIO_CACHE_VARIANT_ENABLED:
                client.set_cache_variant(decision.cache_variant)
            log(
                f"[STRATEGY] mode={'forced' if forced_scenario else 'auto'}, "
                f"scenario={decision.scenario.value}, "
                f"cache_variant={decision.cache_variant}, "
                f"options_override={decision.options_override}"
            )
        else:
            if doc_context:
                client.system_prompt = f"{base_system_prompt}\n\nDocument context: {doc_context}"
            else:
                client.system_prompt = base_system_prompt

        # ------------------------------------------------------------------
        # Phase 0: Term Extraction
        # ------------------------------------------------------------------
        _phase0_hook = None
        _glossary_terms_holder: list = []  # shared across per-file scope; populated by _phase0_hook
        if term_db is not None:
            from app.backend.services.term_extractor import run_phase0_multi, SCENARIO_TO_DOMAIN
            from app.backend.config import (
                TERM_EMBEDDING_MODEL,
                TERM_EMBEDDING_THRESHOLD,
                TERM_EXTRACTION_MODEL,
            )
            _scenario_name = (
                (scenario.value if hasattr(scenario, "value") else str(scenario))
                if DYNAMIC_SCENARIO_STRATEGY_ENABLED else "general"
            )
            _domain = SCENARIO_TO_DOMAIN.get(_scenario_name, "general")
            _source_lang = src_lang or "Chinese"
            _target_langs = targets if targets else ["English"]

            # Resolve PANJIT config for the DB-first embedding path.
            # Reuse the already-loaded provider config (panjit provider).
            _panjit_base_url: Optional[str] = None
            _panjit_api_key: str = ""
            _panjit_tls_verify: bool = False
            try:
                from app.backend.config import load_providers_config
                _p_cfg = load_providers_config()
                if _p_cfg:
                    _p_providers = {p["id"]: p for p in _p_cfg.get("providers", [])}
                    _panjit_prov = _p_providers.get("panjit")
                    if _panjit_prov and _panjit_prov.get("enabled") is True:
                        _panjit_base_url = _panjit_prov.get("base_url") or None
                        _panjit_api_key = _panjit_prov.get("api_key", "")
                        _panjit_tls_verify = bool(_panjit_prov.get("tls_verify", False))
            except Exception as _p_exc:
                logger.warning("[PHASE0] Could not resolve PANJIT config for embedding: %s", _p_exc)

            if extraction_only:
                # Standalone extraction: use 2K-char chunks from _extract_all_segments
                # extraction_only keeps the legacy Ollama path (AC-7, out of scope).
                phase0_segments = _extract_all_segments(src)
                try:
                    term_summary = run_phase0_multi(
                        segments=phase0_segments,
                        source_lang=_source_lang,
                        target_langs=_target_langs,
                        scenario=_scenario_name,
                        document_context=doc_context,
                        term_db=term_db,
                        model=DEFAULT_MODEL,
                        base_url=OLLAMA_BASE_URL,
                        timeout=timeout_config,
                        log=log,
                    )
                    for k in aggregate_term_summary:
                        aggregate_term_summary[k] += term_summary.get(k, 0)
                except Exception as _ph0_exc:
                    log(f"[PHASE0] Unexpected error (non-fatal): {_ph0_exc}")
                    logger.warning("[PHASE0] Unexpected error: %s", _ph0_exc)
            else:
                # Translation mode: build hook for processors to call with their actual segments.
                # Uses the DB-first PANJIT embedding-gated flow when panjit config is available.
                def _phase0_hook(uniq_texts: List[str]) -> None:
                    """Phase 0 hook: extract terms from translation segments, inject into prompts."""
                    try:
                        _ts = run_phase0_multi(
                            segments=uniq_texts,
                            source_lang=_source_lang,
                            target_langs=_target_langs,
                            scenario=_scenario_name,
                            document_context=doc_context,
                            term_db=term_db,
                            model=DEFAULT_MODEL,
                            base_url=OLLAMA_BASE_URL,
                            timeout=timeout_config,
                            log=log,
                            panjit_base_url=_panjit_base_url,
                            panjit_api_key=_panjit_api_key,
                            panjit_tls_verify=_panjit_tls_verify,
                            embedding_model=TERM_EMBEDDING_MODEL,
                            extraction_model=TERM_EXTRACTION_MODEL,
                            embedding_threshold=TERM_EMBEDDING_THRESHOLD,
                        )
                        for k in aggregate_term_summary:
                            aggregate_term_summary[k] += _ts.get(k, 0)
                    except Exception as _exc:
                        log(f"[PHASE0] Unexpected error (non-fatal): {_exc}")
                        logger.warning("[PHASE0] Unexpected error: %s", _exc)
                        return

                    # Retrieve approved/conf=1.0 terms for this document
                    _doc_source_texts = _ts.get("extracted_source_texts", [])
                    _seen_keys: set = set()
                    _top_terms: list = []
                    for _lang in _target_langs:
                        for _t in term_db.get_document_terms(_lang, _domain, _doc_source_texts):
                            _k = (_t.source_text, _t.target_text)
                            if _k not in _seen_keys:
                                _seen_keys.add(_k)
                                _top_terms.append(_t)
                    # Share approved terms with translate_texts for glossary substitution (P2-5)
                    _glossary_terms_holder[:] = _top_terms
                    if not _top_terms:
                        return

                    # Phase 1: inject unless a dedicated translation variant (backward-compat guard)
                    if not client._is_translategemma_model():
                        _p1_ctx = TRANSLATION_NUM_CTX if client._is_translation_dedicated() else GENERAL_NUM_CTX
                        _p1_terms = _cap_terms_by_budget(_top_terms, _p1_ctx, client.system_prompt or "")
                        if _p1_terms:
                            _block = build_terminology_block(_p1_terms)
                            client.system_prompt = (
                                client.system_prompt.rstrip() + "\n\n" + _block
                                if client.system_prompt.strip()
                                else _block
                            )
                            log(f"[PHASE0] Injected {len(_p1_terms)}/{len(_top_terms)} terms into Phase 1 (budget: {_p1_ctx} ctx)")
                        elif _top_terms:
                            log(f"[PHASE0] WARNING: No room for terms in Phase 1 (ctx={_p1_ctx})")

        # extraction_only mode: skip translation
        if extraction_only:
            processed_count += 1
            log(f"[EXTRACTION] Done: {src.name} (extraction only)")
            continue

        try:
            if ext == ".docx":
                stopped = translate_docx(
                    str(src),
                    str(out_path),
                    targets,
                    src_lang,
                    client,
                    include_headers_shapes_via_com=include_headers_shapes_via_com,
                    stop_flag=stop_flag,
                    log=log,
                    max_batch_chars=max_batch_chars,
                    pre_translate_hook=_phase0_hook,
                    post_translate_hook=post_translate_hook,
                    terms_getter=lambda: list(_glossary_terms_holder),
                    output_mode=effective_output_mode,
                    block_overrides=block_overrides,
                )
            elif ext == ".doc":
                tmp_docx = str(output_dir / f"{src.stem}__tmp.docx")
                if is_libreoffice_available():
                    log("[DOC] Converting to .docx via LibreOffice")
                    doc_to_docx(str(src), tmp_docx)
                elif is_win32com_available():
                    log("[DOC] Converting to .docx via COM")
                    word_convert(str(src), tmp_docx, 16)
                else:
                    log(
                        "[DOC] Cannot convert .doc: neither LibreOffice nor "
                        "Word COM is available. Install LibreOffice: "
                        "sudo apt install libreoffice-core (Linux) / "
                        "brew install --cask libreoffice (macOS)"
                    )
                    continue
                try:
                    stopped = translate_docx(
                        tmp_docx,
                        str(out_path),
                        targets,
                        src_lang,
                        client,
                        include_headers_shapes_via_com=include_headers_shapes_via_com,
                        stop_flag=stop_flag,
                        log=log,
                        max_batch_chars=max_batch_chars,
                        pre_translate_hook=_phase0_hook,
                        post_translate_hook=post_translate_hook,
                        terms_getter=lambda: list(_glossary_terms_holder),
                        output_mode=effective_output_mode,
                        block_overrides=block_overrides,
                    )
                finally:
                    try:
                        os.remove(tmp_docx)
                    except OSError:
                        pass
            elif ext == ".pptx":
                stopped = translate_pptx(
                    str(src),
                    str(out_path),
                    targets,
                    src_lang,
                    client,
                    stop_flag=stop_flag,
                    log=log,
                    max_batch_chars=max_batch_chars,
                    pre_translate_hook=_phase0_hook,
                    post_translate_hook=post_translate_hook,
                    terms_getter=lambda: list(_glossary_terms_holder),
                    output_mode=effective_output_mode,
                    block_overrides=block_overrides,
                )
            elif ext in (".xlsx", ".xls"):
                stopped = translate_xlsx_xls(
                    str(src),
                    str(out_path),
                    targets,
                    src_lang,
                    client,
                    stop_flag=stop_flag,
                    log=log,
                    max_batch_chars=max_batch_chars,
                    pre_translate_hook=_phase0_hook,
                    post_translate_hook=post_translate_hook,
                    terms_getter=lambda: list(_glossary_terms_holder),
                    block_overrides=block_overrides,
                )
            elif ext == ".pdf":
                log(f"[PDF] Using output_format={output_format}, layout_mode={layout_mode}")
                stopped = translate_pdf(
                    str(src),
                    str(out_path),
                    targets,
                    src_lang,
                    client,
                    stop_flag=stop_flag,
                    log=log,
                    skip_header_footer=PDF_SKIP_HEADER_FOOTER,
                    output_format=output_format or "docx",
                    layout_mode=layout_mode,
                    pre_translate_hook=_phase0_hook,
                    post_translate_hook=post_translate_hook,
                    block_overrides=block_overrides,
                )
            else:
                log(f"[SKIP] Unsupported file: {src.name}")
                continue
            processed_count += 1
            if stopped:
                log(f"[STOP] file interrupted: {src.name}")
                break
            log(f"Done: {src.name} -> {out_path.name}")
        except Exception as exc:
            log(f"[ERROR] {src.name}: {exc}")
        finally:
            client.system_prompt = base_system_prompt
            client.set_runtime_options_override(None)
            client.set_cache_variant(None)
    if stopped:
        log(f"[STOP] job stopped after {processed_count}/{total_count} files")
    else:
        log(f"[DONE] job complete: {processed_count}/{total_count} files")
    # p1-cloud-providers: record winning provider (ollama-local or cloud ID)
    winning_provider = _provider_id or "ollama-local"
    return processed_count, total_count, stopped, client, aggregate_term_summary, winning_provider
