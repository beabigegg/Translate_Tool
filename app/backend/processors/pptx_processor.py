"""PPTX translation processor."""

from __future__ import annotations

import os
import re
import threading
import zipfile
from typing import Any, Callable, Dict, List, Optional, Tuple, Union
from xml.etree import ElementTree as ET

import pptx
from pptx.table import _Cell
from pptx.util import Pt as PPTPt

from app.backend.clients.ollama_client import OllamaClient
from app.backend.config import DEFAULT_MAX_BATCH_CHARS, MAX_SEGMENTS, MAX_TEXT_LENGTH
from app.backend.services.translation_service import translate_texts
from app.backend.utils.exceptions import check_document_size_limits
from app.backend.utils.text_utils import normalize_text, should_translate


# Segment types for tracking
SEGMENT_TEXT_FRAME = "text_frame"
SEGMENT_TABLE_CELL = "table_cell"
SEGMENT_SMARTART = "smartart"

# XML namespaces for SmartArt
SMARTART_NS = {
    "dgm": "http://schemas.openxmlformats.org/drawingml/2006/diagram",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
}


def _ppt_text_of_tf(tf: Any) -> str:
    return "\n".join([p.text for p in tf.paragraphs])


def _ppt_tail_equals(tf: Any, translations: List[str]) -> bool:
    if len(tf.paragraphs) < len(translations):
        return False
    tail = tf.paragraphs[-len(translations):]
    for para, expect in zip(tail, translations):
        if normalize_text(para.text) != normalize_text(expect):
            return False
        if any((r.font.italic is not True) and (r.text or "").strip() for r in para.runs):
            return False
    return True


def _ppt_append(tf: Any, text_block: str) -> None:
    p = tf.add_paragraph()
    p.text = text_block
    for r in p.runs:
        r.font.italic = True
        r.font.size = PPTPt(12)


def _cell_text(cell: _Cell) -> str:
    """Get text from a table cell."""
    return cell.text.strip()


def _cell_tail_equals(cell: _Cell, translations: List[str]) -> bool:
    """Check if translation already appended to cell."""
    tf = cell.text_frame
    if len(tf.paragraphs) < len(translations):
        return False
    tail = tf.paragraphs[-len(translations):]
    for para, expect in zip(tail, translations):
        if normalize_text(para.text) != normalize_text(expect):
            return False
        if any((r.font.italic is not True) and (r.text or "").strip() for r in para.runs):
            return False
    return True


def _cell_append(cell: _Cell, text_block: str) -> None:
    """Append translated text to a table cell."""
    tf = cell.text_frame
    p = tf.add_paragraph()
    p.text = text_block
    for r in p.runs:
        r.font.italic = True
        r.font.size = PPTPt(10)  # Slightly smaller for table cells


def _extract_smartart_texts(pptx_path: str) -> List[Tuple[str, str, str]]:
    """Extract text from SmartArt diagrams in a PPTX file.

    Args:
        pptx_path: Path to the PPTX file.

    Returns:
        List of tuples: (diagram_file, xpath, text)
    """
    smartart_texts: List[Tuple[str, str, str]] = []

    with zipfile.ZipFile(pptx_path, "r") as zf:
        # Find all diagram data files
        diagram_files = [f for f in zf.namelist() if f.startswith("ppt/diagrams/data") and f.endswith(".xml")]

        for diagram_file in diagram_files:
            try:
                xml_content = zf.read(diagram_file).decode("utf-8")
                root = ET.fromstring(xml_content)

                # Find all text elements in the diagram
                # SmartArt text is typically in <a:t> tags within <dgm:t> or directly
                for idx, t_elem in enumerate(root.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}t")):
                    text = t_elem.text
                    if text and text.strip():
                        # Store file, index, and text
                        smartart_texts.append((diagram_file, str(idx), text.strip()))
            except (ET.ParseError, KeyError, UnicodeDecodeError):
                continue

    return smartart_texts


def _update_smartart_texts(
    pptx_path: str,
    out_path: str,
    translations: Dict[str, str],
    output_mode: str = "append",
) -> None:
    """Update SmartArt texts with translations.

    Args:
        pptx_path: Path to the original PPTX file.
        out_path: Path to save the updated PPTX file.
        translations: Dict mapping original text to translated text (appended).
        output_mode: ``"replace"`` to overwrite source text; ``"append"`` to append in
            parentheses (default).
    """
    import shutil
    import tempfile

    # Copy to temp file first
    temp_dir = tempfile.mkdtemp()
    temp_pptx = os.path.join(temp_dir, "temp.pptx")
    shutil.copy2(pptx_path, temp_pptx)

    # Read and modify the PPTX
    with zipfile.ZipFile(temp_pptx, "r") as zf_in:
        # Get all file contents
        file_contents = {}
        for name in zf_in.namelist():
            file_contents[name] = zf_in.read(name)

    # Modify SmartArt diagram files
    for filename, content in list(file_contents.items()):
        if filename.startswith("ppt/diagrams/data") and filename.endswith(".xml"):
            try:
                xml_content = content.decode("utf-8")
                root = ET.fromstring(xml_content)
                modified = False

                for t_elem in root.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}t"):
                    text = t_elem.text
                    if text and text.strip() in translations:
                        original = text.strip()
                        translated = translations[original]
                        if output_mode == "replace":
                            t_elem.text = translated
                        else:
                            # Append translation in parentheses (default)
                            t_elem.text = f"{text}\n({translated})"
                        modified = True

                if modified:
                    # Re-serialize with proper XML declaration
                    ET.register_namespace("", "http://schemas.openxmlformats.org/drawingml/2006/diagram")
                    ET.register_namespace("a", "http://schemas.openxmlformats.org/drawingml/2006/main")
                    file_contents[filename] = ET.tostring(root, encoding="unicode").encode("utf-8")
            except (ET.ParseError, UnicodeDecodeError):
                continue

    # Write the modified PPTX
    with zipfile.ZipFile(out_path, "w", zipfile.ZIP_DEFLATED) as zf_out:
        for name, content in file_contents.items():
            zf_out.writestr(name, content)

    # Cleanup
    shutil.rmtree(temp_dir, ignore_errors=True)


def translate_pptx(
    in_path: str,
    out_path: str,
    targets: List[str],
    src_lang: Optional[str],
    client: OllamaClient,
    stop_flag: Optional[threading.Event] = None,
    log: Callable[[str], None] = lambda s: None,
    max_segments: int = MAX_SEGMENTS,
    max_text_length: int = MAX_TEXT_LENGTH,
    max_batch_chars: int = DEFAULT_MAX_BATCH_CHARS,
    pre_translate_hook: Optional[Callable[[List[str]], None]] = None,
    post_translate_hook: Optional[Callable[[List[Tuple[str, str, str]]], None]] = None,
    terms_getter: Optional[Callable[[], list]] = None,
    output_mode: str = "append",
    block_overrides: Optional[Dict[str, str]] = None,
    status_callback: Optional[Callable[[Optional[str]], None]] = None,
) -> bool:
    prs = pptx.Presentation(in_path)
    # segs: List of (segment_type, object_ref, text, row, col, table_id)
    # row/col/table_id are None for non-table segments (IP-4 / BR-81).
    segs: List[Tuple[str, Any, str, Optional[int], Optional[int], Optional[int]]] = []
    total_text_length = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            # Handle tables - use has_table property instead of hasattr
            if getattr(shape, "has_table", False):
                try:
                    table = shape.table
                    shape_id = id(shape)  # unique per-table identifier (IP-4)
                    for r_idx, row in enumerate(table.rows):
                        for c_idx, cell in enumerate(row.cells):
                            txt = _cell_text(cell)
                            if txt:
                                segs.append((SEGMENT_TABLE_CELL, cell, txt, r_idx, c_idx, shape_id))
                                total_text_length += len(txt)
                    continue
                except Exception:
                    # Shape reports has_table but doesn't contain one
                    pass

            # Handle text frames
            if not getattr(shape, "has_text_frame", False):
                continue
            tf = shape.text_frame
            txt = _ppt_text_of_tf(tf)
            if txt.strip():
                segs.append((SEGMENT_TEXT_FRAME, tf, txt, None, None, None))
                total_text_length += len(txt)

    # Extract SmartArt texts
    smartart_texts = _extract_smartart_texts(in_path)
    smartart_segs: List[Tuple[str, str, str]] = []  # (file, idx, text)
    for diagram_file, idx, text in smartart_texts:
        smartart_segs.append((diagram_file, idx, text))
        total_text_length += len(text)

    total_segs = len(segs) + len(smartart_segs)

    check_document_size_limits(
        segment_count=total_segs,
        total_text_length=total_text_length,
        max_segments=max_segments,
        max_text_length=max_text_length,
        document_type="PowerPoint document",
    )

    # Count segments by type for logging
    tf_count = sum(1 for seg in segs if seg[0] == SEGMENT_TEXT_FRAME)
    cell_count = sum(1 for seg in segs if seg[0] == SEGMENT_TABLE_CELL)
    smartart_count = len(smartart_segs)
    log(f"[PPTX] segments: {total_segs} (text frames: {tf_count}, table cells: {cell_count}, SmartArt: {smartart_count})")

    # IP-4: separate table cell segments from non-table segments.
    # Text frames and SmartArt use the existing translate_texts path (2-element keys).
    # Table cells use per-table serialization (translate_once; 3-element keys).
    tf_segs = [seg for seg in segs if seg[0] == SEGMENT_TEXT_FRAME]
    cell_segs = [seg for seg in segs if seg[0] == SEGMENT_TABLE_CELL]

    # Collect unique texts from non-table segments + SmartArt for the translate_texts batch.
    all_tf_texts = [seg[2] for seg in tf_segs] + [s for _, _, s in smartart_segs]
    _seen: set[str] = set()
    uniq: list[str] = []
    for t in all_tf_texts:
        if t not in _seen and should_translate(t, (src_lang or "auto")):
            _seen.add(t)
            uniq.append(t)

    if pre_translate_hook:
        pre_translate_hook(uniq)
    _terms = terms_getter() if terms_getter else None

    # p3-llm-judge: block_overrides seam
    _file_stem = os.path.splitext(os.path.basename(in_path))[0]
    stopped = False

    # para_tmap: 2-element key (tgt, text) for text frames + SmartArt.
    # final_tmap: 3-element key (tgt, text, col) for ALL segments:
    #   text frames → col=None; table cells → col=0-based column index.
    para_tmap: Dict = {}

    if block_overrides is not None:
        for idx, src_text in enumerate(uniq):
            block_id = f"pptx:{_file_stem}:{idx}"
            for tgt in targets:
                if block_id in block_overrides:
                    para_tmap[(tgt, src_text)] = block_overrides[block_id]
                else:
                    para_tmap[(tgt, src_text)] = src_text
        log(f"[PPTX] block_overrides applied: {len(block_overrides)} overrides, {len(uniq)} blocks")
    else:
        para_tmap, _, fail_cnt, stopped = translate_texts(
            uniq,
            targets,
            src_lang,
            client,
            max_batch_chars=max_batch_chars,
            stop_flag=stop_flag,
            log=log,
            terms=_terms,
            status_callback=status_callback,
        )

        if fail_cnt and not stopped:
            from app.backend.utils.translation_verification import verify_and_fill_tmap
            verify_and_fill_tmap(para_tmap, client, src_lang, stop_flag=stop_flag, log=log)

    # Build final_tmap: re-key text frames to (tgt, text, None)
    final_tmap: Dict = {}
    for (tgt, text), tr in para_tmap.items():
        final_tmap[(tgt, text, None)] = tr

    # IP-4: translate each table group via serializer → one translate_once call per table.
    if cell_segs and not stopped:
        from collections import defaultdict
        from dataclasses import dataclass as _dc
        from app.backend.utils import table_serializer

        @_dc
        class _PptCellProxy:
            row: int
            col: int
            content: str
            is_numeric: bool = False

        # Group by table_id (shape id)
        table_groups: Dict = defaultdict(list)
        for seg in cell_segs:
            table_groups[seg[5]].append(seg)  # seg[5] = table_id

        for table_id, t_segs in table_groups.items():
            if stop_flag and stop_flag.is_set():
                stopped = True
                break
            num_rows = max(seg[3] for seg in t_segs) + 1  # seg[3] = row
            num_cols = max(seg[4] for seg in t_segs) + 1  # seg[4] = col

            cells_by_pos = {(seg[3], seg[4]): seg[2] for seg in t_segs}
            proxy_cells = [
                _PptCellProxy(row=r, col=c, content=cells_by_pos.get((r, c), ""))
                for r in range(num_rows) for c in range(num_cols)
            ]

            for tgt in targets:
                if stop_flag and stop_flag.is_set():
                    stopped = True
                    break
                src_for_prompt = src_lang or "auto"
                serialized = table_serializer.serialize(proxy_cells)
                prompt = client._build_table_translate_prompt(serialized, src_for_prompt, tgt)
                # grid stays None on any error → fallback always runs (BR-82)
                grid = None
                try:
                    ok, response = client.translate_once(prompt, tgt, src_lang)
                    if ok:
                        grid = table_serializer.parse(response, num_rows, num_cols)
                        if grid is None:
                            from app.backend.utils.logging_utils import logger as _log
                            _log.warning(
                                "[PPTX] Table %s: parse() returned None (expected %d×%d); "
                                "falling back to per-cell batch for target=%s",
                                table_id, num_rows, num_cols, tgt,
                            )
                except Exception as exc:
                    from app.backend.utils.logging_utils import logger as _log
                    _log.warning(
                        "[PPTX] Table %s translate_once failed (target=%s): %s; "
                        "falling back to per-cell batch",
                        table_id, tgt, exc,
                    )
                if grid is not None:
                    for seg in t_segs:
                        r, c, txt = seg[3], seg[4], seg[2]
                        if txt.strip():
                            final_tmap[(tgt, txt, c)] = grid[r][c]
                else:
                    # Fallback: per-cell batch (BR-82)
                    fallback_texts = list(dict.fromkeys(
                        seg[2] for seg in t_segs
                        if seg[2].strip() and should_translate(seg[2], src_for_prompt)
                    ))
                    if fallback_texts:
                        fb_tmap, _, _, _ = translate_texts(
                            fallback_texts, [tgt], src_lang, client,
                            max_batch_chars=max_batch_chars,
                            stop_flag=stop_flag, log=log,
                        )
                        for (fb_tgt, fb_text), fb_tr in fb_tmap.items():
                            for seg in t_segs:
                                if seg[2] == fb_text:
                                    final_tmap[(fb_tgt, fb_text, seg[4])] = fb_tr

    # Apply translations to text frames and table cells
    ok_cnt = skip_cnt = 0
    for seg in segs:
        seg_type, obj_ref, s = seg[0], seg[1], seg[2]
        col = seg[4]  # None for text frames, 0-based int for table cells
        if not all((tgt, s, col) in final_tmap for tgt in targets):
            continue
        trs = [final_tmap[(tgt, s, col)] for tgt in targets]

        if seg_type == SEGMENT_TEXT_FRAME:
            if output_mode == "replace":
                # Overwrite runs in-place with first translation (multi-target clamped to
                # append by orchestrator before reaching here, BR-67).
                replacement = trs[0]
                tf = obj_ref
                # TODO (R2): SmartArt path stays append-only; this covers text frames only.
                all_runs = [r for p in tf.paragraphs for r in p.runs]
                if all_runs:
                    all_runs[0].text = replacement
                    for r in all_runs[1:]:
                        r.text = ""
                elif tf.paragraphs:
                    tf.paragraphs[0].text = replacement
                ok_cnt += 1
            else:
                if _ppt_tail_equals(obj_ref, trs):
                    skip_cnt += 1
                    continue
                for block in trs:
                    _ppt_append(obj_ref, block)
                ok_cnt += 1
        elif seg_type == SEGMENT_TABLE_CELL:
            if output_mode == "replace":
                replacement = trs[0]
                tf = obj_ref.text_frame
                all_runs = [r for p in tf.paragraphs for r in p.runs]
                if all_runs:
                    all_runs[0].text = replacement
                    for r in all_runs[1:]:
                        r.text = ""
                elif tf.paragraphs:
                    tf.paragraphs[0].text = replacement
                ok_cnt += 1
            else:
                if _cell_tail_equals(obj_ref, trs):
                    skip_cnt += 1
                    continue
                for block in trs:
                    _cell_append(obj_ref, block)
                ok_cnt += 1

    prs.save(out_path)

    # Apply SmartArt translations (requires direct XML manipulation)
    if smartart_segs:
        # Build translation map for SmartArt: original -> combined translations
        smartart_tmap: Dict[str, str] = {}
        for _, _, text in smartart_segs:
            if all((tgt, text) in para_tmap for tgt in targets):
                translations = [para_tmap[(tgt, text)] for tgt in targets]
                smartart_tmap[text] = " / ".join(translations)

        if smartart_tmap:
            _update_smartart_texts(out_path, out_path, smartart_tmap, output_mode=output_mode)
            log(f"[PPTX] SmartArt translated: {len(smartart_tmap)} items")

    if post_translate_hook is not None:
        tuples: List[Tuple[str, str, str]] = []
        for idx, src_text in enumerate(uniq):
            for tgt in targets:
                if (tgt, src_text) in para_tmap:
                    tuples.append((f"pptx:{_file_stem}:{idx}", src_text, para_tmap[(tgt, src_text)]))
        if tuples:
            post_translate_hook(tuples)

    if stopped:
        log(f"[PPTX] partial output: {os.path.basename(out_path)}")
    else:
        log(f"[PPTX] output: {os.path.basename(out_path)}")

    return stopped
